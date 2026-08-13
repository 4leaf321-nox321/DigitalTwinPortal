"""
PPT 템플릿 생성 스크립트
플레이스홀더가 포함된 기본 템플릿을 생성합니다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_template():
    """기본 PPT 템플릿 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    primary_color = RGBColor(16, 185, 129)  # #10B981
    white = RGBColor(255, 255, 255)
    text_color = RGBColor(55, 65, 81)  # #374151

    # =====================================
    # 슬라이드 1: 표지
    # =====================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 상단 녹색 배경
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary_color
    shape.line.fill.background()

    # 과제명
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{과제명}}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    # 연도
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{과제년도}}년 과제 보고서"
    p.font.size = Pt(20)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    # 사업부/프로세스
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "사업부: {{사업부}}  |  프로세스: {{프로세스}}"
    p.font.size = Pt(14)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # 과제PL/작성자
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "과제 PL: {{과제PL}}  |  작성자: {{작성자}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # 기간
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "기간: {{시작}}월 ~ {{종료}}월"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # =====================================
    # 슬라이드 2: 과제 개요
    # =====================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "과제 개요"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_color

    # 개요 정보 (테이블 형식 텍스트)
    overview_items = [
        ("과제명", "{{과제명}}"),
        ("사업부", "{{사업부}}"),
        ("프로세스", "{{프로세스}}"),
        ("과제영역", "{{과제영역}}"),
        ("과제구분", "{{과제구분}}"),
        ("기간", "{{시작}}월 ~ {{종료}}월"),
        ("진행상태", "{{진행상태}}"),
        ("진행률", "{{진행률}}%"),
        ("PoC 과제", "{{PoC과제여부}}"),
        ("중점 과제", "{{중점과제여부}}")
    ]

    y_pos = 1.2
    for label, value in overview_items:
        # 라벨
        txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = text_color

        # 값
        txBox = slide2.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(7), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(12)
        p.font.color.rgb = text_color

        y_pos += 0.5

    # =====================================
    # 슬라이드 3: 담당자 정보
    # =====================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "담당자 정보"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_color

    # 담당부서
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "담당부서"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_color

    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{담당부서목록}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color

    # 참여인력
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "과제 참여인력"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_color

    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{과제참여인력목록}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color

    # =====================================
    # 슬라이드 4: 성과 목록
    # =====================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "성과 목록"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_color

    # 성과 목록 플레이스홀더
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{성과목록}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color

    # =====================================
    # 슬라이드 5: 액션아이템 및 비고
    # =====================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "액션아이템 및 비고"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_color

    # 액션아이템
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "액션아이템"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_color

    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{액션아이템목록}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color

    # 상세설명
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "과제 상세설명"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_color

    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{과제상세설명}}"
    p.font.size = Pt(12)
    p.font.color.rgb = text_color

    # =====================================
    # 슬라이드 6: 상세 과제 정보 (스타일 샘플)
    # — 이 슬라이드는 보고서 생성 시 활성화된 섹션 수만큼 복제됩니다.
    # — 각 shape의 name으로 스타일을 추출한 뒤 동적 콘텐츠를 생성합니다.
    # =====================================
    slide_detail = prs.slides.add_slide(prs.slide_layouts[6])

    # 섹션 제목 스타일 샘플
    txBox = slide_detail.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    txBox.name = "DETAIL_TITLE_STYLE"
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "{{섹션제목}}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_color

    # 상위 항목 (- 레벨) 스타일 샘플
    txBox = slide_detail.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.35))
    txBox.name = "DETAIL_PARENT_STYLE"
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "- 샘플 상위 항목"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color

    # 하위 항목 (· 레벨) 스타일 샘플
    txBox = slide_detail.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.6), Inches(0.3))
    txBox.name = "DETAIL_CHILD_STYLE"
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "· 샘플 하위 항목"
    p.font.size = Pt(10)
    p.font.bold = False
    p.font.color.rgb = RGBColor(107, 114, 128)

    # 이미지 캡션 스타일 샘플
    txBox = slide_detail.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4), Inches(0.3))
    txBox.name = "DETAIL_CAPTION_STYLE"
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "샘플 캡션"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(107, 114, 128)
    p.alignment = PP_ALIGN.CENTER

    # =====================================
    # 슬라이드 7: 플레이스홀더 레퍼런스
    # =====================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 — 연한 회색
    bg_shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(249, 250, 251)
    bg_shape.line.fill.background()

    # 제목
    txBox = slide6.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "사용 가능한 플레이스홀더 목록"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color

    subtitle_color = RGBColor(30, 64, 175)   # 파란 소제목
    label_color = RGBColor(75, 85, 99)        # 회색 본문
    small = Pt(7)
    section_font = Pt(9)

    def add_section(slide, x, y, w, title, items):
        """섹션 제목 + 항목 리스트를 하나의 텍스트박스로 추가"""
        # 줄 수 기반 높이 계산
        line_h = 0.17
        h = 0.3 + len(items) * line_h
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True

        # 섹션 제목
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = section_font
        p.font.bold = True
        p.font.color.rgb = subtitle_color
        p.space_after = Pt(2)

        # 항목들
        for placeholder_name, desc in items:
            p = tf.add_paragraph()
            p.text = f"{placeholder_name}  —  {desc}"
            p.font.size = small
            p.font.color.rgb = label_color
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    # ─── 왼쪽 열: 기본 정보 / 참여인력 ───
    col1_x = 0.3
    col_w = 4.5

    basic_items = [
        ("{{과제명}}", "과제명"),
        ("{{과제년도}}", "과제년도"),
        ("{{사업부}}", "사업부"),
        ("{{프로세스}}", "프로세스"),
        ("{{과제PL}}", "과제 PL"),
        ("{{작성자}}", "작성자"),
        ("{{시작}}", "시작월"),
        ("{{종료}}", "종료월"),
        ("{{과제영역}}", "과제영역"),
        ("{{과제구분}}", "과제구분"),
        ("{{진행상태}}", "진행상태"),
        ("{{진행률}}", "진행률"),
        ("{{PoC과제여부}}", "PoC 과제 여부"),
        ("{{중점과제여부}}", "중점 과제 여부"),
        ("{{과제상세설명}}", "과제 상세설명"),
        ("{{담당부서목록}}", "담당부서 (쉼표 구분)"),
    ]
    add_section(slide6, col1_x, 0.75, col_w, "■ 기본 정보", basic_items)

    personnel_items = [
        ("{{과제참여인력목록}}", "참여인력 전체 (줄바꿈)"),
        ("{{과제참여인력_이름목록}}", "참여인력 이름 (쉼표 구분)"),
        ("{{참여인력수}}", "참여인력 총 수"),
        ("{{참여인력N_이름}}", "N번째 참여인력 이름"),
        ("{{참여인력N_부서}}", "N번째 참여인력 부서"),
        ("{{참여인력N_knoxId}}", "N번째 참여인력 knoxId"),
        ("{{참여인력N_선택사업부}}", "N번째 참여인력 선택사업부"),
    ]
    add_section(slide6, col1_x, 4.0, col_w, "■ 참여인력  (N = 1 ~ 20)", personnel_items)

    # ─── 오른쪽 열: 성과 / 액션아이템 ───
    col2_x = 5.2

    perf_items = [
        ("{{성과목록}}", "성과 전체 (줄바꿈)"),
        ("{{성과수}}", "성과 총 수"),
        ("{{성과N_항목}}", "N번째 성과 항목명"),
        ("{{성과N_대분류}}", "N번째 성과 대분류"),
        ("{{성과N_기여도}}", "N번째 성과 기여도(%)"),
        ("{{성과N_현재수준}}", "N번째 성과 현재수준"),
        ("{{성과N_목표수준}}", "N번째 성과 목표수준"),
        ("{{성과N_실적수준}}", "N번째 성과 실적수준"),
        ("{{성과N_단위}}", "N번째 성과 단위"),
    ]
    add_section(slide6, col2_x, 0.75, col_w, "■ 성과  (N = 1 ~ 20)", perf_items)

    action_items = [
        ("{{액션아이템목록}}", "액션아이템 전체 (줄바꿈)"),
        ("{{액션아이템수}}", "액션아이템 총 수"),
        ("{{액션아이템N_이름}}", "N번째 액션아이템 제목"),
        ("{{액션아이템N_목표일}}", "N번째 액션아이템 목표일"),
        ("{{액션아이템N_완료여부}}", "완료 / 진행중"),
        ("{{액션아이템N_완료일}}", "N번째 액션아이템 완료일"),
        ("{{액션아이템N_세부항목}}", "세부항목 (쉼표 구분)"),
        ("{{액션아이템_마일스톤}}", "마일스톤 타임라인 (도형)"),
    ]
    add_section(slide6, col2_x, 3.0, col_w, "■ 액션아이템  (N = 1 ~ 20)", action_items)

    detail_items = [
        ("{{상세정보_과제개요}}", "과제 개요 (줄바꿈 텍스트)"),
        ("{{상세정보_추진배경}}", "추진 배경 (줄바꿈 텍스트)"),
        ("{{상세정보_과제목표}}", "과제 목표 (줄바꿈 텍스트)"),
        ("{{상세정보_상세내용}}", "상세 내용 (줄바꿈 텍스트)"),
        ("{{상세정보_성과}}", "성과 (텍스트 + 테이블)"),
        ("{{상세정보_산출물}}", "산출물 (줄바꿈 텍스트)"),
        ("{{상세정보_향후계획}}", "향후 계획 (줄바꿈 텍스트)"),
        ("{{이미지_좌측}}", "좌측 이미지 영역"),
        ("{{이미지_우측}}", "우측 이미지 영역"),
    ]
    add_section(slide6, col2_x, 5.2, col_w, "■ 상세 과제 정보", detail_items)

    # 하단 안내문
    txBox = slide6.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(9.4), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    notes = [
        "※ N 자리에 1~20 숫자를 넣으면 해당 순번의 값으로 치환됩니다.  예) {{성과1_항목}}, {{액션아이템3_이름}}",
        "※ 데이터가 없는 번호는 빈 문자열로 치환됩니다.",
        "※ 이 페이지는 레퍼런스용입니다. 템플릿 수정 후 이 슬라이드를 삭제해도 무방합니다.",
    ]
    p = tf.paragraphs[0]
    p.text = notes[0]
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(107, 114, 128)
    for note in notes[1:]:
        p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(107, 114, 128)
        p.space_before = Pt(2)

    # 저장
    template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'templates', 'ppt', 'project_report_template.pptx')
    os.makedirs(os.path.dirname(template_path), exist_ok=True)
    prs.save(template_path)
    print(f"템플릿이 생성되었습니다: {template_path}")
    return template_path


if __name__ == '__main__':
    create_template()
