"""
Digital Twin Dashboard Routes
"""
import io
import os
import zipfile
from datetime import datetime, timedelta
from flask import request, send_file
from flask_jwt_extended import jwt_required, get_jwt_identity
from app.modules.digital_twin_dashboard import bp
from app.modules.digital_twin_dashboard.models import (
    ModuleSettings, Division, Department, ProcessCategory,
    ProjectDomain, TaskCategory, TaskStatus,
    PerformanceCategory, PerformanceSubcategory,
    DashboardData, DashboardSnapshot, DashboardActivityLog,
    ProjectAttachment, ReportImage, UPLOAD_FOLDER,
    KPICategory, KPI, KPIDashboardCard
)
from app.modules.auth.models import User
from app.shared.responses import success_response, error_response
from app.extensions import db
from sqlalchemy.orm.attributes import flag_modified
from app.modules.digital_twin_dashboard.v2_sync import request_sync
from app.shared.timeutil import iso_kst, now_utc_iso_z


@bp.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    return success_response({'status': 'healthy', 'module': 'digital-twin-dashboard'})


@bp.route('/projects', methods=['GET'])
def get_projects():
    """Get all projects."""
    # TODO: Implement with database
    return success_response([])


@bp.route('/projects/<int:project_id>', methods=['GET'])
def get_project(project_id):
    """Get a single project."""
    # TODO: Implement with database
    return success_response({})


# ============================================
# 시스템 설정 API
# ============================================

@bp.route('/settings', methods=['GET'])
@jwt_required()
def get_system_settings():
    """시스템 설정 조회"""
    try:
        # 각 테이블에서 데이터 조회
        divisions = Division.query.filter_by(is_active=True).order_by(Division.order).all()
        departments = Department.query.filter_by(is_active=True).order_by(Department.order).all()
        processes = ProcessCategory.query.filter_by(is_active=True).order_by(ProcessCategory.order).all()
        project_domains = ProjectDomain.query.filter_by(is_active=True).order_by(ProjectDomain.order).all()
        task_categories = TaskCategory.query.filter_by(is_active=True).order_by(TaskCategory.order).all()
        statuses = TaskStatus.query.filter_by(is_active=True).order_by(TaskStatus.order).all()
        perf_categories = PerformanceCategory.query.filter_by(is_active=True).order_by(PerformanceCategory.order).all()
        perf_subcategories = PerformanceSubcategory.query.filter_by(is_active=True).order_by(PerformanceSubcategory.order).all()

        settings_data = {
            'divisions': [d.to_dict() for d in divisions],
            'departments': [d.to_dict() for d in departments],
            'processes': [p.to_dict() for p in processes],
            'projectDomains': [p.to_dict() for p in project_domains],
            'taskCategories': [t.to_dict() for t in task_categories],
            'statuses': [s.to_dict() for s in statuses],
            'performanceCategories': [p.to_dict() for p in perf_categories],
            'performanceSubcategories': [p.to_dict() for p in perf_subcategories]
        }

        # ModuleSettings에서 추가 설정 조회 (열 배치 설정 등)
        column_settings = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='groupedViewColumnSettings'
        ).first()
        if column_settings and column_settings.settings_data:
            settings_data['groupedViewColumnSettings'] = column_settings.settings_data

        # 피봇 설정 조회
        pivot_settings = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='pivotViewSettings'
        ).first()
        if pivot_settings and pivot_settings.settings_data:
            settings_data['pivotViewSettings'] = pivot_settings.settings_data

        # 조치 사항 설정 조회
        perf_evaluations = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='performanceEvaluations'
        ).first()
        if perf_evaluations and perf_evaluations.settings_data:
            settings_data['performanceEvaluations'] = perf_evaluations.settings_data
            print(f"[Settings GET] performanceEvaluations loaded: {perf_evaluations.settings_data}")
        else:
            settings_data['performanceEvaluations'] = []
            print(f"[Settings GET] performanceEvaluations not found in DB")

        # 로직 템플릿 조회
        logic_templates = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='logicTemplates'
        ).first()
        if logic_templates and logic_templates.settings_data:
            settings_data['logicTemplates'] = logic_templates.settings_data
        else:
            settings_data['logicTemplates'] = []

        # 보고 현황 설정 조회
        report_statuses = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='reportStatuses'
        ).first()
        if report_statuses and report_statuses.settings_data:
            settings_data['reportStatuses'] = report_statuses.settings_data
        else:
            settings_data['reportStatuses'] = []

        # 단위 환산 설정 조회
        unit_conversions = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='unitConversions'
        ).first()
        if unit_conversions and unit_conversions.settings_data:
            settings_data['unitConversions'] = unit_conversions.settings_data
        else:
            settings_data['unitConversions'] = []

        # 경영진 보고 설정 조회 (KPI 선택, 카드 선택, 단위 환산 활성화 등)
        exec_report_settings = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='executiveReportSettings'
        ).first()
        if exec_report_settings and exec_report_settings.settings_data:
            settings_data['executiveReportSettings'] = exec_report_settings.settings_data
        else:
            settings_data['executiveReportSettings'] = {}

        # 이슈 현황 사업부별 사무국 코멘트 조회 ({ 연도: { 사업부: 코멘트 } })
        issue_secretariat = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='issueSecretariatComments'
        ).first()
        if issue_secretariat and issue_secretariat.settings_data:
            settings_data['issueSecretariatComments'] = issue_secretariat.settings_data
        else:
            settings_data['issueSecretariatComments'] = {}

        # 결과 보고서 사무국 최종확인(인장) 상태 ({ 과제uuid: {status, by, byName, at, comment, hash} })
        report_confirms = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='reportConfirmations'
        ).first()
        if report_confirms and report_confirms.settings_data:
            settings_data['reportConfirmations'] = report_confirms.settings_data
        else:
            settings_data['reportConfirmations'] = {}

        # KPI 매트릭스 표시 설정 ({ excludedKpiIds: [1, 5, ...] })
        # 경영진 보고의 `executiveReportSettings.excludedKpis` 와 **일부러 나눠 둔다** —
        # 보고서에서 감춘 지표가 운영 화면(매트릭스)에서도 조용히 사라지면
        # 그 지표에 붙은 과제가 통째로 안 보인다. 목적이 다른 두 결정이다.
        kpi_matrix_settings = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='kpiMatrixSettings'
        ).first()
        if kpi_matrix_settings and kpi_matrix_settings.settings_data:
            settings_data['kpiMatrixSettings'] = kpi_matrix_settings.settings_data
        else:
            settings_data['kpiMatrixSettings'] = {}

        # KPI 기여방법 사전 ({ "<kpiDefinitionId>": ["방법A", "방법B"] })
        # 과제 편집창의 'DX KPI 연결' 이 이 목록에서 고르게 한다. 자유 텍스트로 두면
        # 같은 뜻이 사람마다 다른 문장이 되어 나중에 묶어 세지 못한다.
        kpi_methods = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='kpiContributionMethods'
        ).first()
        if kpi_methods and kpi_methods.settings_data:
            settings_data['kpiContributionMethods'] = kpi_methods.settings_data
        else:
            settings_data['kpiContributionMethods'] = {}

        # 결과 보고서 챔피언 보고 상태 ({ 과제uuid: {status, by, byName, at} })
        # 사무국 확인과 **별개 표시**다 — 둘 다 찍힐 수 있고 서로를 대체하지 않는다.
        champion_reports = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='championReports'
        ).first()
        if champion_reports and champion_reports.settings_data:
            settings_data['championReports'] = champion_reports.settings_data
        else:
            settings_data['championReports'] = {}

        # 단위 환산 기본 체크 상태 ({ 소스단위: 환산id }) — KPI 대시보드 상단 토글 기본값
        default_conversions = ModuleSettings.query.filter_by(
            module_name='digital_twin_dashboard',
            settings_key='defaultActiveConversions'
        ).first()
        if default_conversions and default_conversions.settings_data:
            settings_data['defaultActiveConversions'] = default_conversions.settings_data
        else:
            settings_data['defaultActiveConversions'] = {}

        return success_response(settings_data)

    except Exception as e:
        print(f"[Settings Error] Get settings failed: {str(e)}")
        return error_response(f'설정 조회 실패: {str(e)}', status_code=500)


@bp.route('/settings', methods=['PUT'])
@jwt_required()
def update_system_settings():
    """시스템 설정 저장"""
    try:
        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        # 디버그: 받은 데이터의 키 출력
        print(f"[Settings PUT] Received keys: {list(data.keys())}")
        if 'performanceEvaluations' in data:
            print(f"[Settings PUT] performanceEvaluations data: {data['performanceEvaluations']}")

        # 각 설정 항목 업데이트 (데이터가 명시적으로 전달된 경우에만)
        #
        # ⚠️ **부모 먼저.** 부서는 사업부를, 소분류는 대분류를 가리키는데, 화면이
        #    보내는 참조가 'vd' 같은 문자열일 수 있다. 부모를 먼저 동기화해서
        #    {화면 id → DB id} 맵을 만들고 그걸 자식에게 넘겨야 풀린다.
        #    (넘기지 않으면 참조가 None 이 되어 사업부 연결이 조용히 사라진다)
        division_map = _sync_divisions(data['divisions']) if 'divisions' in data else {}
        if 'departments' in data:
            _sync_departments(data['departments'], division_map)
        if 'processes' in data:
            _sync_processes(data['processes'])
        if 'projectDomains' in data:
            _sync_project_domains(data['projectDomains'])
        if 'taskCategories' in data:
            _sync_task_categories(data['taskCategories'])
        if 'statuses' in data:
            _sync_statuses(data['statuses'])
        perf_cat_map = (_sync_performance_categories(data['performanceCategories'])
                        if 'performanceCategories' in data else {})
        cascade_notes = []
        if 'performanceSubcategories' in data:
            # 소분류의 단위가 바뀌면 **그 소분류를 쓰는 성과들도 같이** 바꾼다.
            # 안 하면 그 성과들은 각자 저장될 때까지 옛 단위를 달고 있는데,
            # 화면이 그 칸을 잠그고 있어 사람이 고칠 수도 없다. (_cascade_perf_unit 참고)
            _before_units = _snapshot_subcategory_units()
            _sync_performance_subcategories(data['performanceSubcategories'], perf_cat_map)
            db.session.flush()
            cascade_notes = _cascade_perf_unit(_before_units)

        # 열 배치 설정 저장 (ModuleSettings 테이블 사용)
        if 'groupedViewColumnSettings' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'groupedViewColumnSettings',
                data['groupedViewColumnSettings'],
                '그룹별 보기 열 배치 설정'
            )

        # 피봇 설정 저장 (ModuleSettings 테이블 사용)
        if 'pivotViewSettings' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'pivotViewSettings',
                data['pivotViewSettings'],
                '피봇 보기 설정'
            )

        # 조치 사항 설정 저장 (ModuleSettings 테이블 사용)
        if 'performanceEvaluations' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'performanceEvaluations',
                data['performanceEvaluations'],
                '조치 사항 옵션 목록'
            )

        # 로직 템플릿 저장 (ModuleSettings 테이블 사용)
        if 'logicTemplates' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'logicTemplates',
                data['logicTemplates'],
                '성과 로직 계산 템플릿'
            )

        # 보고 현황 설정 저장 (ModuleSettings 테이블 사용)
        if 'reportStatuses' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'reportStatuses',
                data['reportStatuses'],
                '보고 현황 옵션 목록'
            )

        # 단위 환산 설정 저장 (ModuleSettings 테이블 사용)
        if 'unitConversions' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'unitConversions',
                data['unitConversions'],
                '단위 환산 계수 설정'
            )

        # 경영진 보고 설정 저장 (ModuleSettings 테이블 사용)
        if 'executiveReportSettings' in data:
            _save_module_setting(
                'digital_twin_dashboard',
                'executiveReportSettings',
                data['executiveReportSettings'],
                '경영진 보고 표시 설정 (KPI 선택, 카드 선택, 단위 환산)'
            )

        # 이슈 현황 사업부별 사무국 코멘트 저장 (매니저 이상만 허용)
        if 'issueSecretariatComments' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role not in ('admin', 'manager', 'dt_office') and not user.is_admin):
                return error_response('사무국 코멘트 수정 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'issueSecretariatComments',
                data['issueSecretariatComments'],
                '이슈 현황 사업부별 사무국 코멘트'
            )

        # 결과 보고서 사무국 최종확인(인장) — 사무국(admin, dt_office)만 허용
        if 'reportConfirmations' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role not in ('admin', 'dt_office') and not user.is_admin):
                return error_response('보고서 최종 확인 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'reportConfirmations',
                data['reportConfirmations'],
                '결과 보고서 사무국 최종확인(인장) 상태'
            )

        # KPI 매트릭스 표시 설정 — 관리자만
        if 'kpiMatrixSettings' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role != 'admin' and not user.is_admin):
                return error_response('KPI 표시 설정 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'kpiMatrixSettings',
                data['kpiMatrixSettings'],
                'KPI 매트릭스 표시 설정 (지표 선택)'
            )

        # KPI 기여방법 사전 — 사무국도 정의할 수 있어야 한다(지표 운영이 그쪽 일이다).
        # 매트릭스 '표시할 지표'(admin 전용)와 달리 **무엇이 보이느냐가 아니라 용어**라,
        # 잠가 두면 현장에서 자유 텍스트로 되돌아간다.
        if 'kpiContributionMethods' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role not in ('admin', 'dt_office') and not user.is_admin):
                return error_response('KPI 기여방법 편집 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'kpiContributionMethods',
                data['kpiContributionMethods'],
                'KPI 기여방법 사전 (지표별 기여 방법 목록)'
            )

        # 결과 보고서 챔피언 보고 — 사무국 확인과 **같은 권한**(admin, dt_office)
        if 'championReports' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role not in ('admin', 'dt_office') and not user.is_admin):
                # ⚠️ `status_code=` 를 **반드시 이름으로 준다.** 두 번째 위치 인자는
                #    `errors` 다(shared/responses.py). `error_response(msg, status_code=403)` 이라고
                #    쓰면 403 이 errors 에 들어가고 **HTTP 는 400 으로 나간다.**
                #    이 파일의 다른 권한 검사들이 그 상태다(§아래 주석 참조).
                return error_response('챔피언 보고 표시 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'championReports',
                data['championReports'],
                '결과 보고서 챔피언 보고 상태'
            )

        # 단위 환산 기본 체크 상태 저장 (KPI 대시보드) — 관리자만 허용
        if 'defaultActiveConversions' in data:
            user = User.query.get(get_jwt_identity())
            if not user or (user.role != 'admin' and not user.is_admin):
                return error_response('단위 환산 기본값 설정 권한이 없습니다.', status_code=403)
            _save_module_setting(
                'digital_twin_dashboard',
                'defaultActiveConversions',
                data['defaultActiveConversions'],
                '단위 환산 기본 체크 상태 (KPI 대시보드)'
            )

        db.session.commit()

        # 서버가 함께 바꾼 것은 **반드시 알린다.** 소분류 하나를 고쳤는데 성과 수백 건이
        # 같이 바뀔 수 있어서, 누른 사람이 그 사실을 모르면 안 된다.
        msg = '설정이 저장되었습니다.'
        if cascade_notes:
            msg += ' ' + ' '.join(cascade_notes)
        return success_response({'message': msg, 'cascaded': cascade_notes})

    except Exception as e:
        db.session.rollback()
        print(f"[Settings Error] Update settings failed: {str(e)}")
        return error_response(f'설정 저장 실패: {str(e)}', status_code=500)


def _save_module_setting(module_name, settings_key, settings_data, description=''):
    """ModuleSettings 테이블에 설정 저장 (upsert)"""
    print(f"[_save_module_setting] Saving {module_name}:{settings_key}")
    print(f"[_save_module_setting] Data: {settings_data}")

    existing = ModuleSettings.query.filter_by(
        module_name=module_name,
        settings_key=settings_key
    ).first()

    if existing:
        print(f"[_save_module_setting] Updating existing record (id={existing.id})")
        existing.settings_data = settings_data
        existing.description = description
        # SQLAlchemy JSON 변경 감지를 위해 flag_modified 사용
        flag_modified(existing, 'settings_data')
    else:
        print(f"[_save_module_setting] Creating new record")
        new_setting = ModuleSettings(
            module_name=module_name,
            settings_key=settings_key,
            settings_data=settings_data,
            description=description
        )
        db.session.add(new_setting)

    print(f"[_save_module_setting] Done")


# ─────────────────────────────────────────────────────────────────────────────
# 설정 동기화 (사업부·부서·프로세스 …)
#
# 2026-08-02 재작성. 예전 방식이 만든 사고
#     각 `_sync_*` 가 **id 가 숫자일 때만** 기존 항목으로 인정했다.
#         if item_id and str(item_id).isdigit(): ... 갱신
#         else:                                  ... 새로 만들기
#     그런데 화면의 기본 설정(sampleData.js)은 문자열 id 다.
#         { id: 'vd_mecha_solution', name: 'Mecha Solution', divisionId: 'vd' }
#     서버 설정이 비어 있으면(신규 설치·DB 초기화) 화면이 이 기본값을 들고 있다가
#     그대로 저장한다. 그러면
#         ① 전부 새 행이 되고, 직전 세대는 마지막 줄에서 비활성으로 내려간다
#            → 저장 한 번에 설정 테이블이 통째로 한 벌 복제된다
#         ② `divisionId: 'vd'` 도 숫자가 아니라서 **None 이 되어 사업부 연결이 사라진다**
#
#     실측(개발서버 2026-08-02): 부서 91개 중 88개가 사업부 없음. 설정 테이블 7개가
#     전부 정확히 3벌. 2025-11-24 10:06:34 / 10:07:07 / 2026-01-26 03:51:31 —
#     **33초 간격의 저장 두 번**이 두 세대를 만든 게 시각에 그대로 남아 있다.
#
#     ②가 조용하다. 집계는 과제가 들고 있는 부서 '문자열' 로 그룹을 만들어서 총계는
#     맞고, 사업부로 접을 때만 '공통' 으로 샌다. 게다가 그 부서에 사람이 없으면
#     '공통' 자체가 화면에서 빠져(`totalCount > 0` 필터) 아무 증상도 안 보인다.
#     manager 는 더 나쁘다 — `actor_division_id` 가 사업부를 못 풀면 fail closed 라
#     **자기 사업부 과제를 못 고친다.** 원인이 설정 데이터라 추적이 어렵다.
#
# 그래서 규칙을 하나로 모았다 (`_sync_setting_rows`)
#     · id 가 숫자면 그 행. **아니면 이름으로 찾는다** — 문자열 id 여도 복제되지 않는다
#     · 찾은 행은 **is_active=True 로 되살린다.** 예전에는 비활성 행을 id 로 찾아
#       갱신만 하고 활성화는 안 해서, 되살린 항목이 화면에서 계속 안 보였다
#     · 부모 참조(divisionId·categoryId)는 **같은 요청 안의 부모 항목**과 먼저 맞춘다.
#       'vd' 는 divisions 에 `{id:'vd', name:'VD'}` 로 같이 오므로 이걸로 풀린다
#       (이름만으로 맞추면 `medical` → `의료기기` 는 못 푼다)
# ─────────────────────────────────────────────────────────────────────────────


def _norm_key(value):
    """이름 비교용 정규화. 앞뒤 공백과 대소문자만 무시한다."""
    return (str(value) if value is not None else '').strip().lower()


def _sync_setting_rows(model, items, apply_fields, default_name):
    """
    설정 한 종류를 동기화하고 `{화면이 보낸 id: DB id}` 를 돌려준다.

    돌려주는 맵은 자식(부서·소분류)이 부모를 찾을 때 쓴다.
    """
    rows = model.query.all()
    by_id = {r.id: r for r in rows}
    # 이름 → 후보들. **활성 우선, 그다음 최신 id** — 되살릴 때 옛 세대보다
    # 지금 쓰는 행을 먼저 잡아야 한다.
    by_name = {}
    for r in sorted(rows, key=lambda r: (not bool(r.is_active), -r.id)):
        by_name.setdefault(_norm_key(r.name), []).append(r)

    used, client_map = set(), {}

    for idx, item in enumerate(items or []):
        raw_id = item.get('id')
        raw_key = str(raw_id).strip() if raw_id not in (None, '') else None

        row = None
        if raw_key and raw_key.isdigit():
            cand = by_id.get(int(raw_key))
            if cand is not None and cand.id not in used:
                row = cand
        if row is None:
            for cand in by_name.get(_norm_key(item.get('name')), []):
                if cand.id not in used:
                    row = cand
                    break
        if row is None:
            row = model(name=(item.get('name') or default_name))
            db.session.add(row)

        apply_fields(row, item)
        row.order = idx
        row.is_active = True
        db.session.flush()

        used.add(row.id)
        if raw_key:
            client_map[raw_key] = row.id

    # 이번에 안 온 것은 내린다. 지우지 않는 이유는 과제에 남은 이름과 대조할
    # 근거를 남기기 위해서다. (빈 집합이면 notin_ 이 이상하게 도니 자리표를 넣는다)
    model.query.filter(
        model.id.notin_(used or {0}),
        model.is_active.is_(True)
    ).update({'is_active': False}, synchronize_session=False)

    return client_map


def _resolve_parent(raw, parent_map, parent_model):
    """
    자식이 가리키는 부모의 DB id.

    ① 같은 요청 안의 부모 항목 (화면이 준 id 그대로: 'vd', '18', …)
    ② 숫자면 그 id 로 직접 — 단 **그 행이 비활성이면 같은 이름의 활성 행으로 옮긴다**
    ③ 이름으로 (활성 우선)
    못 풀면 None. **값을 지어내지 않는다.**

    ②의 단서가 중요하다. 이 DB 에는 같은 이름의 사업부가 세대별로 여러 벌 있다
    (MX 가 id 1·9·17, 그중 17 만 활성 — 옛 저장 버그가 남긴 잔재).
    비활성 id 를 그대로 물리면 화면의 사업부 목록(활성만 내려간다)에 없는 값이 되어
    **부서가 '공통' 으로 새고, manager 권한 비교도 어긋난다.** 조용해서 더 나쁘다.
    """
    if raw in (None, ''):
        return None
    key = str(raw).strip()

    if key in parent_map:
        return parent_map[key]

    rows = parent_model.query.all()
    active_by_name = {_norm_key(r.name): r.id for r in rows if r.is_active}

    if key.isdigit():
        row = next((r for r in rows if r.id == int(key)), None)
        if row is not None:
            if row.is_active:
                return row.id
            # 비활성 행을 가리켰다.
            #   · 같은 이름의 활성 행이 있으면 그쪽으로 승격 (세대 중복 잔재를 가리킨 경우)
            #   · 없으면 **원래 값을 그대로 둔다.**
            #     ⚠️ 여기서 None 을 주면 관리자가 손으로 넣은 사업부 지정이 저장 한 번에
            #        지워진다. 잠시 내려간 사업부일 수도 있고(되살리면 링크가 살아난다),
            #        무엇보다 **모르겠다고 데이터를 버리면 안 된다.**
            #        (2026-08-02 왕복 시험에서 실측 — 이 분기가 division_id 를 날렸다)
            return active_by_name.get(_norm_key(row.name), row.id)

    hit = active_by_name.get(_norm_key(key))
    if hit is not None:
        return hit
    # 활성 중에 없으면 비활성이라도 (이름이 완전히 새로운 경우가 아니라면)
    for row in sorted(rows, key=lambda r: -r.id):
        if _norm_key(row.name) == _norm_key(key):
            return row.id
    return None


def _sync_divisions(items):
    """사업부 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.color = item.get('color', row.color or '#64748B')
        row.description = item.get('description', row.description)
    return _sync_setting_rows(Division, items, apply, '새 사업부')


def _sync_departments(items, division_map=None):
    """부서 동기화. `divisionId` 는 같은 요청의 사업부 → 숫자 → 이름 순으로 푼다."""
    division_map = division_map or {}

    def apply(row, item):
        row.name = item.get('name', row.name)
        row.division_id = _resolve_parent(item.get('divisionId'), division_map, Division)
        row.description = item.get('description', row.description)
    return _sync_setting_rows(Department, items, apply, '새 부서')


def _sync_processes(items):
    """프로세스 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.description = item.get('description', row.description)
    return _sync_setting_rows(ProcessCategory, items, apply, '새 프로세스')


def _sync_project_domains(items):
    """과제 영역 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.description = item.get('description', row.description)
    return _sync_setting_rows(ProjectDomain, items, apply, '새 과제 영역')


def _sync_task_categories(items):
    """과제 구분 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.description = item.get('description', row.description)
    return _sync_setting_rows(TaskCategory, items, apply, '새 과제 구분')


def _sync_statuses(items):
    """진행상태 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.color = item.get('color', row.color or '#64748B')
        row.description = item.get('description', row.description)
    return _sync_setting_rows(TaskStatus, items, apply, '새 상태')


def _sync_performance_categories(items):
    """성과 대분류 동기화"""
    def apply(row, item):
        row.name = item.get('name', row.name)
        row.color = item.get('color', row.color or '#64748B')
        row.description = item.get('description', row.description)
    return _sync_setting_rows(PerformanceCategory, items, apply, '새 성과 분류')


def _snapshot_subcategory_units():
    """(대분류명, 소분류명) → (단위, 달성형여부). 연쇄 갱신 전후를 비교하는 데 쓴다."""
    cats = {c.id: c.name for c in PerformanceCategory.query.all()}
    return {(cats.get(s.category_id), s.name): (s.unit, bool(s.is_achievement_type))
            for s in PerformanceSubcategory.query.all()}


def _cascade_perf_unit(before: dict) -> list:
    """
    소분류의 `단위`·`달성형여부` 가 바뀌었으면 **그 소분류를 쓰는 성과들도 같이** 고친다.

    왜 필요한가
        성과의 `단위` 는 소분류가 정하는 값이고, 성과 행에 있는 것은 **사본**이다
        (`_derive_perf_from_subcategory`). 그런데 그 파생은 **성과를 저장할 때만** 돈다.
        여기서 소분류 단위를 바꾸면, 그 소분류를 쓰는 기존 성과들은 **각각 저장될
        때까지 옛 단위를 달고 있게 된다.** 화면은 그 칸을 읽기 전용으로 잠그므로
        사람이 손댈 수도 없다 — 2026-08-05 에 고친 그 버그가 그대로 재발한다.
        (실제로 그날 `억`→`억원` 으로 바꾸면서 성과 30건을 따로 맞춰야 했다.)

    ⚠️ **조용히 하지 않는다.** 바꾼 건수를 돌려주고 응답에 실어 보낸다 — 소분류 하나를
       고쳤는데 수백 건이 함께 바뀔 수 있으므로, 누른 사람이 그 사실을 알아야 한다.

    ⚠️ **이름이 바뀐 소분류는 여기서 못 잡는다.** 성과 행은 소분류를 **이름 문자열**로
       들고 있어서, 이름을 바꾸면 짝을 찾지 못한다(그 성과는 파생 대상에서 빠져 옛 값을
       유지한다 — 조용히 틀린 값이 되지는 않는다). 이름변경 파급은 별도 과제다.
    """
    from app.modules.digital_twin_dashboard.history import record_performance_history
    from app.modules.digital_twin_dashboard.models_v2 import Dt2Performance

    cats = {c.id: c.name for c in PerformanceCategory.query.all()}
    notes = []
    for s in PerformanceSubcategory.query.filter(
            PerformanceSubcategory.is_active.is_(True)).all():
        key = (cats.get(s.category_id), s.name)
        old = before.get(key)
        now = (s.unit, bool(s.is_achievement_type))
        if old is None or old == now:
            continue

        rows = Dt2Performance.query.filter(
            Dt2Performance.category == key[0],
            Dt2Performance.subcategory == key[1],
            Dt2Performance.is_deleted.isnot(True)).all()
        n = 0
        for p in rows:
            touched = False
            # 소분류에 단위가 비어 있으면 손대지 않는다 — 그때는 커스텀 입력을 허용한다.
            if now[0] and (p.unit or '') != now[0]:
                p.unit = now[0]
                touched = True
            if bool(p.is_achievement_type) != now[1]:
                p.is_achievement_type = now[1]
                touched = True
            if touched:
                p.row_version = (p.row_version or 1) + 1
                record_performance_history(p, source='server')
                n += 1
        if n:
            notes.append(
                f"소분류 '{key[1]}' 의 단위를 {old[0]!r} → {now[0]!r} 로 바꾸면서 "
                f"그 소분류를 쓰는 성과 {n}건의 단위도 함께 바꿨습니다.")
            print(f"[Settings] 단위 연쇄 갱신 — {key[0]}/{key[1]}: "
                  f"{old[0]!r} → {now[0]!r} · 성과 {n}건")
    return notes


def _sync_performance_subcategories(items, category_map=None):
    """
    성과 소분류 동기화.

    `category_id` 는 NOT NULL 이라 못 풀면 그 항목을 건너뛴다(예전과 같다).
    다만 이제 문자열 id 도 부모 맵으로 풀리므로 **그냥 사라지는 일이 거의 없다** —
    예전에는 `categoryId: 'productivity'` 같은 값이 전부 스킵돼 소분류가 통째로
    비활성이 됐다.
    """
    category_map = category_map or {}

    kept, cat_of = [], {}
    for item in (items or []):
        cid = _resolve_parent(item.get('categoryId'), category_map, PerformanceCategory)
        if cid is None:
            print(f"[Settings] 성과 소분류 '{item.get('name')}' — 대분류를 못 찾아 건너뜀"
                  f" (categoryId={item.get('categoryId')!r})")
            continue
        cat_of[id(item)] = cid
        kept.append(item)

    def apply(row, item):
        row.name = item.get('name', row.name)
        row.category_id = cat_of[id(item)]
        row.unit = item.get('unit', row.unit)
        row.description = item.get('description', row.description)
        row.is_achievement_type = item.get('isAchievementType', False)

    return _sync_setting_rows(PerformanceSubcategory, kept, apply, '새 소분류')


# ============================================
# 서버 데이터 동기화 API
# ============================================

@bp.route('/data', methods=['GET'])
@jwt_required()
def get_dashboard_data():
    """서버 데이터 조회 (버전 포함)"""
    try:
        # 가장 최신 데이터 조회 (ID=1로 싱글톤 패턴 사용)
        data = DashboardData.query.first()

        if not data:
            # 데이터가 없으면 빈 데이터 반환
            return success_response({
                'version': 0,
                'projects': [],
                'performances': [],
                'metadata': {},
                'last_modified_by_name': None,
                'updated_at': None
            })

        # 완전 삭제된 프로젝트 및 삭제된 성과는 조회에서 제외
        result = data.to_dict()
        if result.get('projects'):
            result['projects'] = [
                p for p in result['projects']
                if not p.get('_permanentlyDeleted')
            ]
        if result.get('performances'):
            result['performances'] = [
                p for p in result['performances']
                if not p.get('_deleted')
            ]

        return success_response(result)

    except Exception as e:
        print(f"[Data Error] Get data failed: {str(e)}")
        return error_response(f'데이터 조회 실패: {str(e)}', status_code=500)


# ============================================
# 스냅샷 보존 정책 (Phase 1-1)
# ============================================
# 배경: 업서트마다 전체 데이터 복사본을 스냅샷으로 쌓아 운영에서 4,896건 / 디스크 25.2 GB 까지 누적됐다.
#       업서트 경로의 자동 스냅샷은 제거했고, 남은 자동 스냅샷(전체 덮어쓰기 직전 백업)은 아래 정책으로 제한한다.

SNAPSHOT_AUTO_KEEP_COUNT = 30                      # 자동 생성분 중 최신 N건 보존
SNAPSHOT_KEEP_DAYS = 7                             # 타입 무관, 최근 D일 이내는 보존
SNAPSHOT_PROTECTED_TYPES = ('manual', 'upload')    # 사람이 의도한 스냅샷은 전량 보존


def _prune_auto_snapshots(keep_count=SNAPSHOT_AUTO_KEEP_COUNT, keep_days=SNAPSHOT_KEEP_DAYS):
    """
    자동 생성 스냅샷을 보존 정책에 맞게 정리한다.

    **커밋하지 않는다.** 호출자의 트랜잭션에 참여하므로, 호출자가 롤백하면 이 정리도 함께 취소된다.

    보존 규칙 — 아래 중 하나라도 해당하면 남긴다:
      1. snapshot_type 이 manual / upload   (사람이 만든 것 + 덮어쓰기 직전 백업)
      2. 최근 keep_days 일 이내 생성        (타입 무관)
      3. 자동 생성분 중 최신 keep_count 건

    Returns: 삭제된 건수 (실패 시 0)
    """
    try:
        cutoff = datetime.utcnow() - timedelta(days=keep_days)

        # 규칙 3 — 보존할 자동 스냅샷 id 목록
        recent_ids = [
            row.id for row in (
                DashboardSnapshot.query
                .filter(~DashboardSnapshot.snapshot_type.in_(SNAPSHOT_PROTECTED_TYPES))
                .order_by(DashboardSnapshot.created_at.desc())
                .limit(keep_count)
                .all()
            )
        ]

        # 규칙 1·2에 해당하지 않고 규칙 3에도 없는 것만 삭제
        q = DashboardSnapshot.query.filter(
            ~DashboardSnapshot.snapshot_type.in_(SNAPSHOT_PROTECTED_TYPES),
            DashboardSnapshot.created_at < cutoff,
        )
        if recent_ids:
            q = q.filter(~DashboardSnapshot.id.in_(recent_ids))

        deleted = q.delete(synchronize_session=False)
        if deleted:
            print(f"[Snapshot] 보존 정책으로 자동 스냅샷 {deleted}건 정리 "
                  f"(keep={keep_count}건, {keep_days}일)")
        return deleted
    except Exception as exc:
        # 정리는 부가 기능이다. 실패해도 본 작업(업로드)을 막지 않는다.
        print(f"[Snapshot] 정리 건너뜀: {exc}")
        return 0


@bp.route('/data/upload', methods=['POST'])
@jwt_required()
def upload_dashboard_data():
    """
    서버 데이터 업로드 (버전 체크로 충돌 감지)
    - client_version이 server_version과 일치해야 업로드 가능
    - 불일치 시 409 Conflict 반환
    """
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        req_data = request.get_json()
        if not req_data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        client_version = req_data.get('version', 0)
        projects = req_data.get('projects', [])
        performances = req_data.get('performances', [])
        metadata = req_data.get('metadata', {})
        activity_logs = req_data.get('activityLogs', [])
        snapshot = req_data.get('snapshot')
        force_upload = req_data.get('forceUpload', False)  # 강제 업로드 옵션

        # 현재 서버 데이터 조회
        server_data = DashboardData.query.first()
        server_version = server_data.version if server_data else 0

        # 버전 충돌 체크 (강제 업로드가 아닌 경우)
        if not force_upload and client_version != server_version:
            # ⚠️ 이 호출은 **HTTP 200 을 내보내고 있었다.** 위치 인자가
            #    (message, errors, status_code) 라서 `409` 가 errors 로, 뒤의 dict 가
            #    status_code 로 들어갔고, Flask 는 dict 를 **헤더**로 해석한다.
            #    그래서 화면의 `if (response.ok)` 가 **버전 충돌을 성공으로 처리**했다.
            #    (2026-08-06 발견 — 충돌 상세도 응답 본문이 아니라 헤더로 새고 있었다)
            return error_response(
                '버전 충돌이 발생했습니다. 다른 사용자가 데이터를 수정했습니다.',
                status_code=409,
                errors={
                    'conflict': True,
                    'server_version': server_version,
                    'client_version': client_version,
                    'last_modified_by': server_data.last_modified_by_name if server_data else None,
                    'last_modified_at': iso_kst(server_data.updated_at) if server_data and server_data.updated_at else None
                }
            )

        # 업로드 전 자동 스냅샷 생성 (기존 데이터가 있는 경우)
        # 전체 덮어쓰기는 되돌리기 어려우므로 여기서는 자동 스냅샷을 유지한다.
        # (업서트 경로의 자동 스냅샷은 Phase 1-1에서 제거됨)
        if server_data and server_data.projects:
            auto_snapshot = DashboardSnapshot(
                name=f'자동 백업 (v{server_version})',
                description=f'업로드 전 자동 백업 - {user_name}에 의해 덮어쓰기됨',
                snapshot_data={
                    'projects': server_data.projects,
                    'performances': server_data.performances,
                    'metadata': server_data.data_metadata
                },
                data_version=server_version,
                created_by=server_data.last_modified_by,
                created_by_name=server_data.last_modified_by_name,
                snapshot_type='auto'
            )
            db.session.add(auto_snapshot)

            # [Phase 1-1] 자동 스냅샷이 무한정 쌓이지 않도록 즉시 정리
            #   전체 덮어쓰기는 빈도가 낮으므로 이 시점에 정리해도 부담이 없다.
            _prune_auto_snapshots()

        # 데이터 저장/업데이트
        new_version = server_version + 1

        if server_data:
            server_data.version = new_version
            server_data.projects = projects
            server_data.performances = performances
            server_data.data_metadata = metadata
            server_data.last_modified_by = user_id
            server_data.last_modified_by_name = user_name
        else:
            server_data = DashboardData(
                version=new_version,
                projects=projects,
                performances=performances,
                data_metadata=metadata,
                last_modified_by=user_id,
                last_modified_by_name=user_name
            )
            db.session.add(server_data)

        # 스냅샷 저장 (요청에 포함된 경우)
        if snapshot:
            new_snapshot = DashboardSnapshot(
                name=snapshot.get('name', f'업로드 스냅샷 v{new_version}'),
                description=snapshot.get('description', ''),
                snapshot_data={
                    'projects': projects,
                    'performances': performances,
                    'metadata': metadata
                },
                data_version=new_version,
                created_by=user_id,
                created_by_name=user_name,
                snapshot_type='upload'
            )
            db.session.add(new_snapshot)

        # 활동 로그 저장
        for log in activity_logs:
            activity_log = DashboardActivityLog(
                action=log.get('action', 'UPLOAD'),
                target_type=log.get('targetType', 'DATA'),
                target_id=log.get('targetId'),
                target_name=log.get('targetName'),
                changes=log.get('changes'),
                summary=log.get('summary'),
                user_id=user_id,
                user_name=user_name,
                data_version=new_version,
                source='server'
            )
            db.session.add(activity_log)

        # 업로드 자체에 대한 로그 (동기화 정보만 기록)
        upload_log = DashboardActivityLog(
            action='UPLOAD',
            target_type='DATA',
            target_name=f'v{server_version} → v{new_version} 덮어쓰기',
            summary=f'{user_name}이(가) 서버에 데이터를 덮어씀 (v{new_version}, 과제 {len(projects)}개, 성과 {len(performances)}개)',
            changes={
                'version': {'before': server_version, 'after': new_version},
                'projectCount': len(projects),
                'performanceCount': len(performances)
            },
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(upload_log)

        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('upload')

        return success_response({
            'message': '데이터가 성공적으로 업로드되었습니다.',
            'version': new_version,
            'uploaded_by': user_name
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Upload failed: {str(e)}")
        return error_response(f'데이터 업로드 실패: {str(e)}', status_code=500)


@bp.route('/data/upsert', methods=['POST'])
@jwt_required()
def upsert_dashboard_data():
    """
    서버 데이터 업서트 (업데이트 또는 추가)
    - UUID 기준으로 기존 데이터 업데이트 또는 신규 추가
    - 기존 데이터를 삭제하지 않고 병합함
    """
    try:
        import time as _time  # [타이밍 계측] 저장 병목 진단용 (동작 변경 없음)
        _t0 = _time.perf_counter()
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        req_data = request.get_json()
        if not req_data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        client_version = req_data.get('version', 0)
        new_projects = req_data.get('projects', [])
        new_performances = req_data.get('performances', [])
        metadata = req_data.get('metadata', {})
        activity_logs = req_data.get('activityLogs', [])
        snapshot = req_data.get('snapshot')
        _t_parsed = _time.perf_counter()  # [타이밍] 요청 본문 파싱 완료

        # 현재 서버 데이터 조회
        server_data = DashboardData.query.first()
        server_version = server_data.version if server_data else 0

        # 기존 데이터 가져오기
        existing_projects = server_data.projects if server_data and server_data.projects else []
        existing_performances = server_data.performances if server_data and server_data.performances else []
        _t_loaded = _time.perf_counter()  # [타이밍] 서버 싱글톤 행 읽기(대용량 JSON 역직렬화) 완료

        # UUID 기준으로 병합 (수정 시간 기반 - 더 최신 데이터만 반영)
        # 프로젝트 병합 - UUID를 우선으로 사용하여 다중 사용자 환경에서 ID 충돌 방지
        # uuid -> project 맵 (uuid가 있는 경우)
        project_uuid_map = {p.get('uuid'): p for p in existing_projects if p.get('uuid')}
        # id -> project 맵 (uuid가 없는 레거시 데이터용 - fallback)
        project_id_map = {p.get('id'): p for p in existing_projects if p.get('id') and not p.get('uuid')}
        # id -> project 맵 (모든 프로젝트 - ID 기반 중복 체크용)
        all_project_by_id = {p.get('id'): p for p in existing_projects if p.get('id')}
        # 기존 ID 집합 (ID 충돌 체크용)
        existing_project_ids = {p.get('id') for p in existing_projects if p.get('id')}

        updated_count = 0
        added_count = 0
        skipped_count = 0
        skipped_deleted_count = 0  # 완전 삭제된 항목으로 인한 스킵 카운트
        renumbered_count = 0  # ID 재번호 부여된 프로젝트 수

        import re  # ID 파싱용

        for project in new_projects:
            project_uuid = project.get('uuid')
            project_id = project.get('id')

            # 1. UUID로 먼저 매칭 시도
            if project_uuid and project_uuid in project_uuid_map:
                # 기존 프로젝트 존재 (UUID 매칭)
                existing = project_uuid_map[project_uuid]

                # 완전 삭제된 프로젝트인 경우 - 삭제 시간과 비교
                if existing.get('_permanentlyDeleted'):
                    deleted_at = existing.get('_permanentlyDeletedAt') or ''
                    new_updated = project.get('updatedAt') or project.get('updated_at') or ''

                    if deleted_at >= new_updated:
                        # 삭제 시간이 더 최신이면 업로드 무시 (삭제 유지)
                        skipped_deleted_count += 1
                        continue
                    else:
                        # 업로드 데이터가 삭제 이후에 수정된 경우 - 그래도 삭제 유지 (완전 삭제는 복구 불가)
                        skipped_deleted_count += 1
                        continue

                # 일반적인 경우 - updatedAt 비교하여 더 최신인 경우만 업데이트
                existing_updated = existing.get('updatedAt') or existing.get('updated_at') or ''
                new_updated = project.get('updatedAt') or project.get('updated_at') or ''

                if new_updated >= existing_updated:
                    project_uuid_map[project_uuid] = project
                    updated_count += 1
                else:
                    # 서버 데이터가 더 최신이므로 스킵
                    skipped_count += 1

            # 2. UUID가 없거나 UUID로 못 찾은 경우 - ID로 fallback 매칭 시도
            elif project_id and project_id in all_project_by_id:
                existing = all_project_by_id[project_id]

                # 완전 삭제된 프로젝트인 경우
                if existing.get('_permanentlyDeleted'):
                    # UUID가 다르면 완전히 다른 신규 프로젝트가 같은 ID로 생성된 것
                    # → 새 ID를 부여하고 신규 추가로 처리
                    if project_uuid and project_uuid != existing.get('uuid'):
                        division = project.get('사업부', 'VD')
                        max_num = 0
                        for eid in existing_project_ids:
                            if eid and isinstance(eid, str):
                                match = re.match(rf'^{re.escape(division)}-(\d+)$', eid)
                                if match:
                                    max_num = max(max_num, int(match.group(1)))
                        new_id = f'{division}-{max_num + 1}'
                        project = {**project, 'id': new_id}
                        renumbered_count += 1
                        existing_project_ids.add(new_id)
                        project_uuid_map[project_uuid] = project
                        added_count += 1
                        print(f"[Upsert] 삭제된 프로젝트와 ID 충돌 → 신규 추가: {project_id} -> {new_id} (uuid: {project_uuid})")
                    else:
                        skipped_deleted_count += 1
                    continue

                # updatedAt 비교
                existing_updated = existing.get('updatedAt') or existing.get('updated_at') or ''
                new_updated = project.get('updatedAt') or project.get('updated_at') or ''

                if new_updated >= existing_updated:
                    # 기존 프로젝트의 UUID가 있으면 uuid_map 업데이트, 없으면 id_map 업데이트
                    existing_uuid = existing.get('uuid')
                    if existing_uuid:
                        # 새 데이터에 UUID가 없으면 기존 UUID 유지
                        if not project_uuid:
                            project = {**project, 'uuid': existing_uuid}
                        project_uuid_map[existing_uuid] = project
                    else:
                        project_id_map[project_id] = project
                    updated_count += 1
                else:
                    skipped_count += 1

            # 3. 신규 프로젝트 추가
            else:
                # ID 충돌 체크: 같은 ID가 이미 존재하고 UUID가 다르면 새 ID 부여 (여기 도달하면 안됨)
                if project_id and project_id in existing_project_ids:
                    # ID 충돌 발생 - 새로운 번호 부여
                    division = project.get('사업부', 'VD')
                    max_num = 0
                    for eid in existing_project_ids:
                        if eid and isinstance(eid, str):
                            match = re.match(rf'^{re.escape(division)}-(\d+)$', eid)
                            if match:
                                max_num = max(max_num, int(match.group(1)))
                    new_id = f'{division}-{max_num + 1}'

                    project = {**project, 'id': new_id}
                    renumbered_count += 1
                    existing_project_ids.add(new_id)

                # UUID 맵에 추가 (신규 프로젝트)
                if project_uuid:
                    project_uuid_map[project_uuid] = project
                else:
                    # UUID가 없는 경우 ID 맵에 추가
                    if project_id:
                        project_id_map[project_id] = project

                existing_project_ids.add(project.get('id'))
                added_count += 1

        # 병합된 프로젝트 목록 생성 (UUID 맵 + ID 맵)
        merged_projects = list(project_uuid_map.values()) + list(project_id_map.values())

        if skipped_count > 0:
            print(f"[Upsert] {skipped_count}개 프로젝트가 서버 데이터가 더 최신이어서 스킵됨")
        if skipped_deleted_count > 0:
            print(f"[Upsert] {skipped_deleted_count}개 프로젝트가 완전 삭제된 상태여서 업로드 스킵됨")

        # 성과 병합 (uuid를 우선으로 사용하여 다중 사용자 환경에서 충돌 방지)
        # uuid -> performance 맵 (기존 데이터)
        performance_uuid_map = {p.get('uuid'): p for p in existing_performances if p.get('uuid')}
        # id -> performance 맵 (id 중복 체크용)
        existing_ids = {p.get('id') for p in existing_performances if p.get('id')}
        # 삭제된 성과 uuid 집합 (부활 방지용)
        deleted_perf_uuids = {p.get('uuid') for p in existing_performances if p.get('uuid') and p.get('_deleted')}
        deleted_perf_ids = {p.get('id') for p in existing_performances if p.get('id') and p.get('_deleted')}

        perf_updated_count = 0
        perf_added_count = 0
        perf_skipped_count = 0
        perf_skipped_deleted_count = 0
        perf_renumbered_count = 0
        renumbered_performances = []  # 재번호 부여된 성과 목록

        for performance in new_performances:
            perf_uuid = performance.get('uuid')
            perf_id = performance.get('id')

            if perf_uuid and perf_uuid in performance_uuid_map:
                existing_perf = performance_uuid_map[perf_uuid]

                # 삭제된 성과인 경우 - 삭제 시간과 비교하여 부활 방지
                if existing_perf.get('_deleted'):
                    deleted_at = existing_perf.get('_deletedAt') or ''
                    new_updated = performance.get('updatedAt') or performance.get('updated_at') or ''

                    if deleted_at >= new_updated:
                        # 삭제 시간이 더 최신이면 업로드 무시 (삭제 유지)
                        perf_skipped_deleted_count += 1
                        continue
                    else:
                        # 업로드 데이터가 삭제 이후에 수정된 경우에도 삭제 유지
                        perf_skipped_deleted_count += 1
                        continue

                # uuid가 일치 - updatedAt 비교하여 더 최신인 경우만 업데이트
                existing_updated = existing_perf.get('updatedAt') or existing_perf.get('updated_at') or ''
                new_updated = performance.get('updatedAt') or performance.get('updated_at') or ''

                if new_updated >= existing_updated:
                    performance_uuid_map[perf_uuid] = performance
                    perf_updated_count += 1
                else:
                    # 서버 데이터가 더 최신이므로 스킵
                    perf_skipped_count += 1
            else:
                # uuid로 매칭 안 됨 - 삭제된 성과인지 ID로도 확인 (부활 방지)
                if perf_uuid and perf_uuid in deleted_perf_uuids:
                    perf_skipped_deleted_count += 1
                    continue
                if perf_id and perf_id in deleted_perf_ids:
                    perf_skipped_deleted_count += 1
                    continue

                # 신규 성과 추가
                # id 중복 체크: 같은 id가 이미 존재하고 uuid가 다르면 새 id 부여
                if perf_id and perf_id in existing_ids:
                    # id 중복 발생 - 새로운 번호 부여
                    max_num = 0
                    for eid in existing_ids:
                        match = re.match(r'performance-(\d+)', eid or '')
                        if match:
                            max_num = max(max_num, int(match.group(1)))
                    new_id = f'performance-{max_num + 1}'

                    # 재번호 정보 저장 (사용자에게 알림용)
                    renumbered_performances.append({
                        'old_id': perf_id,
                        'new_id': new_id,
                        'name': performance.get('성과항목', ''),
                        'uuid': perf_uuid
                    })

                    performance = {**performance, 'id': new_id}
                    perf_renumbered_count += 1
                    print(f"[Upsert] Performance ID 충돌: {perf_id} -> {new_id} (uuid: {perf_uuid})")

                # 맵에 추가
                key = perf_uuid or performance.get('id')
                performance_uuid_map[key] = performance
                existing_ids.add(performance.get('id'))
                perf_added_count += 1

        merged_performances = list(performance_uuid_map.values())
        _t_merged = _time.perf_counter()  # [타이밍] 병합 로직 완료

        if perf_skipped_deleted_count > 0:
            print(f"[Upsert] {perf_skipped_deleted_count}개 성과가 삭제된 상태여서 업로드 스킵됨")

        if perf_renumbered_count > 0:
            print(f"[Upsert] {perf_renumbered_count}개 성과 ID가 재번호 부여됨: {renumbered_performances}")

        if perf_skipped_count > 0:
            print(f"[Upsert] {perf_skipped_count}개 성과가 서버 데이터가 더 최신이어서 스킵됨")

        # [Phase 1-1] 업서트 시 자동 스냅샷 생성 중단
        #
        # 이전에는 저장할 때마다 전체 데이터의 복사본을 스냅샷으로 한 건씩 INSERT 했다.
        # 운영 실측(2026-07-28): 데이터 1벌이 약 37 MB → 저장 1회의 디스크 쓰기가 2배가 되고,
        # dashboard_snapshots 가 4,896건 / 디스크 25.2 GB 까지 누적되어 DB 최대 용량 소비원이 되었다.
        #
        # 업서트는 uuid 기준 병합이라 기존 데이터를 통째로 잃지 않으므로 매회 백업이 필요하지 않다.
        # 되돌리기 어려운 전체 덮어쓰기(POST /data/upload)에는 자동 스냅샷을 그대로 유지하며,
        # 수동 스냅샷(POST /snapshots)도 영향이 없다.

        # [타이밍] 자동 스냅샷 제거 후에는 이 구간이 0에 수렴한다 (계측 지점은 비교를 위해 유지)
        db.session.flush()
        _t_snapshot = _time.perf_counter()

        # 새 버전 계산
        new_version = server_version + 1

        # 데이터 저장/업데이트
        if server_data:
            server_data.version = new_version
            server_data.projects = merged_projects
            server_data.performances = merged_performances
            server_data.data_metadata = metadata
            server_data.last_modified_by = user_id
            server_data.last_modified_by_name = user_name
        else:
            server_data = DashboardData(
                version=new_version,
                projects=merged_projects,
                performances=merged_performances,
                data_metadata=metadata,
                last_modified_by=user_id,
                last_modified_by_name=user_name
            )
            db.session.add(server_data)

        # [타이밍] 싱글톤 행 UPDATE/INSERT 비용만 격리 측정
        db.session.flush()
        _t_singleton = _time.perf_counter()

        # 스냅샷 저장 (요청에 포함된 경우)
        if snapshot:
            new_snapshot = DashboardSnapshot(
                name=snapshot.get('name', f'업서트 스냅샷 v{new_version}'),
                description=snapshot.get('description', ''),
                snapshot_data={
                    'projects': merged_projects,
                    'performances': merged_performances,
                    'metadata': metadata
                },
                data_version=new_version,
                created_by=user_id,
                created_by_name=user_name,
                snapshot_type='upsert'
            )
            db.session.add(new_snapshot)

        # 활동 로그 저장
        for log in activity_logs:
            activity_log = DashboardActivityLog(
                action=log.get('action', 'UPSERT'),
                target_type=log.get('targetType', 'DATA'),
                target_id=log.get('targetId'),
                target_name=log.get('targetName'),
                changes=log.get('changes'),
                summary=log.get('summary'),
                user_id=user_id,
                user_name=user_name,
                data_version=new_version,
                source='server'
            )
            db.session.add(activity_log)

        # 업서트 자체에 대한 로그 (동기화 정보만 기록, 상세 변경 내역은 프론트엔드 개별 로그에 의존)
        upsert_log = DashboardActivityLog(
            action='UPSERT',
            target_type='DATA',
            target_name=f'v{server_version} → v{new_version} 동기화',
            summary=f'{user_name}이(가) 서버와 동기화함 (v{new_version}, 과제 {len(merged_projects)}개, 성과 {len(merged_performances)}개)',
            changes={
                'version': {'before': server_version, 'after': new_version},
                'projectCount': len(merged_projects),
                'performanceCount': len(merged_performances)
            },
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(upsert_log)

        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('upsert')
        _t_committed = _time.perf_counter()  # [타이밍] 커밋(활동로그 등 잔여 INSERT 포함) 완료
        try:
            _ms = lambda a, b: f"{(b - a) * 1000:.0f}ms"
            print(
                f"[Upsert timing] N_proj={len(merged_projects)} N_perf={len(merged_performances)} "
                f"payload={request.content_length or 0}B | "
                f"parse={_ms(_t0, _t_parsed)} load={_ms(_t_parsed, _t_loaded)} "
                f"merge={_ms(_t_loaded, _t_merged)} snapshot_write={_ms(_t_merged, _t_snapshot)} "
                f"singleton_write={_ms(_t_snapshot, _t_singleton)} commit_rest={_ms(_t_singleton, _t_committed)} "
                f"total={_ms(_t0, _t_committed)}"
            )
        except Exception:
            pass

        response_data = {
            'message': '데이터가 성공적으로 업데이트되었습니다.',
            'version': new_version,
            'updated_by': user_name,
            'stats': {
                'projects': {
                    'added': added_count,
                    'updated': updated_count,
                    'skipped': skipped_count,
                    'total': len(merged_projects)
                },
                'performances': {
                    'added': perf_added_count,
                    'updated': perf_updated_count,
                    'skipped': perf_skipped_count,
                    'renumbered': perf_renumbered_count,
                    'total': len(merged_performances)
                }
            }
        }

        # 재번호 부여된 성과가 있으면 정보 추가
        if renumbered_performances:
            response_data['renumbered_performances'] = renumbered_performances

        return success_response(response_data)

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Upsert failed: {str(e)}")
        return error_response(f'데이터 업데이트 실패: {str(e)}', status_code=500)


@bp.route('/data/download', methods=['GET'])
@jwt_required()
def download_dashboard_data():
    """서버 데이터 다운로드 (로그 기록)"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        data = DashboardData.query.first()

        if not data:
            return error_response('서버에 저장된 데이터가 없습니다.', status_code=404)

        # 다운로드 로그 기록 (동기화 정보만 기록)
        projects = data.projects or []
        performances = data.performances or []

        # 완전 삭제된 프로젝트 제외한 개수
        active_projects = [p for p in projects if not p.get('_permanentlyDeleted')]
        project_count = len(active_projects)
        active_performances = [p for p in performances if not p.get('_deleted')]
        performance_count = len(active_performances)

        download_log = DashboardActivityLog(
            action='DOWNLOAD',
            target_type='DATA',
            target_name=f'v{data.version} 다운로드',
            summary=f'{user_name}이(가) 서버에서 데이터를 불러옴 (v{data.version}, 과제 {project_count}개, 성과 {performance_count}개)',
            changes={
                'version': data.version,
                'projectCount': project_count,
                'performanceCount': performance_count
            },
            user_id=user_id,
            user_name=user_name,
            data_version=data.version,
            source='server'
        )
        db.session.add(download_log)
        db.session.commit()

        # 완전 삭제된 프로젝트 및 삭제된 성과는 다운로드에서 제외
        result = data.to_dict()
        if result.get('projects'):
            result['projects'] = [
                p for p in result['projects']
                if not p.get('_permanentlyDeleted')
            ]
        if result.get('performances'):
            result['performances'] = [
                p for p in result['performances']
                if not p.get('_deleted')
            ]

        return success_response(result)

    except Exception as e:
        print(f"[Data Error] Download failed: {str(e)}")
        return error_response(f'데이터 다운로드 실패: {str(e)}', status_code=500)


# ============================================
# 과제 소프트 삭제/복구 API
# ============================================

@bp.route('/project/<project_uuid>/delete', methods=['POST'])
@jwt_required()
def soft_delete_project(project_uuid):
    """
    과제 소프트 삭제 (실제 삭제가 아닌 _deleted 플래그 설정)
    """
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        server_data = DashboardData.query.first()
        if not server_data:
            return error_response('서버에 저장된 데이터가 없습니다.', status_code=404)

        projects = server_data.projects or []
        project_found = False
        project_name = ''

        for project in projects:
            if project.get('uuid') == project_uuid or project.get('id') == project_uuid:
                project['_deleted'] = True
                project['_deletedAt'] = now_utc_iso_z()
                project['_deletedBy'] = user_id
                project['_deletedByName'] = user_name
                project['updatedAt'] = now_utc_iso_z()  # 다른 사용자의 오래된 데이터로 덮어쓰기 방지
                project_found = True
                project_name = project.get('과제명', project_uuid)
                break

        if not project_found:
            return error_response('해당 과제를 찾을 수 없습니다.', status_code=404)

        # 버전 업데이트
        new_version = server_data.version + 1
        server_data.version = new_version
        server_data.projects = projects
        server_data.last_modified_by = user_id
        server_data.last_modified_by_name = user_name

        # JSON 필드 변경 감지를 위해 flag_modified 호출
        flag_modified(server_data, 'projects')

        # 활동 로그
        delete_log = DashboardActivityLog(
            action='DELETE',
            target_type='PROJECT',
            target_name=project_name,
            summary=f'{user_name}이(가) 과제 "{project_name}"을(를) 삭제함',
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(delete_log)
        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('delete')

        return success_response({
            'message': f'과제 "{project_name}"이(가) 삭제되었습니다.',
            'version': new_version,
            'deleted_by': user_name
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Soft delete project failed: {str(e)}")
        return error_response(f'과제 삭제 실패: {str(e)}', status_code=500)


@bp.route('/project/<project_uuid>/restore', methods=['POST'])
@jwt_required()
def restore_project(project_uuid):
    """
    삭제된 과제 복구
    """
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        server_data = DashboardData.query.first()
        if not server_data:
            return error_response('서버에 저장된 데이터가 없습니다.', status_code=404)

        projects = server_data.projects or []
        project_found = False
        project_name = ''

        for project in projects:
            if project.get('uuid') == project_uuid or project.get('id') == project_uuid:
                if not project.get('_deleted'):
                    return error_response('이 과제는 삭제된 상태가 아닙니다.', status_code=400)

                project['_deleted'] = False
                project.pop('_deletedAt', None)
                project.pop('_deletedBy', None)
                project.pop('_deletedByName', None)
                project['updatedAt'] = now_utc_iso_z()  # 다른 사용자의 오래된 데이터로 덮어쓰기 방지
                project_found = True
                project_name = project.get('과제명', project_uuid)
                break

        if not project_found:
            return error_response('해당 과제를 찾을 수 없습니다.', status_code=404)

        # 버전 업데이트
        new_version = server_data.version + 1
        server_data.version = new_version
        server_data.projects = projects
        server_data.last_modified_by = user_id
        server_data.last_modified_by_name = user_name

        # JSON 필드 변경 감지를 위해 flag_modified 호출
        flag_modified(server_data, 'projects')

        # 활동 로그
        restore_log = DashboardActivityLog(
            action='RESTORE',
            target_type='PROJECT',
            target_name=project_name,
            summary=f'{user_name}이(가) 과제 "{project_name}"을(를) 복구함',
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(restore_log)
        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('project-restore')

        return success_response({
            'message': f'과제 "{project_name}"이(가) 복구되었습니다.',
            'version': new_version,
            'restored_by': user_name
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Restore project failed: {str(e)}")
        return error_response(f'과제 복구 실패: {str(e)}', status_code=500)


@bp.route('/project/<project_uuid>/permanent-delete', methods=['DELETE'])
@jwt_required()
def permanent_delete_project(project_uuid):
    """
    과제 완전 삭제 (Soft Delete 방식 - 레코드 유지, 플래그만 설정)
    - 다른 사용자가 삭제 전 데이터를 업로드해도 복원되지 않도록 함
    - _permanentlyDeleted: true, _permanentlyDeletedAt 시간 기록
    """
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        server_data = DashboardData.query.first()
        if not server_data:
            return error_response('서버에 저장된 데이터가 없습니다.', status_code=404)

        projects = server_data.projects or []
        project_found = False
        project_name = ''

        # 해당 과제 찾기 및 완전 삭제 플래그 설정 (레코드는 유지)
        for project in projects:
            if project.get('uuid') == project_uuid or project.get('id') == project_uuid:
                project['_deleted'] = True
                project['_permanentlyDeleted'] = True
                project['_permanentlyDeletedAt'] = now_utc_iso_z()
                project['_permanentlyDeletedBy'] = user_id
                project['_permanentlyDeletedByName'] = user_name
                project_found = True
                project_name = project.get('과제명', project_uuid)
                break

        if not project_found:
            return error_response('해당 과제를 찾을 수 없습니다.', status_code=404)

        # 버전 업데이트
        new_version = server_data.version + 1
        server_data.version = new_version
        server_data.projects = projects
        server_data.last_modified_by = user_id
        server_data.last_modified_by_name = user_name

        # JSON 필드 변경 감지를 위해 flag_modified 호출
        flag_modified(server_data, 'projects')

        # 활동 로그
        perm_delete_log = DashboardActivityLog(
            action='PERMANENT_DELETE',
            target_type='PROJECT',
            target_name=project_name,
            summary=f'{user_name}이(가) 과제 "{project_name}"을(를) 완전 삭제함',
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(perm_delete_log)
        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('project-permanent-delete')

        return success_response({
            'message': f'과제 "{project_name}"이(가) 완전히 삭제되었습니다.',
            'version': new_version,
            'deleted_by': user_name
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Permanent delete project failed: {str(e)}")
        return error_response(f'과제 완전 삭제 실패: {str(e)}', status_code=500)


# ============================================
# 성과 삭제 API
# ============================================

@bp.route('/performance/<performance_id>/delete', methods=['DELETE'])
@jwt_required()
def delete_performance(performance_id):
    """
    성과 항목 삭제
    """
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        server_data = DashboardData.query.first()
        if not server_data:
            return error_response('서버에 저장된 데이터가 없습니다.', status_code=404)

        performances = server_data.performances or []
        projects = server_data.projects or []
        performance_name = ''
        performance_found = False

        # 해당 성과 찾기 및 soft-delete 플래그 설정 (레코드 유지)
        # 과제의 _permanentlyDeleted 방식과 동일 - 다른 사용자의 upsert로 부활 방지
        for perf in performances:
            if perf.get('uuid') == performance_id or perf.get('id') == performance_id:
                performance_name = perf.get('성과항목', performance_id)
                perf['_deleted'] = True
                perf['_deletedAt'] = now_utc_iso_z()
                perf['_deletedBy'] = user_id
                perf['_deletedByName'] = user_name
                performance_found = True
                break

        if not performance_found:
            return error_response('해당 성과를 찾을 수 없습니다.', status_code=404)

        # 프로젝트에서 해당 성과 참조 제거
        for project in projects:
            if project.get('성과목록'):
                original_len = len(project['성과목록'])
                project['성과목록'] = [
                    p for p in project['성과목록']
                    if not (
                        (isinstance(p, dict) and (p.get('id') == performance_id or p.get('성과항목ID') == performance_id)) or
                        (isinstance(p, str) and p == performance_id)
                    )
                ]
                # 성과 참조가 제거된 경우 updatedAt 갱신 (다른 사용자의 오래된 데이터로 덮어쓰기 방지)
                if len(project['성과목록']) != original_len:
                    project['updatedAt'] = now_utc_iso_z()

        # 버전 업데이트
        new_version = server_data.version + 1
        server_data.version = new_version
        server_data.performances = performances
        server_data.projects = projects
        server_data.last_modified_by = user_id
        server_data.last_modified_by_name = user_name

        # JSON 필드 변경 감지를 위해 flag_modified 호출
        flag_modified(server_data, 'performances')
        flag_modified(server_data, 'projects')

        # 활동 로그
        delete_log = DashboardActivityLog(
            action='DELETE',
            target_type='PERFORMANCE',
            target_name=performance_name,
            summary=f'{user_name}이(가) 성과 "{performance_name}"을(를) 삭제함',
            user_id=user_id,
            user_name=user_name,
            data_version=new_version,
            source='server'
        )
        db.session.add(delete_log)
        db.session.commit()

        # V2 동기화 요청 — 커밋 **뒤에** 건다. 백그라운드로 돌고 예외를 던지지 않으므로
        # 실패해도 이 저장은 이미 성공한 상태다. (실행계획 7.5-3)
        request_sync('delete')

        return success_response({
            'message': f'성과 "{performance_name}"이(가) 삭제되었습니다.',
            'version': new_version,
            'deleted_by': user_name
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Data Error] Delete performance failed: {str(e)}")
        return error_response(f'성과 삭제 실패: {str(e)}', status_code=500)


# ============================================
# 스냅샷 API
# ============================================

@bp.route('/snapshots', methods=['GET'])
@jwt_required()
def get_snapshots():
    """스냅샷 목록 조회"""
    try:
        snapshots = DashboardSnapshot.query.order_by(
            DashboardSnapshot.created_at.desc()
        ).limit(100).all()

        return success_response([s.to_dict() for s in snapshots])

    except Exception as e:
        print(f"[Snapshot Error] Get snapshots failed: {str(e)}")
        return error_response(f'스냅샷 조회 실패: {str(e)}', status_code=500)


@bp.route('/snapshots', methods=['POST'])
@jwt_required()
def create_snapshot():
    """스냅샷 생성"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        req_data = request.get_json()
        if not req_data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        # 현재 서버 데이터 버전 조회
        server_data = DashboardData.query.first()
        data_version = server_data.version if server_data else 0

        snapshot = DashboardSnapshot(
            name=req_data.get('name', '수동 스냅샷'),
            description=req_data.get('description', ''),
            snapshot_data=req_data.get('snapshotData', {}),
            data_version=data_version,
            created_by=user_id,
            created_by_name=user_name,
            snapshot_type=req_data.get('snapshotType', 'manual')
        )
        db.session.add(snapshot)
        db.session.commit()

        return success_response({
            'message': '스냅샷이 생성되었습니다.',
            'snapshot': snapshot.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Snapshot Error] Create snapshot failed: {str(e)}")
        return error_response(f'스냅샷 생성 실패: {str(e)}', status_code=500)


@bp.route('/snapshots/<int:snapshot_id>', methods=['GET'])
@jwt_required()
def get_snapshot(snapshot_id):
    """특정 스냅샷 조회"""
    try:
        snapshot = DashboardSnapshot.query.get(snapshot_id)
        if not snapshot:
            return error_response('스냅샷을 찾을 수 없습니다.', status_code=404)

        return success_response(snapshot.to_dict())

    except Exception as e:
        print(f"[Snapshot Error] Get snapshot failed: {str(e)}")
        return error_response(f'스냅샷 조회 실패: {str(e)}', status_code=500)


@bp.route('/snapshots/<int:snapshot_id>', methods=['DELETE'])
@jwt_required()
def delete_snapshot(snapshot_id):
    """스냅샷 삭제"""
    try:
        snapshot = DashboardSnapshot.query.get(snapshot_id)
        if not snapshot:
            return error_response('스냅샷을 찾을 수 없습니다.', status_code=404)

        db.session.delete(snapshot)
        db.session.commit()

        return success_response({'message': '스냅샷이 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[Snapshot Error] Delete snapshot failed: {str(e)}")
        return error_response(f'스냅샷 삭제 실패: {str(e)}', status_code=500)


# ============================================
# 활동 로그 API
# ============================================

@bp.route('/activity-logs', methods=['GET'])
@jwt_required()
def get_activity_logs():
    """활동 로그 조회"""
    try:
        # 쿼리 파라미터
        limit = request.args.get('limit', 100, type=int)
        offset = request.args.get('offset', 0, type=int)
        action_filter = request.args.get('action')
        target_type_filter = request.args.get('targetType')

        query = DashboardActivityLog.query

        if action_filter:
            query = query.filter_by(action=action_filter)
        if target_type_filter:
            query = query.filter_by(target_type=target_type_filter)

        logs = query.order_by(
            DashboardActivityLog.created_at.desc()
        ).offset(offset).limit(limit).all()

        total = query.count()

        return success_response({
            'logs': [log.to_dict() for log in logs],
            'total': total,
            'limit': limit,
            'offset': offset
        })

    except Exception as e:
        print(f"[Log Error] Get logs failed: {str(e)}")
        return error_response(f'로그 조회 실패: {str(e)}', status_code=500)


@bp.route('/activity-logs', methods=['POST'])
@jwt_required()
def create_activity_log():
    """활동 로그 생성 (클라이언트에서 직접 로그 전송)"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        req_data = request.get_json()
        if not req_data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        # 현재 서버 데이터 버전 조회
        server_data = DashboardData.query.first()
        data_version = server_data.version if server_data else 0

        log = DashboardActivityLog(
            action=req_data.get('action', 'UPDATE'),
            target_type=req_data.get('targetType', 'DATA'),
            target_id=req_data.get('targetId'),
            target_name=req_data.get('targetName'),
            changes=req_data.get('changes'),
            summary=req_data.get('summary'),
            user_id=user_id,
            user_name=user_name,
            data_version=data_version,
            source=req_data.get('source', 'local')
        )
        db.session.add(log)
        db.session.commit()

        return success_response({
            'message': '로그가 기록되었습니다.',
            'log': log.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Log Error] Create log failed: {str(e)}")
        return error_response(f'로그 생성 실패: {str(e)}', status_code=500)


# ============================================
# 파일 저장 공통 (첨부파일 · 보고서 이미지)
# ============================================

def ensure_upload_folder():
    """Ensure upload folder exists."""
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)


# ============================================
# 보고서 이미지 API (Phase 1-2)
# ============================================
# 기존에는 이미지가 과제 JSON 안에 base64 로 인라인되어, 과제 하나만 수정해도
# 전체 이미지가 통째로 왕복했다(운영: 33.9 MB = payload 의 94.4%).
# 여기서는 파일로 분리하고 과제 JSON 에는 imageId 참조만 남긴다.

REPORT_IMAGE_MAX_BYTES = 10 * 1024 * 1024          # 원본 1장 상한 10 MB
REPORT_IMAGE_ALLOWED_MIME = {
    'image/jpeg', 'image/jpg', 'image/png', 'image/gif', 'image/webp', 'image/bmp',
}

REPORT_IMAGE_SLOTS = ('이미지_좌측', '이미지_우측', '이미지_개요그림',
                      '이미지_상세내용그림', '이미지_향후계획그림')


def hydrate_report_images(projects):
    """
    보고서를 만들기 **전에** `imageId` 만 있는 이미지 원소에 `dataUrl` 을 채워 넣는다.

    왜 필요한가
        Phase 1-2 에서 이미지를 파일로 분리하면서 과제 JSON 에는 imageId 만 남겼는데,
        **보고서 생성기를 같이 고치지 않았다.** PPT 삽입부 6곳과 PDF 템플릿은 지금도
        `dataUrl` 만 읽는다. 그래서 분리 이후의 이미지는 화면에는 보이는데
        **보고서(PPT·PDF)에는 안 나온다** (2026-07-30 리허설에서 발견).

    왜 여기서 한 번에 하나
        삽입 지점이 12곳이다. 각각 고치면 하나를 빠뜨리기 쉽고, 나중에 슬롯이 늘 때마다
        또 12곳을 봐야 한다. **들어올 때 한 번** 채우면 PPT 와 PDF 가 같이 고쳐진다.

    비용
        파일은 서버에 있으므로 **네트워크 전송량은 늘지 않는다.** base64 로 만드는
        메모리 비용만 있고, 그건 Phase 1-2 이전에 원래 치르던 비용이다.

    동작
        · 이미 `dataUrl` 이 있으면 건드리지 않는다 (예전 데이터 호환)
        · 같은 imageId 가 여러 번 나와도 파일은 한 번만 읽는다
        · 파일이나 레코드가 없으면 그 장만 건너뛴다 — 보고서 생성 자체를 막지 않는다

    반환: (채운 장수, 건너뛴 장수)
    """
    import base64 as _b64

    if isinstance(projects, dict):
        projects = [projects]
    if not projects:
        return 0, 0

    # 1) 채워야 할 원소 모으기
    pending = []          # [(원소 dict, imageId)]
    for project in projects:
        if not isinstance(project, dict):
            continue
        for slot in REPORT_IMAGE_SLOTS:
            for element in (project.get(slot) or []):
                if not isinstance(element, dict):
                    continue
                if element.get('dataUrl'):
                    continue
                image_id = element.get('imageId')
                if image_id:
                    pending.append((element, image_id))

    if not pending:
        return 0, 0

    # 2) 메타를 한 번에 조회 (원소 수가 아니라 이미지 수만큼)
    wanted = {image_id for _, image_id in pending}
    rows = ReportImage.query.filter(ReportImage.id.in_(list(wanted))).all()
    meta = {row.id: row for row in rows}

    # 3) 파일은 imageId 당 한 번만 읽는다
    cache = {}
    filled = skipped = 0
    for element, image_id in pending:
        if image_id not in cache:
            row = meta.get(image_id)
            data_url = None
            if row is not None:
                path = os.path.join(UPLOAD_FOLDER, row.stored_filename)
                try:
                    with open(path, 'rb') as fh:
                        raw = fh.read()
                    mime = row.mime_type or 'image/png'
                    data_url = f"data:{mime};base64,{_b64.b64encode(raw).decode('ascii')}"
                except OSError as exc:
                    print(f'[ReportImage] 파일을 읽지 못했습니다 (id={image_id}): {exc}')
            else:
                print(f'[ReportImage] 참조가 가리키는 레코드가 없습니다 (id={image_id})')
            cache[image_id] = data_url

        if cache[image_id]:
            element['dataUrl'] = cache[image_id]
            filled += 1
        else:
            skipped += 1

    print(f'[ReportImage] 보고서용 이미지 {filled}장 채움'
          + (f' / {skipped}장 건너뜀' if skipped else ''))
    return filled, skipped


@bp.route('/report-images/<int:image_id>', methods=['GET'])
@jwt_required()
def get_report_image(image_id):
    """
    보고서 이미지 원본 서빙.

    브라우저 캐시를 허용한다. 이미지는 한 번 저장되면 내용이 바뀌지 않고
    (수정 시 새 레코드가 생긴다), 이 캐시가 보고서 화면 재방문 비용을 없앤다.
    """
    try:
        image = ReportImage.query.get(image_id)
        if not image:
            return error_response('이미지를 찾을 수 없습니다.', status_code=404)

        file_path = os.path.join(UPLOAD_FOLDER, image.stored_filename)
        if not os.path.exists(file_path):
            return error_response('이미지 파일이 존재하지 않습니다.', status_code=404)

        response = send_file(
            file_path,
            mimetype=image.mime_type or 'application/octet-stream',
            as_attachment=False,
            download_name=image.original_filename or image.stored_filename,
            conditional=True,
        )
        # 내용 불변이므로 길게 캐시해도 안전하다
        response.headers['Cache-Control'] = 'private, max-age=604800'
        return response

    except Exception as e:
        print(f"[ReportImage Error] Get image failed: {str(e)}")
        return error_response(f'이미지 조회 실패: {str(e)}', status_code=500)


@bp.route('/report-images/project/<project_id>', methods=['GET'])
@jwt_required()
def list_project_report_images(project_id):
    """특정 과제의 보고서 이미지 메타 목록 (이미지 데이터는 포함하지 않는다)"""
    try:
        images = (
            ReportImage.query
            .filter_by(project_id=project_id)
            .order_by(ReportImage.slot, ReportImage.position, ReportImage.id)
            .all()
        )
        return success_response([img.to_dict() for img in images])

    except Exception as e:
        print(f"[ReportImage Error] List failed: {str(e)}")
        return error_response(f'이미지 목록 조회 실패: {str(e)}', status_code=500)


@bp.route('/report-images/project/<project_id>', methods=['POST'])
@jwt_required()
def upload_report_image(project_id):
    """
    보고서 이미지 업로드 (multipart/form-data)

    form 필드: file(필수), slot(필수, 예 '이미지_좌측'), position, caption
    응답의 id 를 과제 JSON 의 imageId 로 저장한다.
    """
    try:
        import hashlib

        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        if 'file' not in request.files:
            return error_response('파일이 없습니다.', status_code=400)

        file = request.files['file']
        if not file.filename:
            return error_response('파일 이름이 없습니다.', status_code=400)

        slot = (request.form.get('slot') or '').strip()
        if not slot:
            return error_response('slot 이 필요합니다. (예: 이미지_좌측)', status_code=400)

        data = file.read()
        if not data:
            return error_response('빈 파일입니다.', status_code=400)
        if len(data) > REPORT_IMAGE_MAX_BYTES:
            return error_response(
                f'이미지가 너무 큽니다. 최대 {REPORT_IMAGE_MAX_BYTES // (1024 * 1024)} MB 까지 가능합니다.',
                status_code=400
            )

        mime = (file.mimetype or '').lower()
        if mime not in REPORT_IMAGE_ALLOWED_MIME:
            return error_response(f'지원하지 않는 이미지 형식입니다: {mime}', status_code=400)

        ensure_upload_folder()
        stored = ReportImage.generate_stored_filename(mime, file.filename)
        with open(os.path.join(UPLOAD_FOLDER, stored), 'wb') as fh:
            fh.write(data)

        image = ReportImage(
            project_id=project_id,
            slot=slot,
            position=int(request.form.get('position') or 0),
            original_filename=file.filename,
            stored_filename=stored,
            caption=request.form.get('caption') or '',
            mime_type=mime,
            file_size=len(data),
            sha256=hashlib.sha256(data).hexdigest(),
            uploaded_by=user_id,
            uploaded_by_name=user_name,
            source='upload',
        )
        db.session.add(image)
        db.session.commit()

        return success_response({'message': '이미지가 업로드되었습니다.', 'image': image.to_dict()})

    except Exception as e:
        db.session.rollback()
        print(f"[ReportImage Error] Upload failed: {str(e)}")
        return error_response(f'이미지 업로드 실패: {str(e)}', status_code=500)


@bp.route('/report-images/<int:image_id>', methods=['DELETE'])
@jwt_required()
def delete_report_image(image_id):
    """
    보고서 이미지 삭제 (DB 레코드만 제거하고 파일은 남긴다)

    파일을 즉시 지우지 않는 이유: 같은 이미지를 다른 과제가 참조하고 있을 수 있고,
    실수로 지웠을 때 되돌릴 여지를 남기기 위함이다. 파일 정리는 별도 배치로 한다.
    """
    try:
        image = ReportImage.query.get(image_id)
        if not image:
            return error_response('이미지를 찾을 수 없습니다.', status_code=404)

        db.session.delete(image)
        db.session.commit()
        return success_response({'message': '이미지가 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[ReportImage Error] Delete failed: {str(e)}")
        return error_response(f'이미지 삭제 실패: {str(e)}', status_code=500)


# ============================================
# 첨부파일 API
# ============================================

@bp.route('/attachments/project/<project_id>', methods=['GET'])
@jwt_required()
def get_project_attachments(project_id):
    """특정 과제의 첨부파일 목록 조회"""
    try:
        attachments = ProjectAttachment.query.filter_by(project_id=project_id).all()
        return success_response([att.to_dict() for att in attachments])

    except Exception as e:
        print(f"[Attachment Error] Get attachments failed: {str(e)}")
        return error_response(f'첨부파일 조회 실패: {str(e)}', status_code=500)


@bp.route('/attachments/project/<project_id>', methods=['POST'])
@jwt_required()
def upload_project_attachment(project_id):
    """과제에 첨부파일 업로드"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        user_name = user.name if user else 'Unknown'

        if 'file' not in request.files:
            return error_response('파일이 없습니다.', status_code=400)

        file = request.files['file']
        if file.filename == '':
            return error_response('파일 이름이 없습니다.', status_code=400)

        ensure_upload_folder()

        original_filename = file.filename
        stored_filename = ProjectAttachment.generate_stored_filename(original_filename)
        file_path = os.path.join(UPLOAD_FOLDER, stored_filename)

        # Save file to disk
        file.save(file_path)

        # Get file size
        file_size = os.path.getsize(file_path)

        # Create attachment record
        attachment = ProjectAttachment(
            project_id=project_id,
            original_filename=original_filename,
            stored_filename=stored_filename,
            file_size=file_size,
            mime_type=file.content_type,
            uploaded_by=user_id,
            uploaded_by_name=user_name
        )
        db.session.add(attachment)
        db.session.commit()

        return success_response({
            'message': '파일이 업로드되었습니다.',
            'attachment': attachment.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[Attachment Error] Upload failed: {str(e)}")
        return error_response(f'파일 업로드 실패: {str(e)}', status_code=500)


@bp.route('/attachments/<int:attachment_id>', methods=['GET'])
@jwt_required()
def download_attachment(attachment_id):
    """첨부파일 다운로드 (ZIP 압축으로 DRM 차단 우회)"""
    try:
        attachment = ProjectAttachment.query.get(attachment_id)
        if not attachment:
            return error_response('첨부파일을 찾을 수 없습니다.', status_code=404)

        file_path = os.path.join(UPLOAD_FOLDER, attachment.stored_filename)
        if not os.path.exists(file_path):
            return error_response('파일이 존재하지 않습니다.', status_code=404)

        # Create ZIP in memory
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            zip_file.write(file_path, attachment.original_filename)
        zip_buffer.seek(0)

        # Generate ZIP filename
        base_name = os.path.splitext(attachment.original_filename)[0]
        zip_filename = f'{base_name}.zip'

        return send_file(
            zip_buffer,
            as_attachment=True,
            download_name=zip_filename,
            mimetype='application/zip'
        )

    except Exception as e:
        print(f"[Attachment Error] Download failed: {str(e)}")
        return error_response(f'파일 다운로드 실패: {str(e)}', status_code=500)


@bp.route('/attachments/<int:attachment_id>/raw', methods=['GET'])
@jwt_required()
def download_attachment_raw(attachment_id):
    """첨부파일 원본 다운로드 (ZIP 변환 없이 - 다중 파일 ZIP 생성용)"""
    try:
        attachment = ProjectAttachment.query.get(attachment_id)
        if not attachment:
            return error_response('첨부파일을 찾을 수 없습니다.', status_code=404)

        file_path = os.path.join(UPLOAD_FOLDER, attachment.stored_filename)
        if not os.path.exists(file_path):
            return error_response('파일이 존재하지 않습니다.', status_code=404)

        return send_file(
            file_path,
            as_attachment=True,
            download_name=attachment.original_filename
        )

    except Exception as e:
        print(f"[Attachment Error] Raw download failed: {str(e)}")
        return error_response(f'파일 다운로드 실패: {str(e)}', status_code=500)


@bp.route('/attachments/<int:attachment_id>', methods=['DELETE'])
@jwt_required()
def delete_attachment(attachment_id):
    """첨부파일 삭제"""
    try:
        attachment = ProjectAttachment.query.get(attachment_id)
        if not attachment:
            return error_response('첨부파일을 찾을 수 없습니다.', status_code=404)

        # Delete file from disk
        file_path = os.path.join(UPLOAD_FOLDER, attachment.stored_filename)
        if os.path.exists(file_path):
            os.remove(file_path)

        # Delete record
        db.session.delete(attachment)
        db.session.commit()

        return success_response({'message': '첨부파일이 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[Attachment Error] Delete failed: {str(e)}")
        return error_response(f'첨부파일 삭제 실패: {str(e)}', status_code=500)


# ============================================
# KPI 대시보드 API
# ============================================

@bp.route('/kpi/categories', methods=['GET'])
@jwt_required()
def get_kpi_categories():
    """KPI 카테고리 목록 조회"""
    try:
        year = request.args.get('year', type=int)
        if not year:
            year = 2025  # 기본 연도

        categories = KPICategory.query.filter_by(
            year=year,
            is_active=True
        ).order_by(KPICategory.order).all()

        return success_response([cat.to_dict() for cat in categories])

    except Exception as e:
        print(f"[KPI Error] Get categories failed: {str(e)}")
        return error_response(f'카테고리 조회 실패: {str(e)}', status_code=500)


@bp.route('/kpi/categories', methods=['POST'])
@jwt_required()
def create_kpi_category():
    """KPI 카테고리 생성"""
    try:
        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        name = data.get('name')
        if not name:
            return error_response('카테고리 이름은 필수입니다.', status_code=400)

        year = data.get('year')
        if not year:
            return error_response('연도는 필수입니다.', status_code=400)

        # 현재 최대 order 값 조회
        max_order = db.session.query(db.func.max(KPICategory.order)).filter_by(
            year=year, is_active=True
        ).scalar() or 0

        category = KPICategory(
            name=name,
            color=data.get('color', '#3b82f6'),
            bg_color=data.get('bgColor', '#eff6ff'),
            hover_color=data.get('hoverColor', '#dbeafe'),
            order=max_order + 1,
            year=year
        )
        db.session.add(category)
        db.session.commit()

        return success_response({
            'message': '카테고리가 생성되었습니다.',
            'category': category.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Create category failed: {str(e)}")
        return error_response(f'카테고리 생성 실패: {str(e)}', status_code=500)


@bp.route('/kpi/categories/<int:category_id>', methods=['PUT'])
@jwt_required()
def update_kpi_category(category_id):
    """KPI 카테고리 수정"""
    try:
        category = KPICategory.query.get(category_id)
        if not category:
            return error_response('카테고리를 찾을 수 없습니다.', status_code=404)

        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        if 'name' in data:
            category.name = data['name']
        if 'color' in data:
            category.color = data['color']
        if 'bgColor' in data:
            category.bg_color = data['bgColor']
        if 'hoverColor' in data:
            category.hover_color = data['hoverColor']
        if 'order' in data:
            category.order = data['order']

        db.session.commit()

        return success_response({
            'message': '카테고리가 수정되었습니다.',
            'category': category.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Update category failed: {str(e)}")
        return error_response(f'카테고리 수정 실패: {str(e)}', status_code=500)


@bp.route('/kpi/categories/<int:category_id>', methods=['DELETE'])
@jwt_required()
def delete_kpi_category(category_id):
    """KPI 카테고리 삭제"""
    try:
        category = KPICategory.query.get(category_id)
        if not category:
            return error_response('카테고리를 찾을 수 없습니다.', status_code=404)

        # 해당 카테고리에 KPI가 있는지 확인
        kpi_count = KPI.query.filter_by(category_id=category_id, is_active=True).count()
        if kpi_count > 0:
            return error_response(
                f'이 카테고리에 {kpi_count}개의 KPI가 등록되어 있어 삭제할 수 없습니다.',
                status_code=400
            )

        # 소프트 삭제
        category.is_active = False
        db.session.commit()

        return success_response({'message': '카테고리가 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Delete category failed: {str(e)}")
        return error_response(f'카테고리 삭제 실패: {str(e)}', status_code=500)


@bp.route('/kpi', methods=['GET'])
@jwt_required()
def get_kpis():
    """KPI 목록 조회"""
    try:
        year = request.args.get('year', type=int)
        if not year:
            year = 2025  # 기본 연도

        kpis = KPI.query.filter_by(
            year=year,
            is_active=True
        ).order_by(KPI.category_id, KPI.order).all()

        return success_response([kpi.to_dict() for kpi in kpis])

    except Exception as e:
        print(f"[KPI Error] Get KPIs failed: {str(e)}")
        return error_response(f'KPI 조회 실패: {str(e)}', status_code=500)


@bp.route('/kpi', methods=['POST'])
@jwt_required()
def create_kpi():
    """KPI 생성"""
    try:
        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        # 필수 필드 검증
        required_fields = ['name', 'category', 'unit', 'targetValue', 'year']
        for field in required_fields:
            if field not in data or data[field] is None:
                return error_response(f'{field}는 필수입니다.', status_code=400)

        # 카테고리 존재 확인
        category_id = int(data['category'])
        category = KPICategory.query.get(category_id)
        if not category:
            return error_response('카테고리를 찾을 수 없습니다.', status_code=404)

        # 현재 최대 order 값 조회
        max_order = db.session.query(db.func.max(KPI.order)).filter_by(
            category_id=category_id, year=data['year'], is_active=True
        ).scalar() or 0

        kpi = KPI(
            name=data['name'],
            category_id=category_id,
            unit=data['unit'],
            kpi_type=data.get('type', 'monthly'),
            target_value=float(data['targetValue']),
            actual_value=float(data.get('actualValue', 0)),
            monthly_values=data.get('monthlyValues', [''] * 12),
            direction=data.get('direction', 'higher'),
            green_threshold=int(data.get('greenThreshold', 90)),
            yellow_threshold=int(data.get('yellowThreshold', 70)),
            year=int(data['year']),
            order=max_order + 1
        )
        db.session.add(kpi)
        db.session.commit()

        return success_response({
            'message': 'KPI가 생성되었습니다.',
            'kpi': kpi.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Create KPI failed: {str(e)}")
        return error_response(f'KPI 생성 실패: {str(e)}', status_code=500)


@bp.route('/kpi/<int:kpi_id>', methods=['PUT'])
@jwt_required()
def update_kpi(kpi_id):
    """KPI 수정"""
    try:
        kpi = KPI.query.get(kpi_id)
        if not kpi:
            return error_response('KPI를 찾을 수 없습니다.', status_code=404)

        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        if 'name' in data:
            kpi.name = data['name']
        if 'category' in data:
            kpi.category_id = int(data['category'])
        if 'unit' in data:
            kpi.unit = data['unit']
        if 'type' in data:
            kpi.kpi_type = data['type']
        if 'targetValue' in data:
            kpi.target_value = float(data['targetValue'])
        if 'actualValue' in data:
            kpi.actual_value = float(data['actualValue'])
        if 'monthlyValues' in data:
            kpi.monthly_values = data['monthlyValues']
        if 'direction' in data:
            kpi.direction = data['direction']
        if 'greenThreshold' in data:
            kpi.green_threshold = int(data['greenThreshold'])
        if 'yellowThreshold' in data:
            kpi.yellow_threshold = int(data['yellowThreshold'])
        if 'order' in data:
            kpi.order = int(data['order'])

        db.session.commit()

        return success_response({
            'message': 'KPI가 수정되었습니다.',
            'kpi': kpi.to_dict()
        })

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Update KPI failed: {str(e)}")
        return error_response(f'KPI 수정 실패: {str(e)}', status_code=500)


@bp.route('/kpi/<int:kpi_id>', methods=['DELETE'])
@jwt_required()
def delete_kpi(kpi_id):
    """KPI 삭제"""
    try:
        kpi = KPI.query.get(kpi_id)
        if not kpi:
            return error_response('KPI를 찾을 수 없습니다.', status_code=404)

        # 소프트 삭제
        kpi.is_active = False
        db.session.commit()

        return success_response({'message': 'KPI가 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Error] Delete KPI failed: {str(e)}")
        return error_response(f'KPI 삭제 실패: {str(e)}', status_code=500)


# ============================================
# KPI 대시보드 카드 API (새 KPI 대시보드)
# ============================================

@bp.route('/kpi/cards', methods=['GET'])
@jwt_required()
def get_kpi_dashboard_cards():
    """KPI 대시보드 카드 목록 조회"""
    try:
        year = request.args.get('year', type=int)
        if not year:
            year = 2025

        cards = KPIDashboardCard.query.filter_by(
            year=year,
            is_active=True
        ).order_by(KPIDashboardCard.order).all()

        return success_response([card.to_dict() for card in cards])

    except Exception as e:
        print(f"[KPI Card Error] Get cards failed: {str(e)}")
        return error_response(f'KPI 카드 조회 실패: {str(e)}', status_code=500)


@bp.route('/kpi/cards', methods=['POST'])
@jwt_required()
def create_kpi_dashboard_card():
    """KPI 대시보드 카드 생성"""
    try:
        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        name = data.get('name')
        if not name:
            return error_response('KPI 명칭은 필수입니다.', status_code=400)

        year = data.get('year')
        if not year:
            return error_response('연도는 필수입니다.', status_code=400)

        max_order = db.session.query(db.func.max(KPIDashboardCard.order)).filter_by(
            year=year, is_active=True
        ).scalar() or 0

        card = KPIDashboardCard(
            name=name,
            division=data.get('division', '전체'),
            category=data.get('category', '전체'),
            subcategories=data.get('subcategories', []),
            logic=data.get('logic', '합계'),
            selected_perf_keys=data.get('selectedPerfKeys', []),
            year=year,
            order=max_order + 1,
            treemap_enabled=bool(data.get('treemapEnabled', True))
        )
        db.session.add(card)
        db.session.commit()

        return success_response(card.to_dict(), 201)

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Card Error] Create card failed: {str(e)}")
        return error_response(f'KPI 카드 생성 실패: {str(e)}', status_code=500)


@bp.route('/kpi/cards/<int:card_id>', methods=['PUT'])
@jwt_required()
def update_kpi_dashboard_card(card_id):
    """KPI 대시보드 카드 수정"""
    try:
        card = KPIDashboardCard.query.get(card_id)
        if not card or not card.is_active:
            return error_response('KPI 카드를 찾을 수 없습니다.', status_code=404)

        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        if 'name' in data:
            card.name = data['name']
        if 'division' in data:
            card.division = data['division']
        if 'category' in data:
            card.category = data['category']
        if 'subcategories' in data:
            card.subcategories = data['subcategories']
        if 'logic' in data:
            card.logic = data['logic']
        if 'selectedPerfKeys' in data:
            card.selected_perf_keys = data['selectedPerfKeys']
        if 'order' in data:
            card.order = data['order']
        if 'treemapEnabled' in data:
            card.treemap_enabled = bool(data['treemapEnabled'])

        db.session.commit()

        return success_response(card.to_dict())

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Card Error] Update card failed: {str(e)}")
        return error_response(f'KPI 카드 수정 실패: {str(e)}', status_code=500)


@bp.route('/kpi/cards/reorder', methods=['PUT'])
@jwt_required()
def reorder_kpi_dashboard_cards():
    """KPI 대시보드 카드 순서 일괄 변경"""
    try:
        data = request.get_json()
        if not data or 'orderedIds' not in data:
            return error_response('orderedIds가 필요합니다.', status_code=400)

        ordered_ids = data['orderedIds']
        for idx, card_id in enumerate(ordered_ids):
            card = KPIDashboardCard.query.get(card_id)
            if card and card.is_active:
                card.order = idx

        db.session.commit()
        return success_response({'message': '순서가 변경되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Card Error] Reorder cards failed: {str(e)}")
        return error_response(f'KPI 카드 순서 변경 실패: {str(e)}', status_code=500)


@bp.route('/kpi/cards/<int:card_id>', methods=['DELETE'])
@jwt_required()
def delete_kpi_dashboard_card(card_id):
    """KPI 대시보드 카드 삭제"""
    try:
        card = KPIDashboardCard.query.get(card_id)
        if not card:
            return error_response('KPI 카드를 찾을 수 없습니다.', status_code=404)

        card.is_active = False
        db.session.commit()

        return success_response({'message': 'KPI 카드가 삭제되었습니다.'})

    except Exception as e:
        db.session.rollback()
        print(f"[KPI Card Error] Delete card failed: {str(e)}")
        return error_response(f'KPI 카드 삭제 실패: {str(e)}', status_code=500)


# ============================================
# 전체 첨부파일 다운로드 API (관리자 전용)
# ============================================

@bp.route('/attachments/all', methods=['GET'])
@jwt_required()
def get_all_attachments_with_projects():
    """모든 프로젝트의 첨부파일 목록 조회 (관리자 전용)"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)

        # 관리자 권한 확인
        if not user or (user.role != 'admin' and not user.is_admin):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        # 대시보드 데이터에서 프로젝트 정보 가져오기
        dashboard_data = DashboardData.query.first()
        projects_map = {}

        if dashboard_data and dashboard_data.projects:
            for project in dashboard_data.projects:
                project_id = project.get('id') or project.get('uuid')
                project_name = project.get('과제명', '알 수 없는 과제')
                if project_id:
                    projects_map[project_id] = project_name

        # 모든 첨부파일 조회
        attachments = ProjectAttachment.query.all()

        # 프로젝트별로 그룹화
        result = []
        for att in attachments:
            project_name = projects_map.get(att.project_id, f'Unknown ({att.project_id})')
            result.append({
                **att.to_dict(),
                'project_name': project_name
            })

        return success_response(result)

    except Exception as e:
        print(f"[Attachment Error] Get all attachments failed: {str(e)}")
        return error_response(f'첨부파일 조회 실패: {str(e)}', status_code=500)


@bp.route('/attachments/download-all', methods=['GET'])
@jwt_required()
def download_all_attachments_zip():
    """모든 프로젝트의 첨부파일을 ZIP으로 다운로드 (관리자 전용)"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)

        # 관리자 권한 확인
        if not user or (user.role != 'admin' and not user.is_admin):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        # 대시보드 데이터에서 프로젝트 정보 가져오기
        dashboard_data = DashboardData.query.first()
        projects_map = {}

        if dashboard_data and dashboard_data.projects:
            for project in dashboard_data.projects:
                project_name = project.get('과제명', '알 수 없는 과제')
                safe_name = "".join(c for c in project_name if c not in r'\/:*?"<>|')
                # id와 uuid 둘 다 매핑 (첨부파일이 어느 쪽으로 저장되었든 매칭)
                for key in ('id', 'uuid'):
                    val = project.get(key)
                    if val:
                        projects_map[val] = safe_name

        # 모든 첨부파일 조회
        attachments = ProjectAttachment.query.all()

        if not attachments:
            return error_response('다운로드할 첨부파일이 없습니다.', status_code=404)

        # ZIP 파일 생성 (과제명/파일명 구조, 중복 파일명 방지)
        zip_buffer = io.BytesIO()
        used_paths = set()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for att in attachments:
                file_path = os.path.join(UPLOAD_FOLDER, att.stored_filename)
                if os.path.exists(file_path):
                    project_name = projects_map.get(att.project_id, f'기타_{att.project_id[:8]}')
                    zip_path = f'{project_name}/{att.original_filename}'
                    # 같은 폴더에 동일 파일명이 있으면 번호 붙이기
                    if zip_path in used_paths:
                        name, ext = os.path.splitext(att.original_filename)
                        counter = 1
                        while zip_path in used_paths:
                            zip_path = f'{project_name}/{name}_{counter}{ext}'
                            counter += 1
                    used_paths.add(zip_path)
                    zip_file.write(file_path, zip_path)

        zip_buffer.seek(0)

        # 다운로드 로그 기록
        download_log = DashboardActivityLog(
            action='DOWNLOAD',
            target_type='ATTACHMENT',
            target_name=f'전체 첨부파일 ({len(attachments)}개)',
            summary=f'{user.name}이(가) 모든 과제 첨부파일을 다운로드함',
            user_id=user_id,
            user_name=user.name,
            source='server'
        )
        db.session.add(download_log)
        db.session.commit()

        return send_file(
            zip_buffer,
            as_attachment=True,
            download_name='all_project_attachments.zip',
            mimetype='application/zip'
        )

    except Exception as e:
        print(f"[Attachment Error] Download all attachments failed: {str(e)}")
        return error_response(f'첨부파일 다운로드 실패: {str(e)}', status_code=500)


# ============================================
# PPT 보고서 생성 API (관리자 전용)
# ============================================

TEMPLATE_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
    'templates', 'ppt'
)


def get_template_path(filename=None):
    """안전하게 템플릿 파일 경로를 결정한다.
    filename이 없으면 기본 템플릿을 사용한다."""
    if not filename:
        filename = 'project_report_template.pptx'
    safe_name = os.path.basename(filename)
    if not safe_name.endswith('.pptx'):
        raise ValueError('유효하지 않은 템플릿 파일 형식입니다.')
    path = os.path.join(TEMPLATE_DIR, safe_name)
    if not os.path.exists(path):
        raise FileNotFoundError(f'템플릿 파일을 찾을 수 없습니다: {safe_name}')
    return path


@bp.route('/report/templates', methods=['GET'])
@jwt_required()
def get_report_templates():
    """PPT 템플릿 목록 조회"""
    try:
        import glob as glob_mod
        pattern = os.path.join(TEMPLATE_DIR, '*.pptx')
        files = glob_mod.glob(pattern)
        templates = []
        for f in sorted(files):
            basename = os.path.basename(f)
            # 임시 파일(~$) 제외
            if basename.startswith('~$'):
                continue
            name = os.path.splitext(basename)[0]
            templates.append({'filename': basename, 'name': name})
        return success_response(templates)
    except Exception as e:
        return error_response(f'템플릿 목록 조회 실패: {str(e)}', status_code=500)


@bp.route('/report/ppt', methods=['POST'])
@jwt_required()
def generate_ppt_report():
    """프로젝트 PPT 보고서 생성 (관리자 전용)"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        user_id = get_jwt_identity()
        user = User.query.get(user_id)

        # viewer 제외 모든 역할 허용
        if not user or user.role == 'viewer':
            return error_response('보고서 저장 권한이 없습니다.', status_code=403)

        # 프로젝트 데이터 받기
        project_data = request.get_json()
        if not project_data:
            return error_response('프로젝트 데이터가 없습니다.', status_code=400)

        # imageId 만 온 이미지에 dataUrl 을 채운다 — 아래 삽입부는 dataUrl 만 읽는다
        hydrate_report_images(project_data)

        # 템플릿 파일 경로
        try:
            template_path = get_template_path(project_data.get('template'))
        except (ValueError, FileNotFoundError) as e:
            return error_response(str(e), status_code=400)

        # 템플릿 로드
        prs = Presentation(template_path)

        # HTML 태그 제거 함수
        import re
        def strip_html(html_text):
            if not html_text:
                return '-'
            # HTML 태그 제거
            text = re.sub(r'<[^>]+>', '', html_text)
            # HTML 엔티티 변환
            text = text.replace('&nbsp;', ' ')
            text = text.replace('&lt;', '<')
            text = text.replace('&gt;', '>')
            text = text.replace('&amp;', '&')
            text = text.replace('&quot;', '"')
            # 연속 공백 정리
            text = re.sub(r'\s+', ' ', text).strip()
            return text if text else '-'

        # 플레이스홀더 매핑
        placeholders = {
            '{{과제명}}': project_data.get('과제명', '-'),
            '{{과제년도}}': str(project_data.get('과제년도', '-')),
            '{{사업부}}': project_data.get('사업부', '-'),
            '{{프로세스}}': project_data.get('프로세스', '-'),
            '{{과제PL}}': project_data.get('과제PL', '-'),
            '{{작성자}}': project_data.get('작성자', '-'),
            '{{시작}}': str(project_data.get('시작', '-')),
            '{{종료}}': str(project_data.get('종료', '-')),
            '{{과제영역}}': project_data.get('과제영역', '-'),
            '{{과제구분}}': project_data.get('과제구분', '-'),
            '{{진행상태}}': project_data.get('진행상태', '-'),
            '{{진행률}}': str(project_data.get('진행률', 0)),
            '{{PoC과제여부}}': '예' if project_data.get('PoC과제여부') else '아니오',
            '{{중점과제여부}}': '예' if project_data.get('중점과제여부') else '아니오',
            '{{과제상세설명}}': strip_html(project_data.get('과제상세설명', '')),
        }

        # 담당부서 목록 처리
        dept_list = project_data.get('담당부서목록', [])
        placeholders['{{담당부서목록}}'] = ', '.join(dept_list) if dept_list else '-'

        # 참여인력 목록 처리
        personnel_list = project_data.get('과제참여인력목록', [])
        if personnel_list:
            personnel_text = '\n'.join([
                f"• {p.get('이름', '-')} ({p.get('부서', '-')})"
                for p in personnel_list
            ])
        else:
            personnel_text = '등록된 참여인력이 없습니다.'
        placeholders['{{과제참여인력목록}}'] = personnel_text
        placeholders['{{과제참여인력_이름목록}}'] = ', '.join([
            p.get('이름', '-') for p in personnel_list
        ]) if personnel_list else '-'
        # 소속별 그룹화: 부서 1개면 이름만, 2개 이상이면 "소속A 이름1, 이름2 / 소속B 이름3"
        if personnel_list:
            from collections import OrderedDict
            dept_groups = OrderedDict()
            for p in personnel_list:
                dept = p.get('부서', '-')
                name = p.get('이름', '-')
                dept_groups.setdefault(dept, []).append(name)
            if len(dept_groups) == 1:
                placeholders['{{과제참여인력_소속별목록}}'] = ', '.join(list(dept_groups.values())[0])
            else:
                placeholders['{{과제참여인력_소속별목록}}'] = ' / '.join(
                    f"{dept} {', '.join(names)}" for dept, names in dept_groups.items()
                )
        else:
            placeholders['{{과제참여인력_소속별목록}}'] = '-'
        placeholders['{{참여인력수}}'] = str(len(personnel_list))
        # 참여인력 개별 플레이스홀더 (최대 20명)
        for i in range(20):
            idx = i + 1
            if i < len(personnel_list):
                p = personnel_list[i]
                email = p.get('knoxId', '') or ''
                placeholders[f'{{{{참여인력{idx}_이름}}}}'] = p.get('이름', '-')
                placeholders[f'{{{{참여인력{idx}_부서}}}}'] = p.get('부서', '-')
                placeholders[f'{{{{참여인력{idx}_knoxId}}}}'] = email
                placeholders[f'{{{{참여인력{idx}_선택사업부}}}}'] = p.get('선택사업부', '-')
            else:
                placeholders[f'{{{{참여인력{idx}_이름}}}}'] = ''
                placeholders[f'{{{{참여인력{idx}_부서}}}}'] = ''
                placeholders[f'{{{{참여인력{idx}_knoxId}}}}'] = ''
                placeholders[f'{{{{참여인력{idx}_선택사업부}}}}'] = ''

        # 성과 목록 처리
        performance_list = project_data.get('성과목록', [])
        if performance_list:
            perf_lines = []
            for i, p in enumerate(performance_list, 1):
                line = f"{i}. {p.get('성과항목', '-')}"
                line += f"\n   - 대분류: {p.get('대분류', '-')}"
                line += f"\n   - 기여도: {p.get('과제기여도', '-')}%"
                line += f"\n   - 현재: {p.get('현재수준', '-')} / 목표: {p.get('목표수준', '-')} / 실적: {p.get('실적수준', '-')} ({p.get('단위', '-')})"
                perf_lines.append(line)
            perf_text = '\n\n'.join(perf_lines)
        else:
            perf_text = '등록된 성과가 없습니다.'
        placeholders['{{성과목록}}'] = perf_text
        placeholders['{{성과수}}'] = str(len(performance_list))
        # 성과 개별 플레이스홀더 (최대 20개)
        for i in range(20):
            idx = i + 1
            if i < len(performance_list):
                p = performance_list[i]
                placeholders[f'{{{{성과{idx}_항목}}}}'] = p.get('성과항목', '-')
                placeholders[f'{{{{성과{idx}_대분류}}}}'] = p.get('대분류', '-')
                placeholders[f'{{{{성과{idx}_기여도}}}}'] = str(p.get('과제기여도', '-'))
                placeholders[f'{{{{성과{idx}_현재수준}}}}'] = str(p.get('현재수준', '-'))
                placeholders[f'{{{{성과{idx}_목표수준}}}}'] = str(p.get('목표수준', '-'))
                placeholders[f'{{{{성과{idx}_실적수준}}}}'] = str(p.get('실적수준', '-'))
                placeholders[f'{{{{성과{idx}_단위}}}}'] = p.get('단위', '-')
            else:
                placeholders[f'{{{{성과{idx}_항목}}}}'] = ''
                placeholders[f'{{{{성과{idx}_대분류}}}}'] = ''
                placeholders[f'{{{{성과{idx}_기여도}}}}'] = ''
                placeholders[f'{{{{성과{idx}_현재수준}}}}'] = ''
                placeholders[f'{{{{성과{idx}_목표수준}}}}'] = ''
                placeholders[f'{{{{성과{idx}_실적수준}}}}'] = ''
                placeholders[f'{{{{성과{idx}_단위}}}}'] = ''

        # 액션아이템 목록 처리
        action_list = project_data.get('액션아이템목록', [])
        if action_list:
            action_lines = []
            for item in action_list:
                content = item.get('제목') or item.get('content') or item.get('내용', '-')
                due_date = item.get('목표일') or item.get('dueDate') or item.get('완료일', '-')
                done = '완료' if item.get('완료여부') else '진행중'
                status = item.get('status') or item.get('상태') or done
                action_lines.append(f"• {content} (완료일: {due_date}, 상태: {status})")
            action_text = '\n'.join(action_lines)
        else:
            action_text = '등록된 액션아이템이 없습니다.'
        placeholders['{{액션아이템목록}}'] = action_text
        placeholders['{{액션아이템수}}'] = str(len(action_list))
        # 액션아이템 개별 플레이스홀더 (최대 20개)
        for i in range(20):
            idx = i + 1
            if i < len(action_list):
                item = action_list[i]
                content = item.get('제목') or item.get('content') or item.get('내용', '-')
                due_date = item.get('목표일') or item.get('dueDate') or item.get('완료일', '-')
                done_flag = item.get('완료여부', False)
                completion_date = item.get('완료일', '-')
                sub_items = item.get('세부항목목록', [])
                sub_text = ', '.join([s.get('내용', '-') for s in sub_items]) if sub_items else '-'
                placeholders[f'{{{{액션아이템{idx}_이름}}}}'] = content
                placeholders[f'{{{{액션아이템{idx}_목표일}}}}'] = str(due_date)
                placeholders[f'{{{{액션아이템{idx}_완료여부}}}}'] = '완료' if done_flag else '진행중'
                placeholders[f'{{{{액션아이템{idx}_완료일}}}}'] = str(completion_date)
                placeholders[f'{{{{액션아이템{idx}_세부항목}}}}'] = sub_text
            else:
                placeholders[f'{{{{액션아이템{idx}_이름}}}}'] = ''
                placeholders[f'{{{{액션아이템{idx}_목표일}}}}'] = ''
                placeholders[f'{{{{액션아이템{idx}_완료여부}}}}'] = ''
                placeholders[f'{{{{액션아이템{idx}_완료일}}}}'] = ''
                placeholders[f'{{{{액션아이템{idx}_세부항목}}}}'] = ''

        # 상세 과제 정보 텍스트 플레이스홀더
        DETAIL_LABEL_MAP = {
            '과제개요': '과제 개요',
            '추진배경': '추진 배경',
            '과제목표': '과제 목표',
            '상세내용': '상세 내용',
            '성과': '기술/경영 성과',
            '산출물': '산출물',
            '향후계획': '향후 계획',
        }

        def format_detail_section(raw, label):
            """상세 과제 정보 섹션을 포맷팅된 텍스트로 변환"""
            if not raw or not raw.get('enabled'):
                return ''
            items = raw.get('items', [])
            lines = [f"□ {label}"]
            for item in items:
                if isinstance(item, str):
                    item = {'text': item, 'children': []}
                text = item.get('text', '')
                if not text:
                    continue
                lines.append(f"- {text}")
                for child in item.get('children', []):
                    child_text = child.get('text', '') if isinstance(child, dict) else str(child)
                    if child_text:
                        lines.append(f"  · {child_text}")
            return '\n'.join(lines) if len(lines) > 1 else ''

        for section_key, section_label in DETAIL_LABEL_MAP.items():
            raw = project_data.get(f'상세정보_{section_key}')
            placeholders[f'{{{{상세정보_{section_key}}}}}'] = format_detail_section(raw, section_label)

        def _reset_autofit_font_scale(prs_obj):
            """normAutofit(넘치면 텍스트 크기 조정)이 적용된 텍스트박스를
            noAutofit으로 전환하고, 텍스트가 박스에 맞도록 폰트 크기를 직접 계산하여 설정한다.
            spAutoFit(도형을 텍스트 크기에 맞춤)은 건드리지 않는다."""
            from pptx.oxml.ns import qn
            from pptx.util import Emu, Pt as PtUtil
            from lxml.etree import SubElement

            for slide in prs_obj.slides:
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    tf = shape.text_frame
                    bodyPr = tf._txBody.find(qn('a:bodyPr'))
                    if bodyPr is None:
                        continue
                    normAutofit = bodyPr.find(qn('a:normAutofit'))
                    if normAutofit is None:
                        continue

                    # 현재 폰트 크기 수집 (기본값)
                    base_font_size = None
                    for p in tf.paragraphs:
                        for run in p.runs:
                            if run.font.size:
                                base_font_size = run.font.size
                                break
                        if base_font_size:
                            break
                    if not base_font_size:
                        base_font_size = PtUtil(10)

                    # 박스 내부 여백 계산
                    def _emu_attr(attr, default=45720):
                        val = bodyPr.get(attr)
                        return int(val) if val else default
                    margin_top = _emu_attr('tIns')
                    margin_bottom = _emu_attr('bIns')
                    margin_left = _emu_attr('lIns', 91440)
                    margin_right = _emu_attr('rIns', 91440)

                    usable_h = shape.height - margin_top - margin_bottom
                    usable_w = shape.width - margin_left - margin_right
                    if usable_h <= 0 or usable_w <= 0:
                        continue

                    # 주어진 폰트 크기(EMU)에서 필요한 총 높이를 계산
                    import math
                    def _calc_needed_h(fs_emu):
                        lh = int(fs_emu * 1.2)
                        if lh <= 0:
                            return 0
                        lines = 0
                        for p in tf.paragraphs:
                            text = p.text
                            if not text:
                                lines += 1
                                continue
                            avg_cw = fs_emu * 0.85
                            cpl = max(1, int(usable_w / avg_cw))
                            lines += max(1, math.ceil(len(text) / cpl))
                        return lines * lh

                    needed_h = _calc_needed_h(base_font_size)
                    if needed_h <= usable_h:
                        # 텍스트가 박스에 들어감 → normAutofit 제거만
                        bodyPr.remove(normAutofit)
                        SubElement(bodyPr, qn('a:noAutofit'))
                        continue

                    # 축소할 폰트 크기를 이진 탐색으로 찾기
                    min_pt = 6
                    base_pt = base_font_size / 12700
                    best_pt = min_pt
                    lo, hi = min_pt, int(base_pt)
                    while lo <= hi:
                        mid = (lo + hi) // 2
                        if _calc_needed_h(mid * 12700) <= usable_h:
                            best_pt = mid
                            lo = mid + 1
                        else:
                            hi = mid - 1

                    # 모든 run에 폰트 크기 적용 (비율 유지)
                    scale = best_pt / base_pt if base_pt > 0 else 1
                    for p in tf.paragraphs:
                        for run in p.runs:
                            current_size = run.font.size or base_font_size
                            current_pt = current_size / 12700
                            new_pt = max(min_pt, round(current_pt * scale))
                            run.font.size = PtUtil(new_pt)

                    # normAutofit → noAutofit 전환
                    bodyPr.remove(normAutofit)
                    SubElement(bodyPr, qn('a:noAutofit'))

        # paragraph 단위 placeholder 치환 (PowerPoint가 run을 쪼개도 동작)
        def replace_paragraph_placeholders(paragraph, placeholders):
            """paragraph의 전체 텍스트에서 placeholder를 찾아 치환한다.
            PowerPoint는 같은 문자열을 여러 run으로 쪼갤 수 있으므로,
            run 단위가 아닌 paragraph 전체 텍스트로 매칭한 뒤
            run 텍스트를 재조립한다."""
            runs = paragraph.runs
            if not runs:
                return
            full_text = ''.join(run.text for run in runs)

            # 이 paragraph에 치환할 placeholder가 있는지 확인
            has_match = False
            for ph in placeholders:
                if ph in full_text:
                    has_match = True
                    break
            if not has_match:
                return

            # 치환 수행
            for ph, val in placeholders.items():
                full_text = full_text.replace(ph, str(val))

            # 치환된 텍스트를 첫 번째 run에 넣고 나머지 run은 비운다
            runs[0].text = full_text
            for run in runs[1:]:
                run.text = ''

        def _add_performance_table(slide, project_data, left, top, width):
            """성과 섹션 아래에 성과목록 테이블을 추가"""
            from pptx.util import Pt as PtUtil, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.oxml.ns import qn

            perf_list = project_data.get('성과목록', [])
            if not perf_list:
                return

            border_color = RGBColor(209, 213, 219)  # #d1d5db
            header_bg = RGBColor(209, 213, 219)      # #d1d5db
            cell_bg = RGBColor(255, 255, 255)         # white
            text_color = RGBColor(0, 0, 0)            # black

            rows = len(perf_list) + 1
            cols = 5  # 소분류, 성과항목명, 기존, 목표, 실적
            row_h = Emu(230000)
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
            tbl = table_shape.table

            # 테이블 기본 스타일 완전 제거 (tableStyleId, 밴딩 등)
            tbl_pr = tbl._tbl.tblPr
            tbl_pr.attrib.pop('bandRow', None)
            tbl_pr.attrib.pop('bandCol', None)
            tbl_pr.attrib.pop('firstRow', None)
            tbl_pr.attrib.pop('lastRow', None)
            for child in list(tbl_pr):
                if child.tag.endswith('tableStyleId') or child.tag.endswith('tblStyle'):
                    tbl_pr.remove(child)

            def set_cell_style(cell, fill_rgb, border_color):
                """셀의 배경색과 4면 테두리를 한번에 설정 (올바른 XML 순서 보장)"""
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                # 기존 테두리/배경 모두 제거
                for tag in ('lnL', 'lnR', 'lnT', 'lnB', 'solidFill'):
                    el = tcPr.find(qn(f'a:{tag}'))
                    if el is not None:
                        tcPr.remove(el)
                # 테두리 (ln*) 먼저
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = tcPr.makeelement(qn(f'a:{edge}'), {})
                    ln.set('w', '6350')  # 0.5pt
                    sf = ln.makeelement(qn('a:solidFill'), {})
                    srgb = sf.makeelement(qn('a:srgbClr'), {'val': border_color})
                    sf.append(srgb)
                    ln.append(sf)
                    tcPr.append(ln)
                # 배경 (solidFill) 나중에
                bg = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = bg.makeelement(qn('a:srgbClr'), {
                    'val': f'{fill_rgb.red:02X}{fill_rgb.green:02X}{fill_rgb.blue:02X}' if hasattr(fill_rgb, 'red') else str(fill_rgb)
                })
                bg.append(srgb)
                tcPr.append(bg)

            def fmt_delta(base_val, compare_val):
                """기존 대비 변화값을 줄바꿈 후 (+X) / (-X) 형식으로 반환"""
                try:
                    b = float(base_val)
                    c = float(compare_val)
                    d = c - b
                    if d >= 0:
                        return f'\n(+{d:g})'
                    else:
                        return f'\n({d:g})'
                except (ValueError, TypeError):
                    return ''

            col_widths = [0.22, 0.34, 0.15, 0.15, 0.14]
            for ci, ratio in enumerate(col_widths):
                tbl.columns[ci].width = int(width * ratio)

            headers = ['성과분류', '성과항목명', '기존', '목표', '실적']
            for ci, h in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = PtUtil(8)
                    p.font.bold = True
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.CENTER
                set_cell_style(cell, header_bg, 'D1D5DB')

            for ri, perf in enumerate(perf_list):
                unit = perf.get('단위', '')
                unit_suffix = f' {unit}' if unit else ''
                base_raw = perf.get('현재수준', '-')
                target_raw = perf.get('목표수준', '-')
                actual_raw = perf.get('실적수준', '-') or '-'

                import re as _re
                perf_name = _re.sub(r'^\[.+?\]\s*', '', perf.get('성과항목', '-'))
                row_data = [
                    perf.get('소분류', '-'),
                    perf_name or '-',
                    f'{base_raw}{unit_suffix}',
                    f'{target_raw}{unit_suffix}{fmt_delta(base_raw, target_raw)}',
                    f'{actual_raw}{unit_suffix}{fmt_delta(base_raw, actual_raw)}',
                ]
                for ci, val in enumerate(row_data):
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        p.font.size = PtUtil(8)
                        p.font.color.rgb = text_color
                        p.alignment = PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.CENTER
                    set_cell_style(cell, cell_bg, 'D1D5DB')

        # 상세정보 플레이스홀더를 멀티 paragraph로 특별 처리 (제목 줄 스타일 분리)
        def replace_detail_placeholders(prs_obj, project_data):
            from pptx.util import Pt as PtUtil
            detail_keys = {
                f'{{{{상세정보_{k}}}}}': (k, v) for k, v in DETAIL_LABEL_MAP.items()
            }
            for slide in prs_obj.slides:
                perf_tables_to_add = []  # (slide, left, top, width)
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    tf = shape.text_frame
                    full_text = ''.join(run.text for p in tf.paragraphs for run in p.runs)
                    matched = [(pk, sk, sl) for pk, (sk, sl) in detail_keys.items() if pk in full_text]
                    if not matched:
                        continue

                    base_font_name = None
                    base_font_size = None
                    base_font_color = None
                    for p in tf.paragraphs:
                        for run in p.runs:
                            if run.font.name:
                                base_font_name = run.font.name
                            if run.font.size:
                                base_font_size = run.font.size
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    base_font_color = run.font.color.rgb
                            except AttributeError:
                                pass
                            break
                        if not base_font_size and p.font and p.font.size:
                            base_font_size = p.font.size
                        if not base_font_name and p.font and p.font.name:
                            base_font_name = p.font.name
                        break

                    body_size = base_font_size if base_font_size else PtUtil(12)
                    title_size = body_size + PtUtil(2)

                    # 줄간격 추출
                    base_space_before = None
                    base_space_after = None
                    base_line_spacing = None
                    for p in tf.paragraphs:
                        if p.space_before is not None:
                            base_space_before = p.space_before
                        if p.space_after is not None:
                            base_space_after = p.space_after
                        if p.line_spacing is not None:
                            base_line_spacing = p.line_spacing
                        break

                    def apply_spacing(para):
                        if base_space_before is not None:
                            para.space_before = base_space_before
                        if base_space_after is not None:
                            para.space_after = base_space_after
                        if base_line_spacing is not None:
                            para.line_spacing = base_line_spacing

                    for p in list(tf.paragraphs)[1:]:
                        p._element.getparent().remove(p._element)
                    tf.paragraphs[0].clear()
                    tf.paragraphs[0].font.size = None
                    first_para_used = False

                    for ph_key, section_key, section_label in matched:
                        raw = project_data.get(f'상세정보_{section_key}')
                        if not raw or not raw.get('enabled'):
                            placeholders[ph_key] = ''
                            continue

                        items = raw.get('items', [])

                        if not first_para_used:
                            p = tf.paragraphs[0]
                            first_para_used = True
                        else:
                            blank = tf.add_paragraph()
                            blank.font.size = None
                            apply_spacing(blank)
                            br = blank.add_run()
                            br.text = ''
                            br.font.size = PtUtil(6)
                            p = tf.add_paragraph()

                        p.font.size = None
                        apply_spacing(p)
                        run = p.add_run()
                        run.text = f'□ {section_label}'
                        run.font.size = title_size
                        run.font.bold = True
                        if base_font_name:
                            run.font.name = base_font_name
                        if base_font_color:
                            run.font.color.rgb = base_font_color

                        for item in items:
                            if isinstance(item, str):
                                item = {'text': item, 'children': []}
                            text = item.get('text', '')
                            if not text:
                                continue
                            p = tf.add_paragraph()
                            p.font.size = None
                            apply_spacing(p)
                            run = p.add_run()
                            run.text = f'- {text}'
                            run.font.size = body_size
                            run.font.bold = False
                            if base_font_name:
                                run.font.name = base_font_name
                            if base_font_color:
                                run.font.color.rgb = base_font_color

                            for child in item.get('children', []):
                                child_text = child.get('text', '') if isinstance(child, dict) else str(child)
                                if not child_text:
                                    continue
                                p = tf.add_paragraph()
                                p.font.size = None
                                apply_spacing(p)
                                run = p.add_run()
                                run.text = f'  · {child_text}'
                                run.font.size = body_size
                                run.font.bold = False
                                if base_font_name:
                                    run.font.name = base_font_name
                                if base_font_color:
                                    run.font.color.rgb = base_font_color

                        # 성과 섹션이면 테이블 추가 예약
                        if section_key == '성과':
                            from pptx.util import Emu, Pt as _Pt
                            # 전체 텍스트프레임의 paragraph 수로 높이 추정
                            # 한 줄당 약 0.28인치 (폰트+줄간격+여백 고려)
                            para_count = len(tf.paragraphs)
                            line_h = Emu(254000)  # 약 0.28인치/줄
                            estimated_h = para_count * line_h
                            actual_h = max(shape.height, estimated_h)
                            table_top = shape.top + actual_h + Emu(180000)
                            perf_tables_to_add.append((slide, shape.left, table_top, shape.width))

                        placeholders[ph_key] = ''

                # 성과 테이블 실제 추가
                for _, tbl_left, tbl_top, tbl_width in perf_tables_to_add:
                    _add_performance_table(slide, project_data, tbl_left, tbl_top, tbl_width)

        def replace_detail_images(prs_obj, project_data, label_map, placeholders):
            """{{이미지_좌측}} / {{이미지_우측}} 플레이스홀더를 찾아 해당 텍스트박스 위치에 이미지를 배치
            - 이미지 전용 텍스트박스: shape 삭제 후 이미지로 대체
            - 텍스트와 혼합된 텍스트박스: 플레이스홀더 텍스트만 제거, 텍스트박스 하단에 이미지 배치
            이전 {{상세정보_XXX_이미지}} 플레이스홀더도 호환성을 위해 빈 문자열로 처리
            """
            import base64
            from pptx.util import Emu
            from pptx.util import Pt as PtUtil
            from pptx.enum.text import PP_ALIGN

            # 이전 형식 플레이스홀더 제거 (호환성)
            for key in label_map:
                old_ph = f'{{{{상세정보_{key}_이미지}}}}'
                placeholders[old_ph] = ''

            # 이미지 플레이스홀더 (기존 좌측/우측 + 카테고리별)
            img_phs = {
                '{{이미지_좌측}}': '좌측',
                '{{이미지_우측}}': '우측',
                '{{이미지_개요그림}}': '개요그림',
                '{{이미지_상세내용그림}}': '상세내용그림',
                '{{이미지_향후계획그림}}': '향후계획그림',
            }
            for ph in img_phs:
                placeholders[ph] = ''

            for slide in prs_obj.slides:
                shapes_to_remove = []
                items_to_add = []

                for shape in list(slide.shapes):
                    if not hasattr(shape, 'text_frame'):
                        continue
                    full_text = ''.join(run.text for p in shape.text_frame.paragraphs for run in p.runs)

                    # 이 shape에서 매칭되는 이미지 플레이스홀더 모두 수집
                    matched = []
                    for ph, side in img_phs.items():
                        if ph in full_text:
                            matched.append((ph, side))

                    if not matched:
                        continue

                    # 이미지 플레이스홀더만 있는 전용 텍스트박스인지 확인
                    stripped = full_text
                    for ph, _ in matched:
                        stripped = stripped.replace(ph, '')
                    is_image_only = stripped.strip() == ''

                    area_left = shape.left
                    area_top = shape.top
                    area_width = shape.width
                    area_height = shape.height

                    if is_image_only:
                        all_images = []
                        for ph, side in matched:
                            images = project_data.get(f'이미지_{side}')
                            if images and isinstance(images, list):
                                for img in images:
                                    all_images.append(img)

                        if all_images:
                            gap = Emu(72000)
                            img_count = min(len(all_images), 6)
                            if img_count == 1:
                                positions = [(area_left, area_top, area_width, area_height)]
                            elif img_count == 2:
                                w = (area_width - gap) // 2
                                positions = [
                                    (area_left, area_top, w, area_height),
                                    (area_left + w + gap, area_top, w, area_height),
                                ]
                            elif img_count == 3:
                                w = (area_width - gap * 2) // 3
                                positions = [
                                    (area_left, area_top, w, area_height),
                                    (area_left + w + gap, area_top, w, area_height),
                                    (area_left + (w + gap) * 2, area_top, w, area_height),
                                ]
                            elif img_count == 4:
                                w = (area_width - gap) // 2
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left, area_top + h + gap, w, h),
                                    (area_left + w + gap, area_top + h + gap, w, h),
                                ]
                            else:
                                # 5~6개: 상단 3개, 하단 나머지
                                w = (area_width - gap * 2) // 3
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left + (w + gap) * 2, area_top, w, h),
                                ]
                                bottom_count = img_count - 3
                                if bottom_count == 1:
                                    positions.append((area_left, area_top + h + gap, area_width, h))
                                elif bottom_count == 2:
                                    bw = (area_width - gap) // 2
                                    positions.append((area_left, area_top + h + gap, bw, h))
                                    positions.append((area_left + bw + gap, area_top + h + gap, bw, h))
                                else:
                                    positions.append((area_left, area_top + h + gap, w, h))
                                    positions.append((area_left + w + gap, area_top + h + gap, w, h))
                                    positions.append((area_left + (w + gap) * 2, area_top + h + gap, w, h))

                            for i in range(img_count):
                                img_data = all_images[i]
                                data_url = img_data.get('dataUrl', '')
                                caption = img_data.get('caption', '')
                                px, py, pw, p_h = positions[i]
                                cap_h = Emu(200000) if caption else 0
                                img_h = p_h - cap_h

                                if data_url and ',' in data_url:
                                    try:
                                        img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                        img_stream = io.BytesIO(img_bytes)
                                        items_to_add.append(('picture', img_stream, px, py, pw, img_h))
                                    except Exception:
                                        pass
                                if caption:
                                    items_to_add.append(('caption', caption, px, py + img_h, pw, cap_h))

                        shapes_to_remove.append(shape)

                    else:
                        # 텍스트와 혼합: 플레이스홀더 텍스트만 제거 (이미지는 shape 바로 아래에 배치)
                        img_y = area_top + area_height + Emu(36000)
                        for ph, side in matched:
                            for p in shape.text_frame.paragraphs:
                                for run in p.runs:
                                    if ph in run.text:
                                        run.text = run.text.replace(ph, '')

                            images = project_data.get(f'이미지_{side}')
                            if not images or not isinstance(images, list) or len(images) == 0:
                                continue

                            images = images[:3]
                            gap = Emu(72000)
                            img_h = Emu(2000000)

                            if len(images) == 1:
                                positions = [(area_left, img_y, area_width, img_h)]
                            elif len(images) == 2:
                                w = (area_width - gap) // 2
                                positions = [
                                    (area_left, img_y, w, img_h),
                                    (area_left + w + gap, img_y, w, img_h),
                                ]
                            else:
                                w = (area_width - gap * 2) // 3
                                positions = [
                                    (area_left, img_y, w, img_h),
                                    (area_left + w + gap, img_y, w, img_h),
                                    (area_left + (w + gap) * 2, img_y, w, img_h),
                                ]

                            for i, img_data in enumerate(images):
                                data_url = img_data.get('dataUrl', '')
                                caption = img_data.get('caption', '')
                                px, py, pw, p_h = positions[i]
                                cap_h = Emu(200000) if caption else 0
                                pic_h = p_h - cap_h

                                if data_url and ',' in data_url:
                                    try:
                                        img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                        img_stream = io.BytesIO(img_bytes)
                                        items_to_add.append(('picture', img_stream, px, py, pw, pic_h))
                                    except Exception:
                                        pass
                                if caption:
                                    items_to_add.append(('caption', caption, px, py + pic_h, pw, cap_h))

                            img_y += img_h + Emu(72000)

                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)

                for item in items_to_add:
                    if item[0] == 'picture':
                        _, stream, x, y, w, h = item
                        slide.shapes.add_picture(stream, x, y, w, h)
                    elif item[0] == 'caption':
                        _, text, x, y, w, h = item
                        txBox = slide.shapes.add_textbox(x, y, w, h)
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        p.text = text
                        p.font.size = PtUtil(9)
                        p.font.italic = True
                        p.alignment = PP_ALIGN.CENTER

        def replace_milestone_timeline(prs_obj, project_data, placeholders):
            """{{액션아이템_마일스톤}} 플레이스홀더를 찾아 마일스톤 타임라인 도형으로 대체"""
            from pptx.util import Inches as InUtil, Pt as PtUtil, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from datetime import datetime

            milestone_ph = '{{액션아이템_마일스톤}}'
            placeholders[milestone_ph] = ''

            action_list = project_data.get('액션아이템목록', [])
            if not action_list:
                return

            # 액션아이템에서 날짜 파싱 및 정렬
            milestones = []
            for item in action_list:
                content = item.get('제목') or item.get('content') or item.get('내용', '')
                due_str = item.get('목표일') or item.get('dueDate') or ''
                done = item.get('완료여부', False)
                if not due_str or due_str == '-':
                    continue
                try:
                    due_date = datetime.strptime(str(due_str)[:10], '%Y-%m-%d')
                except (ValueError, TypeError):
                    continue
                milestones.append({
                    'name': content,
                    'date': due_date,
                    'done': done,
                    'label': f"{due_date.month}월",
                })

            if not milestones:
                return

            milestones.sort(key=lambda m: m['date'])

            # 같은 날짜 그룹에서 마지막 항목만 "월" 표시
            for i, ms in enumerate(milestones):
                show_label = True
                for j in range(i + 1, len(milestones)):
                    if milestones[j]['date'] == ms['date']:
                        show_label = False
                        break
                ms['show_month'] = show_label

            for slide in prs_obj.slides:
                shapes_to_remove = []

                for shape in list(slide.shapes):
                    if not hasattr(shape, 'text_frame'):
                        continue
                    full_text = ''.join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
                    if milestone_ph not in full_text:
                        continue

                    area_left = shape.left
                    area_top = shape.top
                    area_width = shape.width
                    area_height = shape.height
                    shapes_to_remove.append(shape)

                    # 제목 + 타임라인 배치
                    # 1) 먼저 위쪽 라벨 최대 높이를 계산하여 line_y를 결정
                    title_h = Emu(320000)
                    title_gap = Emu(80000)
                    dot_radius = Emu(54000)
                    label_w = Emu(1200000)
                    label_h = Emu(420000)
                    stem_h = Emu(120000)
                    use_zigzag = len(milestones) > 4
                    n = len(milestones)
                    margin = Emu(400000)

                    # 위쪽 라벨 최대 높이 사전 계산
                    max_top_label_h = Emu(0)
                    for i, ms in enumerate(milestones):
                        is_top = True if not use_zigzag else (i % 2 == 0)
                        if is_top:
                            name_len = len(ms['name'])
                            est_lines = max(1, (name_len + 7) // 8)
                            if ms['show_month']:
                                est_lines += 1
                            h = Emu(150000 + est_lines * 120000)
                            if h > max_top_label_h:
                                max_top_label_h = h

                    # 2) 제목을 플레이스홀더 위치에 배치
                    titleBox = slide.shapes.add_textbox(
                        area_left, area_top, area_width, title_h
                    )
                    title_tf = titleBox.text_frame
                    title_tf.word_wrap = False
                    title_p = title_tf.paragraphs[0]
                    title_p.alignment = PP_ALIGN.LEFT
                    title_p.space_before = PtUtil(0)
                    title_p.space_after = PtUtil(0)
                    title_run = title_p.add_run()
                    title_run.text = '□ 과제 일정'
                    title_run.font.size = PtUtil(14)
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(30, 41, 59)

                    # 3) line_y = 제목 아래 + 위쪽라벨 최대높이 + stem + dot 여유
                    line_y = area_top + title_h + title_gap + max_top_label_h + stem_h + dot_radius

                    line_left = area_left + margin
                    line_width = area_width - margin * 2

                    # 가로선
                    line_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        line_left, line_y - Emu(12000),
                        line_width, Emu(24000)
                    )
                    line_shape.fill.solid()
                    line_shape.fill.fore_color.rgb = RGBColor(156, 163, 175)
                    line_shape.line.fill.background()

                    # 화살표
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.ISOSCELES_TRIANGLE,
                        line_left + line_width - Emu(5000),
                        line_y - Emu(60000),
                        Emu(120000), Emu(120000)
                    )
                    arrow.rotation = 90.0
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(156, 163, 175)
                    arrow.line.fill.background()

                    for i, ms in enumerate(milestones):
                        # 등간격 X 좌표
                        if n == 1:
                            ratio = 0.5
                        else:
                            ratio = i / (n - 1)
                        dot_x = line_left + int(line_width * ratio)

                        color = RGBColor(16, 185, 129) if ms['done'] else RGBColor(59, 130, 246)
                        dot = slide.shapes.add_shape(
                            MSO_SHAPE.OVAL,
                            dot_x - dot_radius, line_y - dot_radius,
                            dot_radius * 2, dot_radius * 2
                        )
                        dot.fill.solid()
                        dot.fill.fore_color.rgb = color
                        dot.line.fill.background()

                        is_top = True if not use_zigzag else (i % 2 == 0)

                        if is_top:
                            stem_top = line_y - dot_radius - stem_h
                        else:
                            stem_top = line_y + dot_radius

                        stem = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            dot_x - Emu(6000), stem_top,
                            Emu(12000), stem_h
                        )
                        stem.fill.solid()
                        stem.fill.fore_color.rgb = color
                        stem.line.fill.background()

                        if is_top:
                            # 텍스트 길이에 따라 높이 동적 계산 (1줄 ~150000, 추가 줄당 ~120000)
                            name_len = len(ms['name'])
                            est_lines = max(1, (name_len + 7) // 8)  # label_w 기준 약 8자/줄
                            if ms['show_month']:
                                est_lines += 1
                            top_label_h = Emu(150000 + est_lines * 120000)
                            label_top = stem_top - top_label_h
                        else:
                            label_top = stem_top + stem_h

                        label_left = dot_x - label_w // 2
                        if label_left < area_left:
                            label_left = area_left
                        if label_left + label_w > area_left + area_width:
                            label_left = area_left + area_width - label_w

                        label_center_x = label_left + label_w // 2
                        dot._element.spPr.xfrm.off.x = label_center_x - dot_radius
                        stem._element.spPr.xfrm.off.x = label_center_x - Emu(6000)

                        actual_label_h = top_label_h if is_top else label_h
                        txBox = slide.shapes.add_textbox(label_left, label_top, label_w, actual_label_h)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        if is_top:
                            tf.paragraphs[0].vertical_anchor = MSO_ANCHOR.BOTTOM

                        if is_top:
                            # 위쪽: 이름(바깥=위) → 월(바 가까이=아래)
                            p = tf.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            p.space_before = PtUtil(0)
                            p.space_after = PtUtil(0)
                            run = p.add_run()
                            run.text = ms['name']
                            run.font.size = PtUtil(8)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(55, 65, 81)

                            if ms['show_month']:
                                p2 = tf.add_paragraph()
                                p2.alignment = PP_ALIGN.CENTER
                                p2.space_before = PtUtil(0)
                                p2.space_after = PtUtil(0)
                                run2 = p2.add_run()
                                run2.text = ms['label']
                                run2.font.size = PtUtil(8)
                                run2.font.color.rgb = color
                        else:
                            # 아래쪽: 월(바 가까이=위) → 이름(바깥=아래)
                            p = tf.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            p.space_before = PtUtil(0)
                            p.space_after = PtUtil(0)
                            if ms['show_month']:
                                run_month = p.add_run()
                                run_month.text = ms['label']
                                run_month.font.size = PtUtil(8)
                                run_month.font.color.rgb = color

                                p2 = tf.add_paragraph()
                                p2.alignment = PP_ALIGN.CENTER
                                p2.space_before = PtUtil(0)
                                p2.space_after = PtUtil(0)
                                run_name = p2.add_run()
                            else:
                                run_name = p.add_run()

                            run_name.text = ms['name']
                            run_name.font.size = PtUtil(8)
                            run_name.font.bold = True
                            run_name.font.color.rgb = RGBColor(55, 65, 81)

                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)

        replace_detail_images(prs, project_data, DETAIL_LABEL_MAP, placeholders)
        replace_milestone_timeline(prs, project_data, placeholders)
        replace_detail_placeholders(prs, project_data)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        replace_paragraph_placeholders(paragraph, placeholders)

        # 텍스트가 긴 텍스트 박스의 폰트를 자동 축소
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                tf = shape.text_frame
                full = ''.join(p.text for p in tf.paragraphs)
                line_count = full.count('\n') + 1
                if line_count > 5:
                    tf.word_wrap = True
                    box_height_inches = shape.height / 914400
                    available_lines = box_height_inches / 0.22
                    if line_count > available_lines:
                        target_pt = max(7, int(10 * available_lines / line_count))
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(target_pt)

        # normAutofit의 fontScale 리셋 (PowerPoint가 열 때 재계산하도록)
        _reset_autofit_font_scale(prs)

        # 파일명 생성 (특수문자 제거)
        safe_filename = "".join(c for c in (project_data.get('과제명') or '과제보고서') if c not in r'\/:*?"<>|')
        filename = f"{safe_filename[:50]}_{project_data.get('과제년도', '')}년_보고서.pptx"

        # 메모리에 저장
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        # 활동 로그 기록
        report_log = DashboardActivityLog(
            action='EXPORT',
            target_type='REPORT',
            target_id=project_data.get('id'),
            target_name=project_data.get('과제명', '알 수 없는 과제'),
            summary=f'{user.name}이(가) PPT 보고서를 생성함',
            user_id=user_id,
            user_name=user.name,
            source='server'
        )
        db.session.add(report_log)
        db.session.commit()

        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        print(f"[Report Error] PPT generation failed: {str(e)}")
        return error_response(f'PPT 보고서 생성 실패: {str(e)}', status_code=500)


@bp.route('/report/ppt/batch', methods=['POST'])
@jwt_required()
def generate_batch_ppt_report():
    """여러 과제 PPT 보고서 일괄 생성 (관리자 전용)
    mode: 'single' — 모든 과제를 하나의 PPTX로 합침
    mode: 'individual' — 과제별 개별 PPTX를 ZIP으로 묶어 반환
    """
    try:
        from pptx import Presentation
        import re
        import zipfile

        user_id = get_jwt_identity()
        user = User.query.get(user_id)

        # manager, dt_office, admin만 허용
        if not user or (user.role not in ('admin', 'manager', 'dt_office') and not user.is_admin):
            return error_response('보고서 저장 권한이 없습니다.', status_code=403)

        body = request.get_json()
        if not body:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        projects_list = body.get('projects', [])
        mode = body.get('mode', 'single')  # 'single' or 'individual'
        project_template_map = body.get('projectTemplateMap', {})  # { projectId: templateFilename }

        if not projects_list:
            return error_response('선택된 과제가 없습니다.', status_code=400)

        # imageId 만 온 이미지에 dataUrl 을 채운다 — 아래 삽입부는 dataUrl 만 읽는다
        hydrate_report_images(projects_list)

        try:
            template_path = get_template_path(body.get('template'))
        except (ValueError, FileNotFoundError) as e:
            return error_response(str(e), status_code=400)

        def strip_html(html_text):
            if not html_text:
                return '-'
            text = re.sub(r'<[^>]+>', '', html_text)
            text = text.replace('&nbsp;', ' ')
            text = text.replace('&lt;', '<')
            text = text.replace('&gt;', '>')
            text = text.replace('&amp;', '&')
            text = text.replace('&quot;', '"')
            text = re.sub(r'\s+', ' ', text).strip()
            return text if text else '-'

        def build_placeholders(project_data):
            placeholders = {
                '{{과제명}}': project_data.get('과제명', '-'),
                '{{과제년도}}': str(project_data.get('과제년도', '-')),
                '{{사업부}}': project_data.get('사업부', '-'),
                '{{프로세스}}': project_data.get('프로세스', '-'),
                '{{과제PL}}': project_data.get('과제PL', '-'),
                '{{작성자}}': project_data.get('작성자', '-'),
                '{{시작}}': str(project_data.get('시작', '-')),
                '{{종료}}': str(project_data.get('종료', '-')),
                '{{과제영역}}': project_data.get('과제영역', '-'),
                '{{과제구분}}': project_data.get('과제구분', '-'),
                '{{진행상태}}': project_data.get('진행상태', '-'),
                '{{진행률}}': str(project_data.get('진행률', 0)),
                '{{PoC과제여부}}': '예' if project_data.get('PoC과제여부') else '아니오',
                '{{중점과제여부}}': '예' if project_data.get('중점과제여부') else '아니오',
                '{{과제상세설명}}': strip_html(project_data.get('과제상세설명', '')),
            }
            dept_list = project_data.get('담당부서목록', [])
            placeholders['{{담당부서목록}}'] = ', '.join(dept_list) if dept_list else '-'

            personnel_list = project_data.get('과제참여인력목록', [])
            if personnel_list:
                placeholders['{{과제참여인력목록}}'] = '\n'.join([
                    f"• {p.get('이름', '-')} ({p.get('부서', '-')})" for p in personnel_list
                ])
                placeholders['{{과제참여인력_이름목록}}'] = ', '.join([
                    p.get('이름', '-') for p in personnel_list
                ])
            else:
                placeholders['{{과제참여인력목록}}'] = '등록된 참여인력이 없습니다.'
                placeholders['{{과제참여인력_이름목록}}'] = '-'
            # 소속별 그룹화: 부서 1개면 이름만, 2개 이상이면 "소속A 이름1, 이름2 / 소속B 이름3"
            if personnel_list:
                from collections import OrderedDict
                dept_groups = OrderedDict()
                for p in personnel_list:
                    dept = p.get('부서', '-')
                    name = p.get('이름', '-')
                    dept_groups.setdefault(dept, []).append(name)
                if len(dept_groups) == 1:
                    placeholders['{{과제참여인력_소속별목록}}'] = ', '.join(list(dept_groups.values())[0])
                else:
                    placeholders['{{과제참여인력_소속별목록}}'] = ' / '.join(
                        f"{dept} {', '.join(names)}" for dept, names in dept_groups.items()
                    )
            else:
                placeholders['{{과제참여인력_소속별목록}}'] = '-'
            placeholders['{{참여인력수}}'] = str(len(personnel_list))
            for i in range(20):
                idx = i + 1
                if i < len(personnel_list):
                    p = personnel_list[i]
                    email = p.get('knoxId', '') or ''
                    placeholders[f'{{{{참여인력{idx}_이름}}}}'] = p.get('이름', '-')
                    placeholders[f'{{{{참여인력{idx}_부서}}}}'] = p.get('부서', '-')
                    placeholders[f'{{{{참여인력{idx}_knoxId}}}}'] = email
                    placeholders[f'{{{{참여인력{idx}_선택사업부}}}}'] = p.get('선택사업부', '-')
                else:
                    placeholders[f'{{{{참여인력{idx}_이름}}}}'] = ''
                    placeholders[f'{{{{참여인력{idx}_부서}}}}'] = ''
                    placeholders[f'{{{{참여인력{idx}_knoxId}}}}'] = ''
                    placeholders[f'{{{{참여인력{idx}_선택사업부}}}}'] = ''

            performance_list = project_data.get('성과목록', [])
            if performance_list:
                perf_lines = []
                for i, p in enumerate(performance_list, 1):
                    line = f"{i}. {p.get('성과항목', '-')}"
                    line += f"\n   - 대분류: {p.get('대분류', '-')}"
                    line += f"\n   - 기여도: {p.get('과제기여도', '-')}%"
                    line += f"\n   - 현재: {p.get('현재수준', '-')} / 목표: {p.get('목표수준', '-')} / 실적: {p.get('실적수준', '-')} ({p.get('단위', '-')})"
                    perf_lines.append(line)
                placeholders['{{성과목록}}'] = '\n\n'.join(perf_lines)
            else:
                placeholders['{{성과목록}}'] = '등록된 성과가 없습니다.'
            placeholders['{{성과수}}'] = str(len(performance_list))
            for i in range(20):
                idx = i + 1
                if i < len(performance_list):
                    p = performance_list[i]
                    placeholders[f'{{{{성과{idx}_항목}}}}'] = p.get('성과항목', '-')
                    placeholders[f'{{{{성과{idx}_대분류}}}}'] = p.get('대분류', '-')
                    placeholders[f'{{{{성과{idx}_기여도}}}}'] = str(p.get('과제기여도', '-'))
                    placeholders[f'{{{{성과{idx}_현재수준}}}}'] = str(p.get('현재수준', '-'))
                    placeholders[f'{{{{성과{idx}_목표수준}}}}'] = str(p.get('목표수준', '-'))
                    placeholders[f'{{{{성과{idx}_실적수준}}}}'] = str(p.get('실적수준', '-'))
                    placeholders[f'{{{{성과{idx}_단위}}}}'] = p.get('단위', '-')
                else:
                    placeholders[f'{{{{성과{idx}_항목}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_대분류}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_기여도}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_현재수준}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_목표수준}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_실적수준}}}}'] = ''
                    placeholders[f'{{{{성과{idx}_단위}}}}'] = ''

            action_list = project_data.get('액션아이템목록', [])
            if action_list:
                action_lines = []
                for item in action_list:
                    content = item.get('제목') or item.get('content') or item.get('내용', '-')
                    due_date = item.get('목표일') or item.get('dueDate') or item.get('완료일', '-')
                    done = '완료' if item.get('완료여부') else '진행중'
                    status = item.get('status') or item.get('상태') or done
                    action_lines.append(f"• {content} (완료일: {due_date}, 상태: {status})")
                placeholders['{{액션아이템목록}}'] = '\n'.join(action_lines)
            else:
                placeholders['{{액션아이템목록}}'] = '등록된 액션아이템이 없습니다.'
            placeholders['{{액션아이템수}}'] = str(len(action_list))
            for i in range(20):
                idx = i + 1
                if i < len(action_list):
                    item = action_list[i]
                    content = item.get('제목') or item.get('content') or item.get('내용', '-')
                    due_date = item.get('목표일') or item.get('dueDate') or item.get('완료일', '-')
                    done_flag = item.get('완료여부', False)
                    completion_date = item.get('완료일', '-')
                    sub_items = item.get('세부항목목록', [])
                    sub_text = ', '.join([s.get('내용', '-') for s in sub_items]) if sub_items else '-'
                    placeholders[f'{{{{액션아이템{idx}_이름}}}}'] = content
                    placeholders[f'{{{{액션아이템{idx}_목표일}}}}'] = str(due_date)
                    placeholders[f'{{{{액션아이템{idx}_완료여부}}}}'] = '완료' if done_flag else '진행중'
                    placeholders[f'{{{{액션아이템{idx}_완료일}}}}'] = str(completion_date)
                    placeholders[f'{{{{액션아이템{idx}_세부항목}}}}'] = sub_text
                else:
                    placeholders[f'{{{{액션아이템{idx}_이름}}}}'] = ''
                    placeholders[f'{{{{액션아이템{idx}_목표일}}}}'] = ''
                    placeholders[f'{{{{액션아이템{idx}_완료여부}}}}'] = ''
                    placeholders[f'{{{{액션아이템{idx}_완료일}}}}'] = ''
                    placeholders[f'{{{{액션아이템{idx}_세부항목}}}}'] = ''

            # 상세 과제 정보 — 일반 텍스트 플레이스홀더 (fallback용, 멀티스타일 처리 전)
            for section_key in DETAIL_SECTION_MAP:
                placeholders[f'{{{{상세정보_{section_key}}}}}'] = ''

            return placeholders

        def replace_paragraph_placeholders(paragraph, placeholders):
            runs = paragraph.runs
            if not runs:
                return
            full_text = ''.join(run.text for run in runs)
            has_match = False
            for ph in placeholders:
                if ph in full_text:
                    has_match = True
                    break
            if not has_match:
                return
            for ph, val in placeholders.items():
                full_text = full_text.replace(ph, str(val))
            runs[0].text = full_text
            for run in runs[1:]:
                run.text = ''

        # ── 상세 과제 정보: 스타일 추출 + 동적 슬라이드 생성 ──

        DETAIL_SECTION_MAP = {
            '과제개요': '과제 개요',
            '추진배경': '추진 배경',
            '과제목표': '과제 목표',
            '상세내용': '상세 내용',
            '성과': '기술/경영 성과',
            '산출물': '산출물',
            '향후계획': '향후 계획',
        }

        def extract_detail_styles(prs_obj):
            """템플릿의 상세 과제 정보 슬라이드에서 스타일을 추출하고 슬라이드 인덱스를 반환"""
            styles = {}
            detail_slide_idx = None
            for idx, slide in enumerate(prs_obj.slides):
                for shape in slide.shapes:
                    if shape.name in ('DETAIL_TITLE_STYLE', 'DETAIL_PARENT_STYLE',
                                      'DETAIL_CHILD_STYLE', 'DETAIL_CAPTION_STYLE'):
                        detail_slide_idx = idx
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                            if run:
                                styles[shape.name] = {
                                    'font_name': run.font.name,
                                    'font_size': run.font.size,
                                    'font_bold': run.font.bold,
                                    'font_italic': run.font.italic,
                                    'font_color': run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                                }
                                para = shape.text_frame.paragraphs[0]
                                styles[shape.name]['alignment'] = para.alignment
                        # 위치/크기 정보도 저장
                        styles.setdefault(shape.name, {})
                        styles[shape.name]['left'] = shape.left
                        styles[shape.name]['width'] = shape.width
            return styles, detail_slide_idx

        def apply_style(run, style_dict):
            """추출한 스타일을 run에 적용"""
            if not style_dict:
                return
            if style_dict.get('font_name'):
                run.font.name = style_dict['font_name']
            if style_dict.get('font_size'):
                run.font.size = style_dict['font_size']
            if style_dict.get('font_bold') is not None:
                run.font.bold = style_dict['font_bold']
            if style_dict.get('font_italic') is not None:
                run.font.italic = style_dict['font_italic']
            if style_dict.get('font_color'):
                run.font.color.rgb = style_dict['font_color']

        def build_detail_slides(prs_obj, project_data, styles, detail_slide_idx):
            """활성화된 상세 과제 정보 섹션마다 슬라이드를 동적 생성하고, 원본 스타일 슬라이드를 삭제"""
            from pptx.util import Inches as InUtil, Pt as PtUtil, Emu
            import base64

            sections_to_generate = []
            for key, label in DETAIL_SECTION_MAP.items():
                field = f'상세정보_{key}'
                raw = project_data.get(field)
                if raw and raw.get('enabled'):
                    sections_to_generate.append((key, label, raw))

            if not sections_to_generate and detail_slide_idx is not None:
                # 활성 섹션 없으면 스타일 슬라이드만 삭제
                rId = prs_obj.slides._sldIdLst[detail_slide_idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId is None:
                    rId = prs_obj.slides._sldIdLst[detail_slide_idx].attrib.get('r:id')
                prs_obj.part.drop_rel(rId)
                del prs_obj.slides._sldIdLst[detail_slide_idx]
                return

            title_style = styles.get('DETAIL_TITLE_STYLE', {})
            parent_style = styles.get('DETAIL_PARENT_STYLE', {})
            child_style = styles.get('DETAIL_CHILD_STYLE', {})
            caption_style = styles.get('DETAIL_CAPTION_STYLE', {})

            # 삽입 위치: 스타일 슬라이드 위치 (나중에 원본 삭제)
            insert_after_idx = detail_slide_idx

            for sec_idx, (key, label, raw) in enumerate(sections_to_generate):
                slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[6])
                items = raw.get('items', [])
                images = raw.get('images', [])

                content_left = title_style.get('left', InUtil(0.5))
                content_width = title_style.get('width', InUtil(9))

                # 제목
                y = InUtil(0.3)
                txBox = slide.shapes.add_textbox(content_left, y, content_width, InUtil(0.6))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = label
                if p.runs:
                    apply_style(p.runs[0], title_style)
                else:
                    run = p.add_run()
                    run.text = label
                    p.runs[0].text = ''
                    apply_style(run, title_style)
                if title_style.get('alignment'):
                    p.alignment = title_style['alignment']

                # 텍스트 항목 배치
                parent_left = parent_style.get('left', InUtil(0.5))
                parent_width = parent_style.get('width', InUtil(9))
                child_left = child_style.get('left', InUtil(0.9))
                child_width = child_style.get('width', InUtil(8.6))

                y = InUtil(1.2)
                line_h_parent = InUtil(0.35)
                line_h_child = InUtil(0.30)

                for item in items:
                    if isinstance(item, str):
                        item = {'text': item, 'children': []}
                    text = item.get('text', '')
                    if not text:
                        continue

                    # 상위 항목
                    txBox = slide.shapes.add_textbox(parent_left, y, parent_width, line_h_parent)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"- {text}"
                    if p.runs:
                        apply_style(p.runs[0], parent_style)
                    if parent_style.get('alignment'):
                        p.alignment = parent_style['alignment']
                    y += line_h_parent

                    # 하위 항목들
                    for child in item.get('children', []):
                        child_text = child.get('text', '') if isinstance(child, dict) else str(child)
                        if not child_text:
                            continue
                        txBox = slide.shapes.add_textbox(child_left, y, child_width, line_h_child)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = f"\u00B7 {child_text}"
                        if p.runs:
                            apply_style(p.runs[0], child_style)
                        if child_style.get('alignment'):
                            p.alignment = child_style['alignment']
                        y += line_h_child

                # 이미지 배치 (하단 영역)
                if images:
                    img_top = max(y + InUtil(0.3), InUtil(4.0))
                    img_h = InUtil(2.5)
                    img_count = min(len(images), 2)

                    if img_count == 1:
                        img_w = InUtil(4.5)
                        img_x = InUtil(2.75)  # 중앙
                        positions = [(img_x, img_top)]
                    else:
                        img_w = InUtil(4.0)
                        positions = [(InUtil(0.5), img_top), (InUtil(5.2), img_top)]

                    for i in range(img_count):
                        img_data = images[i]
                        data_url = img_data.get('dataUrl', '')
                        caption = img_data.get('caption', '')

                        if data_url and ',' in data_url:
                            try:
                                img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                img_stream = io.BytesIO(img_bytes)
                                pic = slide.shapes.add_picture(img_stream, positions[i][0], positions[i][1], img_w, img_h)
                            except Exception:
                                pass  # 이미지 디코딩 실패 시 스킵

                        if caption:
                            cap_y = positions[i][1] + img_h + Emu(36000)
                            txBox = slide.shapes.add_textbox(positions[i][0], cap_y, img_w, InUtil(0.3))
                            tf = txBox.text_frame
                            p = tf.paragraphs[0]
                            p.text = caption
                            if p.runs:
                                apply_style(p.runs[0], caption_style)
                            if caption_style.get('alignment'):
                                p.alignment = caption_style['alignment']

            # 원본 스타일 슬라이드 삭제
            if detail_slide_idx is not None:
                rId = prs_obj.slides._sldIdLst[detail_slide_idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId is None:
                    rId = prs_obj.slides._sldIdLst[detail_slide_idx].attrib.get('r:id')
                prs_obj.part.drop_rel(rId)
                del prs_obj.slides._sldIdLst[detail_slide_idx]


        def replace_detail_images(prs_obj, project_data, label_map, placeholders):
            """{{이미지_좌측}} / {{이미지_우측}} 플레이스홀더를 찾아 해당 텍스트박스 위치에 이미지를 배치
            - 이미지 전용 텍스트박스: shape 삭제 후 이미지로 대체
            - 텍스트와 혼합된 텍스트박스: 플레이스홀더 텍스트만 제거, 텍스트박스 하단에 이미지 배치
            이전 {{상세정보_XXX_이미지}} 플레이스홀더도 호환성을 위해 빈 문자열로 처리
            """
            import base64
            from pptx.util import Emu
            from pptx.util import Pt as PtUtil
            from pptx.enum.text import PP_ALIGN

            # 이전 형식 플레이스홀더 제거 (호환성)
            for key in label_map:
                old_ph = f'{{{{상세정보_{key}_이미지}}}}'
                placeholders[old_ph] = ''

            # 이미지 플레이스홀더 (기존 좌측/우측 + 카테고리별)
            img_phs = {
                '{{이미지_좌측}}': '좌측',
                '{{이미지_우측}}': '우측',
                '{{이미지_개요그림}}': '개요그림',
                '{{이미지_상세내용그림}}': '상세내용그림',
                '{{이미지_향후계획그림}}': '향후계획그림',
            }
            for ph in img_phs:
                placeholders[ph] = ''

            for slide in prs_obj.slides:
                shapes_to_remove = []
                items_to_add = []

                for shape in list(slide.shapes):
                    if not hasattr(shape, 'text_frame'):
                        continue
                    full_text = ''.join(run.text for p in shape.text_frame.paragraphs for run in p.runs)

                    matched = []
                    for ph, side in img_phs.items():
                        if ph in full_text:
                            matched.append((ph, side))

                    if not matched:
                        continue

                    stripped = full_text
                    for ph, _ in matched:
                        stripped = stripped.replace(ph, '')
                    is_image_only = stripped.strip() == ''

                    area_left = shape.left
                    area_top = shape.top
                    area_width = shape.width
                    area_height = shape.height

                    if is_image_only:
                        all_images = []
                        for ph, side in matched:
                            images = project_data.get(f'이미지_{side}')
                            if images and isinstance(images, list):
                                for img in images:
                                    all_images.append(img)

                        if all_images:
                            gap = Emu(72000)
                            img_count = min(len(all_images), 6)
                            if img_count == 1:
                                positions = [(area_left, area_top, area_width, area_height)]
                            elif img_count == 2:
                                w = (area_width - gap) // 2
                                positions = [
                                    (area_left, area_top, w, area_height),
                                    (area_left + w + gap, area_top, w, area_height),
                                ]
                            elif img_count == 3:
                                w = (area_width - gap * 2) // 3
                                positions = [
                                    (area_left, area_top, w, area_height),
                                    (area_left + w + gap, area_top, w, area_height),
                                    (area_left + (w + gap) * 2, area_top, w, area_height),
                                ]
                            elif img_count == 4:
                                w = (area_width - gap) // 2
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left, area_top + h + gap, w, h),
                                    (area_left + w + gap, area_top + h + gap, w, h),
                                ]
                            else:
                                # 5~6개: 상단 3개, 하단 나머지
                                w = (area_width - gap * 2) // 3
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left + (w + gap) * 2, area_top, w, h),
                                ]
                                bottom_count = img_count - 3
                                if bottom_count == 1:
                                    positions.append((area_left, area_top + h + gap, area_width, h))
                                elif bottom_count == 2:
                                    bw = (area_width - gap) // 2
                                    positions.append((area_left, area_top + h + gap, bw, h))
                                    positions.append((area_left + bw + gap, area_top + h + gap, bw, h))
                                else:
                                    positions.append((area_left, area_top + h + gap, w, h))
                                    positions.append((area_left + w + gap, area_top + h + gap, w, h))
                                    positions.append((area_left + (w + gap) * 2, area_top + h + gap, w, h))

                            for i in range(img_count):
                                img_data = all_images[i]
                                data_url = img_data.get('dataUrl', '')
                                caption = img_data.get('caption', '')
                                px, py, pw, p_h = positions[i]
                                cap_h = Emu(200000) if caption else 0
                                img_h = p_h - cap_h

                                if data_url and ',' in data_url:
                                    try:
                                        img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                        img_stream = io.BytesIO(img_bytes)
                                        items_to_add.append(('picture', img_stream, px, py, pw, img_h))
                                    except Exception:
                                        pass
                                if caption:
                                    items_to_add.append(('caption', caption, px, py + img_h, pw, cap_h))

                        shapes_to_remove.append(shape)

                    else:
                        img_y = area_top + area_height + Emu(36000)
                        for ph, side in matched:
                            for p in shape.text_frame.paragraphs:
                                for run in p.runs:
                                    if ph in run.text:
                                        run.text = run.text.replace(ph, '')

                            images = project_data.get(f'이미지_{side}')
                            if not images or not isinstance(images, list) or len(images) == 0:
                                continue

                            images = images[:3]
                            gap = Emu(72000)
                            img_h = Emu(2000000)

                            if len(images) == 1:
                                positions = [(area_left, img_y, area_width, img_h)]
                            elif len(images) == 2:
                                w = (area_width - gap) // 2
                                positions = [
                                    (area_left, img_y, w, img_h),
                                    (area_left + w + gap, img_y, w, img_h),
                                ]
                            else:
                                w = (area_width - gap * 2) // 3
                                positions = [
                                    (area_left, img_y, w, img_h),
                                    (area_left + w + gap, img_y, w, img_h),
                                    (area_left + (w + gap) * 2, img_y, w, img_h),
                                ]

                            for i, img_data in enumerate(images):
                                data_url = img_data.get('dataUrl', '')
                                caption = img_data.get('caption', '')
                                px, py, pw, p_h = positions[i]
                                cap_h = Emu(200000) if caption else 0
                                pic_h = p_h - cap_h

                                if data_url and ',' in data_url:
                                    try:
                                        img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                        img_stream = io.BytesIO(img_bytes)
                                        items_to_add.append(('picture', img_stream, px, py, pw, pic_h))
                                    except Exception:
                                        pass
                                if caption:
                                    items_to_add.append(('caption', caption, px, py + pic_h, pw, cap_h))

                            img_y += img_h + Emu(72000)

                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)

                for item in items_to_add:
                    if item[0] == 'picture':
                        _, stream, x, y, w, h = item
                        slide.shapes.add_picture(stream, x, y, w, h)
                    elif item[0] == 'caption':
                        _, text, x, y, w, h = item
                        txBox = slide.shapes.add_textbox(x, y, w, h)
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        p.text = text
                        p.font.size = PtUtil(9)
                        p.font.italic = True
                        p.alignment = PP_ALIGN.CENTER


        def replace_milestone_timeline(prs_obj, project_data, placeholders):
            """{{액션아이템_마일스톤}} 플레이스홀더를 찾아 마일스톤 타임라인 도형으로 대체"""
            from pptx.util import Inches as InUtil, Pt as PtUtil, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from datetime import datetime

            milestone_ph = '{{액션아이템_마일스톤}}'
            placeholders[milestone_ph] = ''

            action_list = project_data.get('액션아이템목록', [])
            if not action_list:
                return

            # 액션아이템에서 날짜 파싱 및 정렬
            milestones = []
            for item in action_list:
                content = item.get('제목') or item.get('content') or item.get('내용', '')
                due_str = item.get('목표일') or item.get('dueDate') or ''
                done = item.get('완료여부', False)
                if not due_str or due_str == '-':
                    continue
                try:
                    due_date = datetime.strptime(str(due_str)[:10], '%Y-%m-%d')
                except (ValueError, TypeError):
                    continue
                milestones.append({
                    'name': content,
                    'date': due_date,
                    'done': done,
                    'label': f"{due_date.month}월",
                })

            if not milestones:
                return

            milestones.sort(key=lambda m: m['date'])

            # 같은 날짜 그룹에서 마지막 항목만 "월" 표시
            for i, ms in enumerate(milestones):
                show_label = True
                for j in range(i + 1, len(milestones)):
                    if milestones[j]['date'] == ms['date']:
                        show_label = False
                        break
                ms['show_month'] = show_label

            for slide in prs_obj.slides:
                shapes_to_remove = []

                for shape in list(slide.shapes):
                    if not hasattr(shape, 'text_frame'):
                        continue
                    full_text = ''.join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
                    if milestone_ph not in full_text:
                        continue

                    area_left = shape.left
                    area_top = shape.top
                    area_width = shape.width
                    area_height = shape.height
                    shapes_to_remove.append(shape)

                    # 제목 + 타임라인 배치
                    # 1) 먼저 위쪽 라벨 최대 높이를 계산하여 line_y를 결정
                    title_h = Emu(320000)
                    title_gap = Emu(80000)
                    dot_radius = Emu(54000)
                    label_w = Emu(1200000)
                    label_h = Emu(420000)
                    stem_h = Emu(120000)
                    use_zigzag = len(milestones) > 4
                    n = len(milestones)
                    margin = Emu(400000)

                    # 위쪽 라벨 최대 높이 사전 계산
                    max_top_label_h = Emu(0)
                    for i, ms in enumerate(milestones):
                        is_top = True if not use_zigzag else (i % 2 == 0)
                        if is_top:
                            name_len = len(ms['name'])
                            est_lines = max(1, (name_len + 7) // 8)
                            if ms['show_month']:
                                est_lines += 1
                            h = Emu(150000 + est_lines * 120000)
                            if h > max_top_label_h:
                                max_top_label_h = h

                    # 2) 제목을 플레이스홀더 위치에 배치
                    titleBox = slide.shapes.add_textbox(
                        area_left, area_top, area_width, title_h
                    )
                    title_tf = titleBox.text_frame
                    title_tf.word_wrap = False
                    title_p = title_tf.paragraphs[0]
                    title_p.alignment = PP_ALIGN.LEFT
                    title_p.space_before = PtUtil(0)
                    title_p.space_after = PtUtil(0)
                    title_run = title_p.add_run()
                    title_run.text = '□ 과제 일정'
                    title_run.font.size = PtUtil(14)
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(30, 41, 59)

                    # 3) line_y = 제목 아래 + 위쪽라벨 최대높이 + stem + dot 여유
                    line_y = area_top + title_h + title_gap + max_top_label_h + stem_h + dot_radius

                    line_left = area_left + margin
                    line_width = area_width - margin * 2

                    # 가로선
                    line_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        line_left, line_y - Emu(12000),
                        line_width, Emu(24000)
                    )
                    line_shape.fill.solid()
                    line_shape.fill.fore_color.rgb = RGBColor(156, 163, 175)
                    line_shape.line.fill.background()

                    # 화살표
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.ISOSCELES_TRIANGLE,
                        line_left + line_width - Emu(5000),
                        line_y - Emu(60000),
                        Emu(120000), Emu(120000)
                    )
                    arrow.rotation = 90.0
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(156, 163, 175)
                    arrow.line.fill.background()

                    for i, ms in enumerate(milestones):
                        # 등간격 X 좌표
                        if n == 1:
                            ratio = 0.5
                        else:
                            ratio = i / (n - 1)
                        dot_x = line_left + int(line_width * ratio)

                        color = RGBColor(16, 185, 129) if ms['done'] else RGBColor(59, 130, 246)
                        dot = slide.shapes.add_shape(
                            MSO_SHAPE.OVAL,
                            dot_x - dot_radius, line_y - dot_radius,
                            dot_radius * 2, dot_radius * 2
                        )
                        dot.fill.solid()
                        dot.fill.fore_color.rgb = color
                        dot.line.fill.background()

                        is_top = True if not use_zigzag else (i % 2 == 0)

                        if is_top:
                            stem_top = line_y - dot_radius - stem_h
                        else:
                            stem_top = line_y + dot_radius

                        stem = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            dot_x - Emu(6000), stem_top,
                            Emu(12000), stem_h
                        )
                        stem.fill.solid()
                        stem.fill.fore_color.rgb = color
                        stem.line.fill.background()

                        if is_top:
                            # 텍스트 길이에 따라 높이 동적 계산 (1줄 ~150000, 추가 줄당 ~120000)
                            name_len = len(ms['name'])
                            est_lines = max(1, (name_len + 7) // 8)  # label_w 기준 약 8자/줄
                            if ms['show_month']:
                                est_lines += 1
                            top_label_h = Emu(150000 + est_lines * 120000)
                            label_top = stem_top - top_label_h
                        else:
                            label_top = stem_top + stem_h

                        label_left = dot_x - label_w // 2
                        if label_left < area_left:
                            label_left = area_left
                        if label_left + label_w > area_left + area_width:
                            label_left = area_left + area_width - label_w

                        label_center_x = label_left + label_w // 2
                        dot._element.spPr.xfrm.off.x = label_center_x - dot_radius
                        stem._element.spPr.xfrm.off.x = label_center_x - Emu(6000)

                        actual_label_h = top_label_h if is_top else label_h
                        txBox = slide.shapes.add_textbox(label_left, label_top, label_w, actual_label_h)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        if is_top:
                            tf.paragraphs[0].vertical_anchor = MSO_ANCHOR.BOTTOM

                        if is_top:
                            # 위쪽: 이름(바깥=위) → 월(바 가까이=아래)
                            p = tf.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            p.space_before = PtUtil(0)
                            p.space_after = PtUtil(0)
                            run = p.add_run()
                            run.text = ms['name']
                            run.font.size = PtUtil(8)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(55, 65, 81)

                            if ms['show_month']:
                                p2 = tf.add_paragraph()
                                p2.alignment = PP_ALIGN.CENTER
                                p2.space_before = PtUtil(0)
                                p2.space_after = PtUtil(0)
                                run2 = p2.add_run()
                                run2.text = ms['label']
                                run2.font.size = PtUtil(8)
                                run2.font.color.rgb = color
                        else:
                            # 아래쪽: 월(바 가까이=위) → 이름(바깥=아래)
                            p = tf.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            p.space_before = PtUtil(0)
                            p.space_after = PtUtil(0)
                            if ms['show_month']:
                                run_month = p.add_run()
                                run_month.text = ms['label']
                                run_month.font.size = PtUtil(8)
                                run_month.font.color.rgb = color

                                p2 = tf.add_paragraph()
                                p2.alignment = PP_ALIGN.CENTER
                                p2.space_before = PtUtil(0)
                                p2.space_after = PtUtil(0)
                                run_name = p2.add_run()
                            else:
                                run_name = p.add_run()

                            run_name.text = ms['name']
                            run_name.font.size = PtUtil(8)
                            run_name.font.bold = True
                            run_name.font.color.rgb = RGBColor(55, 65, 81)

                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)


        def _reset_autofit_font_scale(prs_obj):
            """모든 텍스트박스의 normAutofit fontScale을 제거하여
            PowerPoint가 파일을 열 때 자동 글자크기 축소를 재계산하도록 한다."""
            from pptx.oxml.ns import qn
            for slide in prs_obj.slides:
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
                    if bodyPr is None:
                        continue
                    normAutofit = bodyPr.find(qn('a:normAutofit'))
                    if normAutofit is not None:
                        normAutofit.attrib.pop('fontScale', None)
                        normAutofit.attrib.pop('lnSpcReduction', None)

        def auto_fit_text_boxes(prs_obj):
            """텍스트가 긴 텍스트 박스의 폰트를 자동 축소"""
            from pptx.util import Pt as PtUtil
            for slide in prs_obj.slides:
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    tf = shape.text_frame
                    full = ''.join(p.text for p in tf.paragraphs)
                    line_count = full.count('\n') + 1
                    if line_count > 5:
                        tf.word_wrap = True
                        box_height_inches = shape.height / 914400
                        available_lines = box_height_inches / 0.22
                        if line_count > available_lines:
                            target_pt = max(7, int(10 * available_lines / line_count))
                            for para in tf.paragraphs:
                                for run in para.runs:
                                    run.font.size = PtUtil(target_pt)

        def _add_performance_table(slide, project_data, left, top, width):
            """성과 섹션 아래에 성과목록 테이블을 추가 (batch용)"""
            from pptx.util import Pt as PtUtil, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.oxml.ns import qn

            perf_list = project_data.get('성과목록', [])
            if not perf_list:
                return

            border_color = RGBColor(209, 213, 219)
            header_bg = RGBColor(209, 213, 219)
            cell_bg = RGBColor(255, 255, 255)
            text_color = RGBColor(0, 0, 0)

            rows = len(perf_list) + 1
            cols = 5
            row_h = Emu(230000)
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
            tbl = table_shape.table

            tbl_pr = tbl._tbl.tblPr
            tbl_pr.attrib.pop('bandRow', None)
            tbl_pr.attrib.pop('bandCol', None)
            tbl_pr.attrib.pop('firstRow', None)
            tbl_pr.attrib.pop('lastRow', None)
            for child in list(tbl_pr):
                if child.tag.endswith('tableStyleId') or child.tag.endswith('tblStyle'):
                    tbl_pr.remove(child)

            def set_cell_style(cell, fill_rgb, b_color):
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for tag in ('lnL', 'lnR', 'lnT', 'lnB', 'solidFill'):
                    el = tcPr.find(qn(f'a:{tag}'))
                    if el is not None:
                        tcPr.remove(el)
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = tcPr.makeelement(qn(f'a:{edge}'), {})
                    ln.set('w', '6350')
                    sf = ln.makeelement(qn('a:solidFill'), {})
                    srgb = sf.makeelement(qn('a:srgbClr'), {'val': b_color})
                    sf.append(srgb)
                    ln.append(sf)
                    tcPr.append(ln)
                bg = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = bg.makeelement(qn('a:srgbClr'), {
                    'val': f'{fill_rgb.red:02X}{fill_rgb.green:02X}{fill_rgb.blue:02X}' if hasattr(fill_rgb, 'red') else str(fill_rgb)
                })
                bg.append(srgb)
                tcPr.append(bg)

            def fmt_delta(base_val, compare_val):
                try:
                    b = float(base_val)
                    c = float(compare_val)
                    d = c - b
                    return f'\n(+{d:g})' if d >= 0 else f'\n({d:g})'
                except (ValueError, TypeError):
                    return ''

            col_widths = [0.22, 0.34, 0.15, 0.15, 0.14]
            for ci, ratio in enumerate(col_widths):
                tbl.columns[ci].width = int(width * ratio)

            headers = ['성과분류', '성과항목명', '기존', '목표', '실적']
            for ci, h in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = PtUtil(8)
                    p.font.bold = True
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.CENTER
                set_cell_style(cell, header_bg, 'D1D5DB')

            for ri, perf in enumerate(perf_list):
                unit = perf.get('단위', '')
                unit_suffix = f' {unit}' if unit else ''
                base_raw = perf.get('현재수준', '-')
                target_raw = perf.get('목표수준', '-')
                actual_raw = perf.get('실적수준', '-') or '-'
                import re as _re
                perf_name = _re.sub(r'^\[.+?\]\s*', '', perf.get('성과항목', '-'))
                row_data = [
                    perf.get('소분류', '-'),
                    perf_name or '-',
                    f'{base_raw}{unit_suffix}',
                    f'{target_raw}{unit_suffix}{fmt_delta(base_raw, target_raw)}',
                    f'{actual_raw}{unit_suffix}{fmt_delta(base_raw, actual_raw)}',
                ]
                for ci, val in enumerate(row_data):
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        p.font.size = PtUtil(8)
                        p.font.color.rgb = text_color
                        p.alignment = PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.CENTER
                    set_cell_style(cell, cell_bg, 'D1D5DB')

        def replace_detail_placeholders_batch(prs_obj, project_data, placeholders):
            from pptx.util import Pt as PtUtil
            detail_keys = {
                f'{{{{상세정보_{k}}}}}': (k, v) for k, v in DETAIL_SECTION_MAP.items()
            }
            for slide in prs_obj.slides:
                perf_tables_to_add = []
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    tf = shape.text_frame
                    full_text = ''.join(run.text for p in tf.paragraphs for run in p.runs)
                    matched = [(pk, sk, sl) for pk, (sk, sl) in detail_keys.items() if pk in full_text]
                    if not matched:
                        continue

                    base_font_name = None
                    base_font_size = None
                    base_font_color = None
                    for p in tf.paragraphs:
                        for run in p.runs:
                            if run.font.name:
                                base_font_name = run.font.name
                            if run.font.size:
                                base_font_size = run.font.size
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    base_font_color = run.font.color.rgb
                            except AttributeError:
                                pass
                            break
                        if not base_font_size and p.font and p.font.size:
                            base_font_size = p.font.size
                        if not base_font_name and p.font and p.font.name:
                            base_font_name = p.font.name
                        break

                    body_size = base_font_size if base_font_size else PtUtil(12)
                    title_size = body_size + PtUtil(2)

                    # 줄간격 추출
                    base_space_before = None
                    base_space_after = None
                    base_line_spacing = None
                    for p in tf.paragraphs:
                        if p.space_before is not None:
                            base_space_before = p.space_before
                        if p.space_after is not None:
                            base_space_after = p.space_after
                        if p.line_spacing is not None:
                            base_line_spacing = p.line_spacing
                        break

                    def apply_spacing(para):
                        if base_space_before is not None:
                            para.space_before = base_space_before
                        if base_space_after is not None:
                            para.space_after = base_space_after
                        if base_line_spacing is not None:
                            para.line_spacing = base_line_spacing

                    for p in list(tf.paragraphs)[1:]:
                        p._element.getparent().remove(p._element)
                    tf.paragraphs[0].clear()
                    tf.paragraphs[0].font.size = None
                    first_para_used = False

                    for ph_key, section_key, section_label in matched:
                        raw = project_data.get(f'상세정보_{section_key}')
                        if not raw or not raw.get('enabled'):
                            placeholders[ph_key] = ''
                            continue
                        items = raw.get('items', [])

                        if not first_para_used:
                            p = tf.paragraphs[0]
                            first_para_used = True
                        else:
                            blank = tf.add_paragraph()
                            blank.font.size = None
                            apply_spacing(blank)
                            br = blank.add_run()
                            br.text = ''
                            br.font.size = PtUtil(6)
                            p = tf.add_paragraph()

                        p.font.size = None
                        apply_spacing(p)
                        run = p.add_run()
                        run.text = f'□ {section_label}'
                        run.font.size = title_size
                        run.font.bold = True
                        if base_font_name:
                            run.font.name = base_font_name
                        if base_font_color:
                            run.font.color.rgb = base_font_color

                        for item in items:
                            if isinstance(item, str):
                                item = {'text': item, 'children': []}
                            text = item.get('text', '')
                            if not text:
                                continue
                            p = tf.add_paragraph()
                            p.font.size = None
                            apply_spacing(p)
                            run = p.add_run()
                            run.text = f'- {text}'
                            run.font.size = body_size
                            run.font.bold = False
                            if base_font_name:
                                run.font.name = base_font_name
                            if base_font_color:
                                run.font.color.rgb = base_font_color

                            for child in item.get('children', []):
                                child_text = child.get('text', '') if isinstance(child, dict) else str(child)
                                if not child_text:
                                    continue
                                p = tf.add_paragraph()
                                p.font.size = None
                                apply_spacing(p)
                                run = p.add_run()
                                run.text = f'  · {child_text}'
                                run.font.size = body_size
                                run.font.bold = False
                                if base_font_name:
                                    run.font.name = base_font_name
                                if base_font_color:
                                    run.font.color.rgb = base_font_color

                        # 성과 섹션이면 테이블 추가 예약
                        if section_key == '성과':
                            from pptx.util import Emu
                            para_count = len(tf.paragraphs)
                            line_h = Emu(254000)  # 약 0.28인치/줄
                            estimated_h = para_count * line_h
                            actual_h = max(shape.height, estimated_h)
                            table_top = shape.top + actual_h + Emu(180000)
                            perf_tables_to_add.append((shape.left, table_top, shape.width))

                        placeholders[ph_key] = ''

                # 성과 테이블 실제 추가
                for tbl_left, tbl_top, tbl_width in perf_tables_to_add:
                    _add_performance_table(slide, project_data, tbl_left, tbl_top, tbl_width)


        def generate_single_pptx(project_data, override_template_path=None):
            """과제 1건에 대한 PPTX BytesIO 반환"""
            prs = Presentation(override_template_path or template_path)

            # 디버그: 상세정보 데이터 확인
            proj_name = project_data.get('과제명', '?')
            for dk in ('과제개요', '추진배경', '상세내용', '향후계획'):
                raw = project_data.get(f'상세정보_{dk}')
                if raw:
                    print(f"[DEBUG Batch] {proj_name} - 상세정보_{dk}: enabled={raw.get('enabled')}, items={len(raw.get('items',[]))}, images={len(raw.get('images',[]))}")
                else:
                    print(f"[DEBUG Batch] {proj_name} - 상세정보_{dk}: None")

            # 상세 과제 정보: 스타일 추출 → 동적 슬라이드 생성 → 원본 삭제
            detail_styles, detail_slide_idx = extract_detail_styles(prs)
            build_detail_slides(prs, project_data, detail_styles, detail_slide_idx)

            ph = build_placeholders(project_data)

            # 상세정보 플레이스홀더: 멀티스타일 처리 (제목 굵게 + 내용 보통)
            replace_detail_images(prs, project_data, DETAIL_SECTION_MAP, ph)
            replace_milestone_timeline(prs, project_data, ph)
            replace_detail_placeholders_batch(prs, project_data, ph)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            replace_paragraph_placeholders(paragraph, ph)
            auto_fit_text_boxes(prs)
            _reset_autofit_font_scale(prs)
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return buf

        def safe_name(name):
            return "".join(c for c in name if c not in r'\/:*?"<>|{}')[:50]

        if mode == 'single':
            # 하나의 Presentation 안에서 템플릿 슬라이드를 복제한 뒤 각각 치환
            # (cross-presentation 복사 대신 within-presentation 복제로 PPTX 손상 방지)
            # NOTE: single 모드에서는 상세 과제 정보 동적 슬라이드가 지원되지 않습니다.
            #       상세 과제 정보가 필요한 경우 individual 모드를 사용하세요.
            from copy import deepcopy
            from collections import Counter

            # single 모드: projectTemplateMap에서 가장 많이 사용되는 템플릿 선택
            single_template_path = template_path
            if project_template_map:
                tmpl_counts = Counter(project_template_map.values())
                most_common = tmpl_counts.most_common(1)[0][0]
                try:
                    single_template_path = get_template_path(most_common)
                except (ValueError, FileNotFoundError):
                    single_template_path = template_path

            prs = Presentation(single_template_path)

            # 상세 과제 정보 스타일 슬라이드 제거 (single 모드에서는 사용하지 않음)
            detail_styles_tmp, detail_idx_tmp = extract_detail_styles(prs)
            if detail_idx_tmp is not None:
                rId = prs.slides._sldIdLst[detail_idx_tmp].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId is None:
                    rId = prs.slides._sldIdLst[detail_idx_tmp].attrib.get('r:id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[detail_idx_tmp]

            template_slide_count = len(prs.slides)

            # 원본 템플릿 슬라이드 참조 (placeholder 치환 전)
            original_slides = list(prs.slides)[:template_slide_count]

            # 추가 과제 수만큼 템플릿 슬라이드 세트를 복제 (치환 전에 먼저 복제)
            for _ in range(len(projects_list) - 1):
                for orig_slide in original_slides:
                    new_slide = prs.slides.add_slide(orig_slide.slide_layout)
                    # add_slide가 자동 생성한 placeholder shape 제거
                    for ph_shape in list(new_slide.placeholders):
                        new_slide.shapes._spTree.remove(ph_shape._element)
                    # 원본 슬라이드의 모든 shape를 deep copy
                    for shape in orig_slide.shapes:
                        new_slide.shapes._spTree.append(deepcopy(shape._element))

            # 각 과제별로 해당 슬라이드 범위에서 placeholder 치환
            for proj_idx, project_data in enumerate(projects_list):
                ph = build_placeholders(project_data)

                # 상세정보 처리 (해당 과제의 슬라이드 범위에서만)
                # 임시 Presentation 대신, 슬라이드 범위를 직접 처리
                start = proj_idx * template_slide_count
                end = start + template_slide_count

                # 상세정보 이미지, 마일스톤, 텍스트 처리 (슬라이드 범위 내)
                for slide_idx in range(start, end):
                    slide = prs.slides[slide_idx]
                    # replace_detail_images 로직 (인라인) - 좌측/우측 이미지
                    # 이전 형식 플레이스홀더 제거 (호환성)
                    for section_key in DETAIL_SECTION_MAP:
                        old_ph = f'{{{{상세정보_{section_key}_이미지}}}}'
                        ph[old_ph] = ''
                    # 새 형식 플레이스홀더
                    img_side_phs = {'{{이미지_좌측}}': '좌측', '{{이미지_우측}}': '우측', '{{이미지_개요그림}}': '개요그림', '{{이미지_상세내용그림}}': '상세내용그림', '{{이미지_향후계획그림}}': '향후계획그림'}
                    for side_ph in img_side_phs:
                        ph[side_ph] = ''

                    for shape in list(slide.shapes):
                        if not hasattr(shape, 'text_frame'):
                            continue
                        full_text = ''.join(run.text for p in shape.text_frame.paragraphs for run in p.runs)

                        # 좌측/우측 이미지 플레이스홀더 처리
                        matched_sides = []
                        for side_ph, side in img_side_phs.items():
                            if side_ph in full_text:
                                matched_sides.append((side_ph, side))

                        if not matched_sides:
                            continue

                        import base64
                        from pptx.util import Emu
                        from pptx.enum.text import PP_ALIGN
                        area_left, area_top = shape.left, shape.top
                        area_width, area_height = shape.width, shape.height
                        gap = Emu(72000)

                        all_images = []
                        for side_ph, side in matched_sides:
                            images = project_data.get(f'이미지_{side}')
                            if images and isinstance(images, list):
                                for img in images:
                                    all_images.append(img)

                        if all_images:
                            img_count = min(len(all_images), 6)
                            if img_count == 1:
                                positions = [(area_left, area_top, area_width, area_height)]
                            elif img_count == 2:
                                w = (area_width - gap) // 2
                                positions = [(area_left, area_top, w, area_height), (area_left + w + gap, area_top, w, area_height)]
                            elif img_count == 3:
                                w = (area_width - gap * 2) // 3
                                positions = [
                                    (area_left, area_top, w, area_height),
                                    (area_left + w + gap, area_top, w, area_height),
                                    (area_left + (w + gap) * 2, area_top, w, area_height),
                                ]
                            elif img_count == 4:
                                w = (area_width - gap) // 2
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left, area_top + h + gap, w, h),
                                    (area_left + w + gap, area_top + h + gap, w, h),
                                ]
                            else:
                                w = (area_width - gap * 2) // 3
                                h = (area_height - gap) // 2
                                positions = [
                                    (area_left, area_top, w, h),
                                    (area_left + w + gap, area_top, w, h),
                                    (area_left + (w + gap) * 2, area_top, w, h),
                                ]
                                bottom_count = img_count - 3
                                if bottom_count == 1:
                                    positions.append((area_left, area_top + h + gap, area_width, h))
                                elif bottom_count == 2:
                                    bw = (area_width - gap) // 2
                                    positions.append((area_left, area_top + h + gap, bw, h))
                                    positions.append((area_left + bw + gap, area_top + h + gap, bw, h))
                                else:
                                    positions.append((area_left, area_top + h + gap, w, h))
                                    positions.append((area_left + w + gap, area_top + h + gap, w, h))
                                    positions.append((area_left + (w + gap) * 2, area_top + h + gap, w, h))
                            for i in range(img_count):
                                img_data = all_images[i]
                                data_url = img_data.get('dataUrl', '')
                                caption = img_data.get('caption', '')
                                px, py, pw, p_h = positions[i]
                                cap_h = Emu(200000) if caption else 0
                                img_h = p_h - cap_h
                                if data_url and ',' in data_url:
                                    try:
                                        img_bytes = base64.b64decode(data_url.split(',', 1)[1])
                                        img_stream = io.BytesIO(img_bytes)
                                        slide.shapes.add_picture(img_stream, px, py, pw, img_h)
                                    except Exception:
                                        pass
                                if caption:
                                    from pptx.util import Pt as PtU
                                    txB = slide.shapes.add_textbox(px, py + img_h, pw, cap_h)
                                    tf2 = txB.text_frame
                                    p2 = tf2.paragraphs[0]
                                    p2.text = caption
                                    p2.font.size = PtU(9)
                                    p2.font.italic = True
                                    p2.alignment = PP_ALIGN.CENTER
                            shape._element.getparent().remove(shape._element)
                            break  # shape 삭제했으므로 다음 shape로

                # 상세정보 텍스트 멀티스타일 처리
                class _TmpPrs:
                    """슬라이드 범위를 prs처럼 보이게 하는 래퍼"""
                    def __init__(self, slides_list):
                        self.slides = slides_list
                _tmp = _TmpPrs([prs.slides[i] for i in range(start, end)])
                replace_detail_placeholders_batch(_tmp, project_data, ph)
                replace_milestone_timeline(_tmp, project_data, ph)

                for slide_idx in range(start, end):
                    slide = prs.slides[slide_idx]
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                replace_paragraph_placeholders(paragraph, ph)

            auto_fit_text_boxes(prs)
            _reset_autofit_font_scale(prs)

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)

            year = projects_list[0].get('과제년도', '')
            filename = f"과제보고서_통합_{year}년_{len(projects_list)}건.pptx"

            report_log = DashboardActivityLog(
                action='EXPORT',
                target_type='REPORT',
                target_id=None,
                target_name=f'통합 보고서 ({len(projects_list)}건)',
                summary=f'{user.name}이(가) 통합 PPT 보고서를 생성함 ({len(projects_list)}건)',
                user_id=user_id,
                user_name=user.name,
                source='server'
            )
            db.session.add(report_log)
            db.session.commit()

            return send_file(
                buf,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        else:
            # individual: 각각 PPTX → ZIP
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for project_data in projects_list:
                    # 과제별 템플릿 매핑이 있으면 해당 템플릿 사용
                    proj_id = project_data.get('id', '')
                    per_project_template = None
                    if proj_id and proj_id in project_template_map:
                        try:
                            per_project_template = get_template_path(project_template_map[proj_id])
                        except (ValueError, FileNotFoundError):
                            per_project_template = None
                    pptx_buf = generate_single_pptx(project_data, per_project_template)
                    name = safe_name(project_data.get('과제명', '과제보고서'))
                    year = project_data.get('과제년도', '')
                    fname = f"{name}_{year}년_보고서.pptx"
                    zf.writestr(fname, pptx_buf.getvalue())

            zip_buffer.seek(0)

            year = projects_list[0].get('과제년도', '')
            zip_filename = f"과제보고서_개별_{year}년_{len(projects_list)}건.zip"

            report_log = DashboardActivityLog(
                action='EXPORT',
                target_type='REPORT',
                target_id=None,
                target_name=f'개별 보고서 ({len(projects_list)}건)',
                summary=f'{user.name}이(가) 개별 PPT 보고서를 생성함 ({len(projects_list)}건)',
                user_id=user_id,
                user_name=user.name,
                source='server'
            )
            db.session.add(report_log)
            db.session.commit()

            return send_file(
                zip_buffer,
                as_attachment=True,
                download_name=zip_filename,
                mimetype='application/zip'
            )

    except Exception as e:
        print(f"[Report Error] Batch PPT generation failed: {str(e)}")
        import traceback
        traceback.print_exc()
        return error_response(f'PPT 보고서 일괄 생성 실패: {str(e)}', status_code=500)


# ============================================
# PDF 보고서 생성 API (HTML→PDF, xhtml2pdf)
# 순수 Python, GTK 등 네이티브 의존성 없음
# ============================================

def _find_korean_font_pair():
    """시스템 한글 TTF 폰트 (regular, bold) 경로를 반환. bold가 없으면 regular 사용."""
    # (regular_path, bold_path) 후보 쌍
    candidates = [
        # Windows - 맑은 고딕
        (r'C:\Windows\Fonts\malgun.ttf', r'C:\Windows\Fonts\malgunbd.ttf'),
        # Windows - 나눔고딕
        (r'C:\Windows\Fonts\NanumGothic.ttf', r'C:\Windows\Fonts\NanumGothicBold.ttf'),
        # Windows - 나눔바른고딕
        (r'C:\Windows\Fonts\NanumBarunGothic.ttf', r'C:\Windows\Fonts\NanumBarunGothicBold.ttf'),
        # Windows - Noto Sans KR
        (r'C:\Windows\Fonts\NotoSansKR-Regular.ttf', r'C:\Windows\Fonts\NotoSansKR-Bold.ttf'),
        # Linux - Nanum
        ('/usr/share/fonts/truetype/nanum/NanumGothic.ttf',
         '/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf'),
        ('/usr/share/fonts/nanum/NanumGothic.ttf',
         '/usr/share/fonts/nanum/NanumGothicBold.ttf'),
        # macOS
        ('/Library/Fonts/AppleGothic.ttf', None),
    ]
    for regular, bold in candidates:
        if os.path.exists(regular):
            bold_path = bold if (bold and os.path.exists(bold)) else regular
            return regular, bold_path
    return None, None


_FONT_REGISTERED = False


def _register_korean_font():
    """ReportLab + xhtml2pdf 기본 폰트 매핑에 한글 폰트 등록. 1회만 실행.

    xhtml2pdf는 내부적으로 Helvetica/Times/Courier 같은 표준 폰트를 호출하므로,
    이 이름들을 한글 폰트로 덮어써서 어떤 CSS 지정이든 한글이 나오게 한다.
    """
    global _FONT_REGISTERED
    if _FONT_REGISTERED:
        return True
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
    except ImportError as e:
        print(f'[PDF Report] reportlab import 실패: {e}')
        return False

    regular, bold = _find_korean_font_pair()
    if not regular:
        print('[PDF Report] 시스템에서 한글 폰트를 찾을 수 없습니다. PDF의 한글이 깨질 수 있습니다.')
        return False

    try:
        # 우리 이름으로 등록
        pdfmetrics.registerFont(TTFont('KFont', regular))
        pdfmetrics.registerFont(TTFont('KFont-Bold', bold))
        registerFontFamily(
            'KFont',
            normal='KFont', bold='KFont-Bold',
            italic='KFont', boldItalic='KFont-Bold',
        )

        # xhtml2pdf가 부르는 표준 폰트 이름을 모두 한글 폰트로 매핑
        # (Helvetica, Times, Courier 등 - xhtml2pdf가 CSS font-family를 처리할 때
        #  결국 이 이름들로 떨어지므로 여기서 한글 폰트를 등록해야 한글이 보임)
        try:
            from xhtml2pdf.default import DEFAULT_FONT
            for key in list(DEFAULT_FONT.keys()):
                DEFAULT_FONT[key] = 'KFont'
            print(f'[PDF Report] xhtml2pdf DEFAULT_FONT 매핑 교체 완료: {len(DEFAULT_FONT)}개 키')
        except Exception as e:
            print(f'[PDF Report] DEFAULT_FONT 매핑 실패: {e}')

        # 안전장치: Helvetica/Times/Courier 이름으로도 한글 폰트 등록
        # (xhtml2pdf가 내부에서 직접 이 이름으로 폰트를 부를 때 대비)
        try:
            for std_name in ('Helvetica', 'Times-Roman', 'Courier'):
                pdfmetrics.registerFont(TTFont(std_name, regular))
            for std_bold in ('Helvetica-Bold', 'Times-Bold', 'Courier-Bold'):
                pdfmetrics.registerFont(TTFont(std_bold, bold))
            registerFontFamily(
                'Helvetica', normal='Helvetica', bold='Helvetica-Bold',
                italic='Helvetica', boldItalic='Helvetica-Bold',
            )
        except Exception as e:
            print(f'[PDF Report] 표준 폰트 오버라이드 일부 실패 (계속 진행): {e}')

        print(f'[PDF Report] 한글 폰트 등록 성공: regular={regular}, bold={bold}')
        _FONT_REGISTERED = True
        return True
    except Exception as e:
        print(f'[PDF Report] 한글 폰트 등록 실패: {e}')
        return False


def _build_pdf_html_template(font_registered=False):
    """폰트 등록 여부에 따라 HTML 템플릿 문자열을 반환.
    레이아웃: '결과 보고서' 페이지를 단열로 풀어놓은 형태.
    구성: 헤더 → 기본정보표 → 과제개요 → 추진배경 → 마일스톤(액션아이템 표) →
          과제목표 → 상세내용 → 기술/경영성과(+성과표) → 산출물 → 향후계획 → 관련이미지
    """
    body_font = "'KFont'" if font_registered else 'Helvetica'

    return r"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
@page {
    size: A4 portrait;
    margin: 1.6cm 1.8cm 2cm 1.8cm;
    @frame footer_frame {
        -pdf-frame-content: footer_content;
        bottom: 0.8cm;
        margin-left: 1.8cm;
        margin-right: 1.8cm;
        height: 0.7cm;
    }
}
body {
    font-family: """ + body_font + r""";
    font-size: 10pt;
    color: #1f2937;
    line-height: 1.5;
}
.project { page-break-after: always; }

/* 헤더 */
.project-header {
    border-bottom: 2pt solid #0066cc;
    padding-bottom: 6pt;
    margin-bottom: 10pt;
}
.project-title {
    font-size: 17pt;
    font-weight: bold;
    color: #0066cc;
    margin: 0 0 4pt 0;
}
.project-meta {
    color: #4b5563;
    font-size: 9pt;
}
.project-meta span { padding-right: 10pt; }
.project-meta strong { color: #1f2937; }
.status-badge {
    background-color: #dbeafe;
    color: #1e40af;
    font-size: 8pt;
    font-weight: bold;
    padding: 1pt 6pt;
}

/* 메타 정보 표 */
.meta-table {
    width: 100%;
    border-collapse: collapse;
    margin-bottom: 10pt;
    font-size: 9pt;
}
.meta-table th, .meta-table td {
    border: 0.5pt solid #d1d5db;
    padding: 4pt 6pt;
    text-align: left;
    vertical-align: top;
}
.meta-table th {
    background-color: #f3f4f6;
    font-weight: bold;
    color: #374151;
    width: 13%;
}
.progress-bar-wrap {
    background-color: #e5e7eb;
    height: 8pt;
    width: 100%;
}
.progress-bar-fill {
    background-color: #0066cc;
    height: 8pt;
}
.tag-poc {
    background-color: #fef3c7;
    color: #92400e;
    font-size: 8pt;
    font-weight: bold;
    padding: 1pt 5pt;
    margin-right: 3pt;
}
.tag-focus {
    background-color: #fce7f3;
    color: #be185d;
    font-size: 8pt;
    font-weight: bold;
    padding: 1pt 5pt;
}

/* 일반 섹션 */
.section { margin-top: 32pt; }
.section-title {
    font-size: 12pt;
    font-weight: bold;
    color: #0066cc;
    border-left: 3pt solid #0066cc;
    padding-left: 6pt;
    padding-top: 4pt;
    padding-bottom: 4pt;
    margin: 0 0 8pt 0;
}
.section ul {
    margin: 0 0 0 14pt;
    padding: 0;
}
.section ul ul {
    margin: 3pt 0 3pt 14pt;
}
.section li { margin-bottom: 3pt; }
.section ul ul li {
    color: #4b5563;
    font-size: 9pt;
}
.empty {
    color: #9ca3af;
    font-style: italic;
    margin: 0;
}
.detail-text {
    margin-top: 3pt;
    color: #374151;
    font-size: 9pt;
}

/* 이미지 */
.section-images { margin-top: 8pt; }
.image-item {
    margin-bottom: 8pt;
    text-align: center;
}
.image-item img {
    max-width: 14cm;
    max-height: 7cm;
}
.image-caption {
    font-size: 8pt;
    color: #6b7280;
    font-style: italic;
    margin-top: 2pt;
}

/* 마일스톤 표 */
.milestone-summary {
    background-color: #f3f4f6;
    color: #374151;
    font-size: 9pt;
    padding: 2pt 7pt;
    margin-left: 4pt;
}
.milestone-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 4pt;
    margin-bottom: 14pt;
    font-size: 9pt;
}
.milestone-table th {
    background-color: #f8fafc;
    border: 0.5pt solid #d1d5db;
    padding: 4pt 6pt;
    text-align: center;
    font-weight: bold;
    color: #475569;
    font-size: 8.5pt;
}
.milestone-table td {
    border: 0.5pt solid #e5e7eb;
    padding: 4pt 6pt;
    vertical-align: middle;
}
.milestone-table td.no { text-align: center; width: 6%; color: #94a3b8; }
.milestone-table td.name { width: 56%; }
.milestone-table td.date { text-align: center; width: 14%; color: #475569; }
.milestone-table td.status { text-align: center; width: 10%; }
.milestone-done {
    background-color: #dcfce7;
    color: #166534;
    font-size: 8pt;
    font-weight: bold;
    padding: 1pt 5pt;
}
.milestone-pending {
    background-color: #fef3c7;
    color: #92400e;
    font-size: 8pt;
    padding: 1pt 5pt;
}
.milestone-name-done {
    color: #166534;
}

/* 성과 표 */
.perf-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 6pt;
    margin-bottom: 14pt;
    font-size: 8.5pt;
}
.perf-table th {
    background-color: #eff6ff;
    border: 0.5pt solid #93c5fd;
    padding: 4pt 5pt;
    text-align: center;
    font-weight: bold;
    color: #1e40af;
}
.perf-table th.col-cat { width: 16%; }
.perf-table th.col-name { width: 44%; }
.perf-table th.col-val { width: 13.3%; }
.perf-table td {
    border: 0.5pt solid #e5e7eb;
    padding: 4pt 5pt;
    vertical-align: top;
}
.perf-table td.center { text-align: center; }
.perf-cat { color: #4b5563; }
.perf-delta {
    color: #6b7280;
    font-size: 7.5pt;
    margin-left: 3pt;
}
.perf-delta-pos { color: #059669; }
.perf-delta-neg { color: #dc2626; }

#footer_content {
    text-align: center;
    font-size: 8pt;
    color: #6b7280;
}
</style>
</head>
<body>
<div id="footer_content">
    <pdf:pagenumber/> / <pdf:pagecount/>
</div>

{% for project in projects %}
<div class="project">
    <!-- 헤더 -->
    <div class="project-header">
        <div class="project-title">{{ project.과제명 or '-' }}</div>
        <div class="project-meta">
            <span><strong>{{ project.사업부 or '-' }}</strong></span>
            <span>{{ project.프로세스 or '-' }}</span>
            <span>PL: {{ project.과제PL or '-' }}</span>
            <span class="status-badge">{{ project.진행상태 or '미착수' }}</span>
            <span>{{ project.시작 or '-' }}월 ~ {{ project.종료 or '-' }}월</span>
        </div>
    </div>

    <!-- 기본 정보 표 -->
    <table class="meta-table">
        <tr>
            <th>과제년도</th><td>{{ project.과제년도 or '-' }}</td>
            <th>사업부</th><td>{{ project.사업부 or '-' }}</td>
            <th>프로세스</th><td>{{ project.프로세스 or '-' }}</td>
        </tr>
        <tr>
            <th>과제영역</th><td>{{ project.과제영역 or '-' }}</td>
            <th>과제구분</th><td>{{ project.과제구분 or '-' }}</td>
            <th>진행상태</th>
            <td><span class="status-badge">{{ project.진행상태 or '미착수' }}</span></td>
        </tr>
        <tr>
            <th>진행률</th>
            <td>
                <span style="font-weight:bold;">{{ project.진행률 or 0 }}%</span>
                <div class="progress-bar-wrap" style="height: 4pt; margin-top: 2pt;">
                    <div class="progress-bar-fill" style="width: {{ project.진행률 or 0 }}%; height: 4pt;"></div>
                </div>
            </td>
            <th>기간</th><td>{{ project.시작 or '-' }}월 ~ {{ project.종료 or '-' }}월</td>
            <th>구분</th>
            <td>
                {% if project.PoC과제여부 %}<span class="tag-poc">PoC</span>{% endif %}
                {% if project.중점과제여부 %}<span class="tag-focus">중점과제</span>{% endif %}
                {% if not project.PoC과제여부 and not project.중점과제여부 %}-{% endif %}
            </td>
        </tr>
        <tr>
            <th>과제PL</th><td>{{ project.과제PL or '-' }}</td>
            <th>작성자</th><td colspan="3">{{ project.작성자 or '-' }}</td>
        </tr>
        <tr>
            <th>담당부서</th><td colspan="5">{{ dept_str(project) }}</td>
        </tr>
        <tr>
            <th>참여인력</th><td colspan="5">{{ personnel_str(project) }}</td>
        </tr>
    </table>

    <!-- 상세 섹션 (순서: 과제개요 → 추진배경 → 마일스톤 → 과제목표 → 상세내용 → 성과 → 산출물 → 향후계획) -->
    {% for section in sections %}
        {% if section.key == '_milestone' %}
            {# 마일스톤 #}
            {% set milestones = project.액션아이템목록 or [] %}
            {% if milestones %}
            <div class="section">
                <div class="section-title">
                    🏁 마일스톤
                    <span class="milestone-summary">
                        {{ project.milestone_summary.completed }}/{{ project.milestone_summary.total }}
                        ({{ project.milestone_summary.rate }}%)
                    </span>
                </div>
                <table class="milestone-table">
                    <thead>
                        <tr>
                            <th>#</th>
                            <th>액션아이템</th>
                            <th>목표일</th>
                            <th>완료일</th>
                            <th>상태</th>
                        </tr>
                    </thead>
                    <tbody>
                        {% for ms in milestones %}
                        <tr>
                            <td class="no">{{ loop.index }}</td>
                            <td class="name">
                                {% if ms.완료여부 %}
                                    <span class="milestone-name-done">{{ ms.제목 or ('액션아이템 ' + loop.index|string) }}</span>
                                {% else %}
                                    {{ ms.제목 or ('액션아이템 ' + loop.index|string) }}
                                {% endif %}
                            </td>
                            <td class="date">{{ ms.목표일_fmt or '-' }}</td>
                            <td class="date">{{ ms.완료일_fmt or '-' }}</td>
                            <td class="status">
                                {% if ms.완료여부 %}
                                    <span class="milestone-done">완료</span>
                                {% else %}
                                    <span class="milestone-pending">진행중</span>
                                {% endif %}
                            </td>
                        </tr>
                        {% endfor %}
                    </tbody>
                </table>
            </div>
            {% endif %}
        {% else %}
            {# 일반 상세 정보 섹션 #}
            {% set data = project['상세정보_' + section.key] %}
            {% if data.enabled %}
                {% set content_items = data['items'] | selectattr('text') | list %}
                {% set section_images = project.get('이미지_' + section.image_category, []) if section.image_category else [] %}
                {% set image_items = section_images | selectattr('dataUrl') | list %}
                {% set perf_list = project.성과목록 if section.key == '성과' else [] %}
                <div class="section">
                    <div class="section-title">{{ section.icon }} {{ section.title }}</div>
                    {% if not content_items and not perf_list and not image_items %}
                        <p class="empty">내용 없음</p>
                    {% endif %}
                    {% if content_items %}
                    <ul>
                        {% for item in content_items %}
                            {% set child_items = item.children | selectattr('text') | list %}
                            <li>
                                {{ item.text }}
                                {% if child_items and not section.parent_only %}
                                <ul>
                                    {% for child in child_items %}
                                    <li>{{ child.text }}</li>
                                    {% endfor %}
                                </ul>
                                {% endif %}
                            </li>
                        {% endfor %}
                    </ul>
                    {% endif %}
                    {# 성과 표 #}
                    {% if section.key == '성과' and perf_list %}
                    <table class="perf-table">
                        <thead>
                            <tr>
                                <th class="col-cat">성과분류</th>
                                <th class="col-name">성과항목명</th>
                                <th class="col-val">기존</th>
                                <th class="col-val">목표</th>
                                <th class="col-val">실적</th>
                            </tr>
                        </thead>
                        <tbody>
                            {% for perf in perf_list %}
                            <tr>
                                <td class="perf-cat">{{ perf.분류_표시 or '-' }}</td>
                                <td>{{ perf.성과항목_표시 or '-' }}</td>
                                <td class="center">{{ perf.기존_표시 or '-' }}</td>
                                <td class="center">
                                    {{ perf.목표_표시 or '-' }}
                                    {% if perf.목표_델타 %}<span class="perf-delta">{{ perf.목표_델타 }}</span>{% endif %}
                                </td>
                                <td class="center">
                                    {{ perf.실적_표시 or '-' }}
                                    {% if perf.실적_델타 %}
                                        <span class="perf-delta {{ 'perf-delta-pos' if perf.실적_달성 else 'perf-delta-neg' }}">{{ perf.실적_델타 }}</span>
                                    {% endif %}
                                </td>
                            </tr>
                            {% endfor %}
                        </tbody>
                    </table>
                    {% endif %}
                    {# 섹션 이미지 #}
                    {% if image_items %}
                    <div class="section-images">
                        {% for img in image_items %}
                        <div class="image-item">
                            <img src="{{ img.dataUrl|safe }}" />
                            {% if img.caption %}<div class="image-caption">{{ img.caption }}</div>{% endif %}
                        </div>
                        {% endfor %}
                    </div>
                    {% endif %}
                </div>
            {% endif %}
        {% endif %}
    {% endfor %}

    {% if project.과제상세설명_clean %}
    <div class="section">
        <div class="section-title">📄 과제 상세 설명</div>
        <div class="detail-text">{{ project.과제상세설명_clean }}</div>
    </div>
    {% endif %}

    {# 관련 이미지 (좌측/우측) #}
    {% set left_imgs = (project.get('이미지_좌측', []) or []) | selectattr('dataUrl') | list %}
    {% set right_imgs = (project.get('이미지_우측', []) or []) | selectattr('dataUrl') | list %}
    {% if left_imgs or right_imgs %}
    <div class="section">
        <div class="section-title">🖼 관련 이미지</div>
        <div class="section-images">
            {% for img in left_imgs %}
            <div class="image-item">
                <img src="{{ img.dataUrl|safe }}" />
                {% if img.caption %}<div class="image-caption">{{ img.caption }}</div>{% endif %}
            </div>
            {% endfor %}
            {% for img in right_imgs %}
            <div class="image-item">
                <img src="{{ img.dataUrl|safe }}" />
                {% if img.caption %}<div class="image-caption">{{ img.caption }}</div>{% endif %}
            </div>
            {% endfor %}
        </div>
    </div>
    {% endif %}
</div>
{% endfor %}
</body>
</html>"""


def _pdf_link_callback(uri, rel):
    """xhtml2pdf 리소스 로드 콜백.
    - data:image/...;base64,... URI는 그대로 통과 (xhtml2pdf가 내부에서 디코드)
    - file:// URI는 로컬 경로로 변환
    - 그 외는 그대로 반환
    """
    if uri.startswith('data:'):
        return uri
    if uri.startswith('file://'):
        return uri[7:] if os.name != 'nt' else uri[8:]
    return uri


def _render_pdf_from_html(html_str):
    """xhtml2pdf로 HTML 문자열을 PDF 바이트로 변환."""
    from xhtml2pdf import pisa
    result = io.BytesIO()
    pisa_status = pisa.CreatePDF(
        src=html_str,
        dest=result,
        encoding='utf-8',
        link_callback=_pdf_link_callback,
    )
    if pisa_status.err:
        raise RuntimeError(f'xhtml2pdf 변환 실패: {pisa_status.err}건의 에러 발생')
    return result.getvalue()


@bp.route('/report/pdf/batch', methods=['POST'])
@jwt_required()
def generate_batch_pdf_report():
    """여러 과제 PDF 보고서 일괄 생성 (HTML→PDF, xhtml2pdf)
    mode: 'single' — 모든 과제를 하나의 PDF로 합침
    mode: 'individual' — 과제별 개별 PDF를 ZIP으로 묶어 반환
    """
    try:
        try:
            from xhtml2pdf import pisa  # noqa: F401
        except ImportError as imp_err:
            return error_response(
                f'xhtml2pdf가 설치되지 않았습니다. '
                f'백엔드 환경에서 다음 명령을 실행해주세요: '
                f'pip install -r requirements.txt  /  상세: {str(imp_err)}',
                status_code=500
            )
        from jinja2 import Environment, BaseLoader, select_autoescape
        import re
        import zipfile

        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        if not user or (user.role not in ('admin', 'manager', 'dt_office') and not user.is_admin):
            return error_response('보고서 저장 권한이 없습니다.', status_code=403)

        body = request.get_json()
        if not body:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        projects_list = body.get('projects', [])
        mode = body.get('mode', 'single')

        if not projects_list:
            return error_response('선택된 과제가 없습니다.', status_code=400)

        # imageId 만 온 이미지에 dataUrl 을 채운다 — PDF 템플릿은 dataUrl 만 읽는다
        hydrate_report_images(projects_list)

        def strip_html(html_text):
            if not html_text:
                return ''
            text = re.sub(r'<[^>]+>', '', html_text)
            text = text.replace('&nbsp;', ' ').replace('&lt;', '<')
            text = text.replace('&gt;', '>').replace('&amp;', '&').replace('&quot;', '"')
            text = re.sub(r'\s+', ' ', text).strip()
            return text

        # 섹션 정의 - 결과 보고서 페이지와 동일 순서 + 마일스톤 삽입
        # parent_only: True면 자식 항목(children) 숨김 (ProjectReportView와 동일)
        # image_category: 해당 섹션 끝에 붙일 이미지 필드 카테고리
        # key='_milestone' 은 마일스톤 액션아이템 표 (특수 처리)
        sections = [
            {'key': '과제개요', 'title': '과제 개요', 'icon': '📌',
             'parent_only': True, 'image_category': '개요그림'},
            {'key': '추진배경', 'title': '추진 배경', 'icon': '🔍',
             'parent_only': True, 'image_category': None},
            {'key': '_milestone', 'title': '마일스톤', 'icon': '🏁',
             'parent_only': False, 'image_category': None},
            {'key': '과제목표', 'title': '과제 목표', 'icon': '🎯',
             'parent_only': True, 'image_category': None},
            {'key': '상세내용', 'title': '상세 내용', 'icon': '📝',
             'parent_only': False, 'image_category': '상세내용그림'},
            {'key': '성과', 'title': '기술/경영 성과', 'icon': '📊',
             'parent_only': False, 'image_category': None},
            {'key': '산출물', 'title': '산출물', 'icon': '📦',
             'parent_only': False, 'image_category': None},
            {'key': '향후계획', 'title': '향후 계획', 'icon': '🗓',
             'parent_only': False, 'image_category': '향후계획그림'},
        ]

        def normalize_section(section_data):
            if not isinstance(section_data, dict):
                return {'enabled': False, 'items': []}
            items_raw = section_data.get('items') or []
            normalized = []
            for item in items_raw:
                if isinstance(item, str):
                    normalized.append({'text': item, 'children': []})
                elif isinstance(item, dict):
                    children = item.get('children') or []
                    norm_children = []
                    for c in children:
                        if isinstance(c, str):
                            norm_children.append({'text': c})
                        elif isinstance(c, dict):
                            norm_children.append({'text': c.get('text', '') or ''})
                    normalized.append({
                        'text': item.get('text', '') or '',
                        'children': norm_children
                    })
            return {
                'enabled': bool(section_data.get('enabled', False)),
                'items': normalized
            }

        def fmt_date(date_str):
            """'2026-06-18' / '06-18' → '6/18' 형태로 변환."""
            if not date_str:
                return ''
            parts = str(date_str).split('-')
            try:
                if len(parts) >= 3:
                    return f'{int(parts[1])}/{int(parts[2])}'
                if len(parts) == 2:
                    return f'{int(parts[0])}/{int(parts[1])}'
            except (ValueError, IndexError):
                pass
            return str(date_str)

        def normalize_milestones(milestones):
            """액션아이템 리스트에 표시용 날짜 포맷 추가."""
            if not isinstance(milestones, list):
                return []
            result = []
            for m in milestones:
                if not isinstance(m, dict):
                    continue
                nm = dict(m)
                nm['목표일_fmt'] = fmt_date(m.get('목표일', ''))
                nm['완료일_fmt'] = fmt_date(m.get('완료일', '')) if m.get('완료여부') else ''
                result.append(nm)
            return result

        def compute_milestone_summary(milestones):
            """완료/전체/% 요약."""
            if not isinstance(milestones, list) or not milestones:
                return {'completed': 0, 'total': 0, 'rate': 0}
            total = 0
            completed = 0
            for m in milestones:
                if not isinstance(m, dict):
                    continue
                details = m.get('세부항목목록') or []
                if isinstance(details, list) and details:
                    total += len(details)
                    completed += sum(1 for d in details if isinstance(d, dict) and d.get('완료여부'))
                else:
                    total += 1
                    if m.get('완료여부'):
                        completed += 1
            rate = round((completed / total) * 100) if total > 0 else 0
            return {'completed': completed, 'total': total, 'rate': rate}

        def try_parse_float(v):
            if v is None or v == '':
                return None
            try:
                return float(v)
            except (ValueError, TypeError):
                return None

        def fmt_num(v):
            """숫자 표시 (정수면 정수로, 소수면 소수점 2자리 이내)."""
            if v is None:
                return ''
            if abs(v - round(v)) < 1e-9:
                return str(int(round(v)))
            return f'{v:.2f}'.rstrip('0').rstrip('.')

        def normalize_performances(perfs):
            """성과 리스트에 표시용 필드 추가 (분류, 항목, 기존/목표/실적, 델타)."""
            if not isinstance(perfs, list):
                return []
            result = []
            for perf in perfs:
                if not isinstance(perf, dict):
                    continue
                np = dict(perf)
                unit = perf.get('단위') or ''
                base = try_parse_float(perf.get('현재수준'))
                target = try_parse_float(perf.get('목표수준'))
                actual = try_parse_float(perf.get('실적수준'))
                base_raw = perf.get('현재수준')
                target_raw = perf.get('목표수준')
                actual_raw = perf.get('실적수준')

                # 분류: 소분류 > 대분류
                np['분류_표시'] = perf.get('소분류') or perf.get('대분류') or '-'
                # 성과항목: [...] 접두사 제거
                item_name = perf.get('성과항목') or '-'
                import re as re_mod
                np['성과항목_표시'] = re_mod.sub(r'^\[.+?\]\s*', '', str(item_name))
                # 기존/목표/실적 표시 (값 + 단위)
                np['기존_표시'] = f'{base_raw} {unit}'.strip() if base_raw not in (None, '') else ''
                np['목표_표시'] = f'{target_raw} {unit}'.strip() if target_raw not in (None, '') else ''
                np['실적_표시'] = f'{actual_raw} {unit}'.strip() if actual_raw not in (None, '') else ''
                # 델타 (vs 기존)
                if base is not None and target is not None:
                    d = target - base
                    np['목표_델타'] = f'+{fmt_num(d)}' if d >= 0 else fmt_num(d)
                else:
                    np['목표_델타'] = ''
                if base is not None and actual is not None:
                    d = actual - base
                    np['실적_델타'] = f'+{fmt_num(d)}' if d >= 0 else fmt_num(d)
                    np['실적_달성'] = (target is None) or (actual >= target)
                else:
                    np['실적_델타'] = ''
                    np['실적_달성'] = False
                result.append(np)
            return result

        def normalize_project(p):
            np = dict(p)
            for sec in sections:
                if sec['key'].startswith('_'):
                    continue
                key = f'상세정보_{sec["key"]}'
                np[key] = normalize_section(p.get(key))
            np['과제상세설명_clean'] = strip_html(p.get('과제상세설명', ''))
            milestones = normalize_milestones(p.get('액션아이템목록', []))
            np['액션아이템목록'] = milestones
            np['milestone_summary'] = compute_milestone_summary(p.get('액션아이템목록', []))
            np['성과목록'] = normalize_performances(p.get('성과목록', []))
            return np

        def personnel_str(project):
            personnel = project.get('과제참여인력목록', []) or []
            if not personnel:
                return '-'
            return ', '.join(
                f'{p.get("이름", "-")}({p.get("부서", "-")})'
                for p in personnel
            )

        def dept_str(project):
            depts = project.get('담당부서목록', []) or []
            return ', '.join(depts) if depts else '-'

        font_registered = _register_korean_font()
        template_str = _build_pdf_html_template(font_registered)
        env = Environment(
            loader=BaseLoader(),
            autoescape=select_autoescape(['html', 'xml'])
        )
        template = env.from_string(template_str)

        if mode == 'single':
            normalized = [normalize_project(p) for p in projects_list]
            html_str = template.render(
                projects=normalized,
                sections=sections,
                personnel_str=personnel_str,
                dept_str=dept_str
            )
            pdf_bytes = _render_pdf_from_html(html_str)
            pdf_buffer = io.BytesIO(pdf_bytes)
            pdf_buffer.seek(0)

            filename = f'과제보고서_통합_{len(projects_list)}건.pdf'

            report_log = DashboardActivityLog(
                action='EXPORT', target_type='REPORT',
                target_name=f'{len(projects_list)}건 통합',
                summary=f'{user.name}이(가) PDF 보고서를 생성함 (통합 {len(projects_list)}건)',
                user_id=user_id, user_name=user.name, source='server'
            )
            db.session.add(report_log)
            db.session.commit()

            return send_file(
                pdf_buffer,
                as_attachment=True,
                download_name=filename,
                mimetype='application/pdf'
            )
        else:
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for p in projects_list:
                    normalized = normalize_project(p)
                    html_str = template.render(
                        projects=[normalized],
                        sections=sections,
                        personnel_str=personnel_str,
                        dept_str=dept_str
                    )
                    pdf_bytes = _render_pdf_from_html(html_str)
                    safe_name = "".join(
                        c for c in (p.get('과제명') or '과제보고서')
                        if c not in r'\/:*?"<>|'
                    )
                    pdf_filename = f"{safe_name[:60]}_{p.get('과제년도', '')}년_보고서.pdf"
                    zf.writestr(pdf_filename, pdf_bytes)
            zip_buffer.seek(0)
            zip_filename = f'과제보고서_개별_{len(projects_list)}건.zip'

            report_log = DashboardActivityLog(
                action='EXPORT', target_type='REPORT',
                target_name=f'{len(projects_list)}건 개별',
                summary=f'{user.name}이(가) PDF 보고서를 생성함 (개별 {len(projects_list)}건)',
                user_id=user_id, user_name=user.name, source='server'
            )
            db.session.add(report_log)
            db.session.commit()

            return send_file(
                zip_buffer,
                as_attachment=True,
                download_name=zip_filename,
                mimetype='application/zip'
            )

    except Exception as e:
        print(f"[Report Error] PDF generation failed: {str(e)}")
        import traceback
        traceback.print_exc()
        return error_response(f'PDF 보고서 생성 실패: {str(e)}', status_code=500)
