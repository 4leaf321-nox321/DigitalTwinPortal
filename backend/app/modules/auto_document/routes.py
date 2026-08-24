"""
문서 자동 작성 모듈 - PPT/Word 템플릿 기반 문서 생성
"""
import os
import re
import io
import json
import uuid
import glob as glob_mod
from datetime import datetime

from flask import request, send_file
from flask_jwt_extended import jwt_required, get_jwt_identity
from pptx import Presentation
from pptx.util import Pt
from docx import Document as DocxDocument

from app.modules.auto_document import bp
from app.modules.auth.models import User
from app.shared.responses import success_response, error_response

# 템플릿 디렉토리
_BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
TEMPLATE_DIR = os.path.join(_BASE_DIR, 'templates', 'document_type_ppt')
TEMPLATE_DIR_WORD = os.path.join(_BASE_DIR, 'templates', 'document_type_word')

# 프리셋 저장 파일
PRESETS_FILE = os.path.join(_BASE_DIR, 'data', 'auto_document_presets.json')


def _load_presets():
    """프리셋 목록 로드"""
    if not os.path.exists(PRESETS_FILE):
        return []
    try:
        with open(PRESETS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError):
        return []


def _save_presets(presets):
    """프리셋 목록 저장"""
    os.makedirs(os.path.dirname(PRESETS_FILE), exist_ok=True)
    with open(PRESETS_FILE, 'w', encoding='utf-8') as f:
        json.dump(presets, f, ensure_ascii=False, indent=2)


def _check_admin(user):
    """관리자 권한 확인"""
    return user and (user.role == 'admin' or user.is_admin)


def _scan_placeholders(template_path):
    """PPT 템플릿 파일에서 {{...}} 플레이스홀더를 스캔하여 반환"""
    prs = Presentation(template_path)
    found = set()
    pattern = re.compile(r'\{\{([^}]+)\}\}')

    for slide in prs.slides:
        for shape in slide.shapes:
            # 일반 텍스트 프레임
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join(run.text for run in paragraph.runs)
                    for match in pattern.finditer(full_text):
                        found.add(match.group(1))
            # 테이블
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = ''.join(run.text for run in paragraph.runs)
                            for match in pattern.finditer(full_text):
                                found.add(match.group(1))

    return sorted(found)


def _scan_placeholders_docx(template_path):
    """Word 템플릿 파일에서 {{...}} 플레이스홀더를 스캔하여 반환"""
    doc = DocxDocument(template_path)
    found = set()
    pattern = re.compile(r'\{\{([^}]+)\}\}')

    # 본문 단락
    for paragraph in doc.paragraphs:
        full_text = ''.join(run.text for run in paragraph.runs)
        for match in pattern.finditer(full_text):
            found.add(match.group(1))

    # 테이블
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    full_text = ''.join(run.text for run in paragraph.runs)
                    for match in pattern.finditer(full_text):
                        found.add(match.group(1))

    # 헤더/푸터
    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer and header_footer.is_linked_to_previous is False:
                for paragraph in header_footer.paragraphs:
                    full_text = ''.join(run.text for run in paragraph.runs)
                    for match in pattern.finditer(full_text):
                        found.add(match.group(1))

    return sorted(found)


def _replace_paragraph_placeholders(paragraph, placeholders):
    """PPT paragraph 전체 텍스트에서 placeholder를 찾아 치환"""
    runs = paragraph.runs
    if not runs:
        return
    full_text = ''.join(run.text for run in runs)

    has_match = False
    for ph in placeholders:
        if ph in full_text:
            has_match = True
            break
    if not has_match:
        return

    for ph, val in placeholders.items():
        full_text = full_text.replace(ph, str(val))

    runs[0].text = full_text
    for run in runs[1:]:
        run.text = ''


def _replace_paragraph_placeholders_docx(paragraph, placeholders):
    """Word paragraph 전체 텍스트에서 placeholder를 찾아 치환 (서식 유지)"""
    runs = paragraph.runs
    if not runs:
        return
    full_text = ''.join(run.text for run in runs)

    has_match = False
    for ph in placeholders:
        if ph in full_text:
            has_match = True
            break
    if not has_match:
        return

    for ph, val in placeholders.items():
        full_text = full_text.replace(ph, str(val))

    runs[0].text = full_text
    for run in runs[1:]:
        run.text = ''


def _replace_all_placeholders_docx(doc, placeholders):
    """Word 문서 전체에서 플레이스홀더를 치환"""
    # 본문 단락
    for paragraph in doc.paragraphs:
        _replace_paragraph_placeholders_docx(paragraph, placeholders)

    # 테이블
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _replace_paragraph_placeholders_docx(paragraph, placeholders)

    # 헤더/푸터
    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer and header_footer.is_linked_to_previous is False:
                for paragraph in header_footer.paragraphs:
                    _replace_paragraph_placeholders_docx(paragraph, placeholders)


def _auto_fit_text(slide, shape):
    """텍스트가 긴 텍스트 박스의 폰트를 자동 축소"""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    full = ''.join(p.text for p in tf.paragraphs)
    line_count = full.count('\n') + 1
    if line_count > 5:
        tf.word_wrap = True
        box_height_inches = shape.height / 914400
        available_lines = box_height_inches / 0.22
        if line_count > available_lines:
            target_pt = max(7, int(10 * available_lines / line_count))
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(target_pt)


# ============================================================
# API 엔드포인트
# ============================================================

@bp.route('/templates', methods=['GET'])
@jwt_required()
def get_templates():
    """PPT/Word 템플릿 목록 조회"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        templates = []

        # PPT 템플릿
        pptx_files = glob_mod.glob(os.path.join(TEMPLATE_DIR, '*.pptx'))
        for f in sorted(pptx_files):
            basename = os.path.basename(f)
            if basename.startswith('~$'):
                continue
            name = os.path.splitext(basename)[0]
            templates.append({'filename': basename, 'name': name, 'type': 'pptx'})

        # Word 템플릿
        docx_files = glob_mod.glob(os.path.join(TEMPLATE_DIR_WORD, '*.docx'))
        for f in sorted(docx_files):
            basename = os.path.basename(f)
            if basename.startswith('~$'):
                continue
            name = os.path.splitext(basename)[0]
            templates.append({'filename': basename, 'name': name, 'type': 'docx'})

        return success_response(templates)
    except Exception as e:
        return error_response(f'템플릿 목록 조회 실패: {str(e)}', status_code=500)


@bp.route('/templates/placeholders', methods=['POST'])
@jwt_required()
def get_template_placeholders():
    """선택한 템플릿의 플레이스홀더 목록 스캔 (PPT/Word)"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        data = request.get_json()
        filename = data.get('template') if data else None
        if not filename:
            return error_response('템플릿 파일명이 필요합니다.', status_code=400)

        safe_name = os.path.basename(filename)

        if safe_name.endswith('.pptx'):
            template_path = os.path.join(TEMPLATE_DIR, safe_name)
            if not os.path.exists(template_path):
                return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)
            placeholders = _scan_placeholders(template_path)
        elif safe_name.endswith('.docx'):
            template_path = os.path.join(TEMPLATE_DIR_WORD, safe_name)
            if not os.path.exists(template_path):
                return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)
            placeholders = _scan_placeholders_docx(template_path)
        else:
            return error_response('유효하지 않은 파일 형식입니다. (.pptx 또는 .docx)', status_code=400)

        return success_response(placeholders)
    except Exception as e:
        return error_response(f'플레이스홀더 스캔 실패: {str(e)}', status_code=500)


@bp.route('/datasources', methods=['GET'])
@jwt_required()
def get_datasources():
    """사용 가능한 데이터 소스 목록 반환"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        from app.modules.digital_twin_dashboard.models import DashboardData

        server_data = DashboardData.query.first()
        projects = []
        performances = []

        if server_data:
            raw_projects = server_data.projects or []
            projects = [p for p in raw_projects if not p.get('_permanentlyDeleted')]
            raw_performances = server_data.performances or []
            performances = [p for p in raw_performances if not p.get('_deleted')]

        # 성과를 uuid/id로 빠르게 찾기 위한 맵
        perf_by_uuid = {}
        perf_by_id = {}
        for p in performances:
            if p.get('uuid'):
                perf_by_uuid[p['uuid']] = p
            if p.get('id'):
                perf_by_id[p['id']] = p

        def _find_perf(ref):
            """성과목록 항목에서 실제 성과 데이터 찾기"""
            uuid_val = ref.get('성과항목UUID') or ref.get('uuid')
            if uuid_val and uuid_val in perf_by_uuid:
                return perf_by_uuid[uuid_val]
            id_val = ref.get('id')
            if id_val and id_val in perf_by_id:
                return perf_by_id[id_val]
            return None

        # 과제별 연결 성과 resolve + 플래트 필드 추가
        perf_sub_fields = ['성과항목', '대분류', '소분류', '단위',
                           '현재수준', '목표수준', '실적수준', '과제기여도']
        enriched_projects = []
        max_perf_count = 0

        for p in projects:
            ep = dict(p)  # 복사
            linked_list = p.get('성과목록', [])
            resolved = []
            for idx, ref in enumerate(linked_list):
                actual = _find_perf(ref)
                if not actual:
                    continue
                num = idx + 1
                resolved.append(actual)
                # 플래트 필드 추가: 성과1_성과항목, 성과1_대분류, ...
                for sf in perf_sub_fields:
                    # 과제기여도는 ref(연결 정보)에서, 나머지는 actual(성과 원본)에서
                    if sf == '과제기여도':
                        ep[f'성과{num}_{sf}'] = ref.get(sf, actual.get(sf, ''))
                    else:
                        ep[f'성과{num}_{sf}'] = actual.get(sf, '')
            ep['_linkedPerformanceCount'] = len(resolved)
            max_perf_count = max(max_perf_count, len(resolved))
            enriched_projects.append(ep)

        # 과제 필드 목록 (기본 + 연결 성과)
        project_fields = [
            '과제명', '과제년도', '사업부', '프로세스', '과제영역', '과제구분',
            '과제PL', '작성자', '시작', '종료', '진행상태', '진행률',
            '과제상세설명', 'PoC과제여부', '중점과제여부'
        ]
        # 연결 성과 플래트 필드 추가
        for i in range(1, max_perf_count + 1):
            for sf in perf_sub_fields:
                project_fields.append(f'성과{i}_{sf}')

        project_items = [
            {
                'id': p.get('uuid') or p.get('id'),
                'label': f"{p.get('id', '')} - {p.get('과제명', '')}",
                'perfCount': p.get('_linkedPerformanceCount', 0)
            }
            for p in enriched_projects
        ]

        # 성과 필드 목록
        performance_fields = [
            '성과항목', '성과년도', '대분류', '소분류', '단위',
            '현재수준', '목표수준', '실적수준', '과제기여도', '설명'
        ]
        performance_items = [
            {'id': p.get('uuid') or p.get('id'), 'label': f"{p.get('id', '')} - {p.get('성과항목', '')}"}
            for p in performances
        ]

        # ─── SPDM 플랫폼 현황 데이터 ───
        from app.modules.spdm_status.models import (
            SpdmIssue, SpdmGroup, SpdmScheduleDepartment, SpdmModule
        )

        # 이슈 데이터
        spdm_issues_raw = SpdmIssue.query.all()
        spdm_groups_map = {g.id: g.name for g in SpdmGroup.query.all()}
        spdm_issues = []
        for issue in spdm_issues_raw:
            d = issue.to_dict()
            d['이슈제목'] = d.get('title', '')
            d['상태'] = d.get('status', '')
            d['개설자'] = d.get('assignee', '')
            d['목표일정'] = d.get('dueDate', '')
            d['그룹명'] = spdm_groups_map.get(issue.group_id, '')
            d['관련사업부'] = ', '.join(d.get('departments', []))
            d['댓글수'] = d.get('historyCount', 0)
            spdm_issues.append(d)

        spdm_issue_fields = [
            '이슈제목', '상태', '개설자', '목표일정',
            '그룹명', '관련사업부', '댓글수'
        ]
        spdm_issue_items = [
            {'id': str(i['id']), 'label': f"#{i['id']} - {i.get('이슈제목', '')}"}
            for i in spdm_issues
        ]

        # 일정 데이터
        spdm_depts = SpdmScheduleDepartment.query.order_by(SpdmScheduleDepartment.order).all()
        spdm_schedules = []
        for dept in spdm_depts:
            dept_dict = dept.to_dict(include_schedules=True)
            for sch in dept_dict.get('schedules', []):
                sch['사업부명'] = dept_dict['departmentName']
                sch['차수'] = sch.get('phase', '')
                sch['일정제목'] = sch.get('title', '')
                sch['시작일'] = sch.get('startDate', '')
                sch['종료일'] = sch.get('endDate', '')
                sch['상태'] = sch.get('status', '')
                spdm_schedules.append(sch)

        spdm_schedule_fields = [
            '사업부명', '차수', '일정제목', '시작일', '종료일', '상태'
        ]
        spdm_schedule_items = [
            {'id': str(s['id']), 'label': f"{s['사업부명']} {s['차수']} - {s['일정제목']}"}
            for s in spdm_schedules
        ]

        # 모듈 데이터
        spdm_modules_raw = SpdmModule.query.order_by(SpdmModule.order).all()
        spdm_modules = []
        for m in spdm_modules_raw:
            d = m.to_dict()
            d['모듈명'] = d.get('name', '')
            d['사업부'] = d.get('divisionId', '')
            d['구현시스템'] = d.get('systemType', '')
            d['연계여부'] = '예' if d.get('isLinked') else '아니오'
            d['설명'] = d.get('description', '')
            spdm_modules.append(d)

        spdm_module_fields = [
            '모듈명', '사업부', '구현시스템', '연계여부', '설명'
        ]
        spdm_module_items = [
            {'id': str(m['id']), 'label': f"{m['사업부']} - {m['모듈명']}"}
            for m in spdm_modules
        ]

        # ─── 디지털 트윈 로드맵 데이터 ───
        from app.modules.digital_twin_reference.models import DtReferenceTask

        roadmap_tasks_raw = DtReferenceTask.query.order_by(DtReferenceTask.order).all()
        roadmap_tasks = []
        for t in roadmap_tasks_raw:
            d = t.to_dict()
            d['과제명'] = d.get('name', '')
            d['사업부'] = d.get('divisionId', '')
            d['테스트항목'] = d.get('testItem', '')
            d['상세내용'] = d.get('testItemDetail', '')
            d['구분'] = d.get('category', '')
            d['적용제품군'] = ', '.join(d.get('productFamily', []))
            d['연도'] = d.get('year', '')
            d['상태'] = d.get('status', '')
            connected = d.get('connectedDtTask', [])
            d['연결과제'] = ', '.join(
                c.get('projectName', '') for c in connected if c.get('projectName')
            )
            d['설명'] = d.get('description', '')
            roadmap_tasks.append(d)

        roadmap_fields = [
            '과제명', '사업부', '테스트항목', '상세내용',
            '구분', '적용제품군', '연도', '상태',
            '연결과제', '설명'
        ]
        roadmap_items = [
            {'id': str(t['id']), 'label': f"{t.get('사업부', '')} - {t['과제명']}"}
            for t in roadmap_tasks
        ]

        sources = [
            {
                'id': 'digital-twin-dashboard',
                'name': '디지털트윈 대시보드',
                'types': [
                    {
                        'id': 'projects',
                        'name': '과제 데이터',
                        'fields': project_fields,
                        'items': project_items,
                        'data': enriched_projects,
                        'perfSubFields': perf_sub_fields,
                        'maxPerfCount': max_perf_count
                    },
                    {
                        'id': 'performances',
                        'name': '성과 데이터',
                        'fields': performance_fields,
                        'items': performance_items,
                        'data': performances
                    }
                ]
            },
            {
                'id': 'spdm-status',
                'name': '플랫폼 현황',
                'types': [
                    {
                        'id': 'spdm-issues',
                        'name': '이슈 데이터',
                        'fields': spdm_issue_fields,
                        'items': spdm_issue_items,
                        'data': spdm_issues
                    },
                    {
                        'id': 'spdm-schedules',
                        'name': '일정 데이터',
                        'fields': spdm_schedule_fields,
                        'items': spdm_schedule_items,
                        'data': spdm_schedules
                    },
                    {
                        'id': 'spdm-modules',
                        'name': '모듈 데이터',
                        'fields': spdm_module_fields,
                        'items': spdm_module_items,
                        'data': spdm_modules
                    }
                ]
            },
            {
                'id': 'digital-twin-reference',
                'name': '개발 디지털 트윈 로드맵',
                'types': [
                    {
                        'id': 'roadmap-tasks',
                        'name': '로드맵 과제',
                        'fields': roadmap_fields,
                        'items': roadmap_items,
                        'data': roadmap_tasks
                    }
                ]
            }
        ]

        return success_response(sources)
    except Exception as e:
        return error_response(f'데이터 소스 조회 실패: {str(e)}', status_code=500)


@bp.route('/generate', methods=['POST'])
@jwt_required()
def generate_document():
    """매핑 정보를 기반으로 PPT/Word 문서 생성"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        template_filename = data.get('template')
        mappings = data.get('mappings', {})

        if not template_filename:
            return error_response('템플릿 파일명이 필요합니다.', status_code=400)

        safe_name = os.path.basename(template_filename)

        # 플레이스홀더 → 값 맵 구성
        placeholders = {}
        for placeholder_name, value in mappings.items():
            key = f'{{{{{placeholder_name}}}}}'
            placeholders[key] = str(value) if value is not None else ''

        if safe_name.endswith('.pptx'):
            template_path = os.path.join(TEMPLATE_DIR, safe_name)
            if not os.path.exists(template_path):
                return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)

            prs = Presentation(template_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            _replace_paragraph_placeholders(paragraph, placeholders)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    _replace_paragraph_placeholders(paragraph, placeholders)
            for slide in prs.slides:
                for shape in slide.shapes:
                    _auto_fit_text(slide, shape)

            doc_name = os.path.splitext(safe_name)[0]
            safe_doc_name = "".join(c for c in doc_name if c not in r'\/:*?"<>|')
            filename = f"{safe_doc_name[:50]}_생성본.pptx"

            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            return send_file(
                buffer,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        elif safe_name.endswith('.docx'):
            template_path = os.path.join(TEMPLATE_DIR_WORD, safe_name)
            if not os.path.exists(template_path):
                return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)

            doc = DocxDocument(template_path)
            _replace_all_placeholders_docx(doc, placeholders)

            doc_name = os.path.splitext(safe_name)[0]
            safe_doc_name = "".join(c for c in doc_name if c not in r'\/:*?"<>|')
            filename = f"{safe_doc_name[:50]}_생성본.docx"

            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)

            return send_file(
                buffer,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )

        else:
            return error_response('유효하지 않은 파일 형식입니다.', status_code=400)

    except Exception as e:
        print(f"[AutoDocument Error] Document generation failed: {str(e)}")
        return error_response(f'문서 생성 실패: {str(e)}', status_code=500)


@bp.route('/generate-batch', methods=['POST'])
@jwt_required()
def generate_batch_document():
    """여러 항목에 대해 동일 매핑으로 PPT/Word 일괄 생성"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        data = request.get_json()
        if not data:
            return error_response('요청 데이터가 없습니다.', status_code=400)

        template_filename = data.get('template')
        batch = data.get('batch', [])

        if not template_filename:
            return error_response('템플릿 파일명이 필요합니다.', status_code=400)
        if not batch:
            return error_response('생성할 항목이 없습니다.', status_code=400)

        safe_name = os.path.basename(template_filename)

        if safe_name.endswith('.docx'):
            return _generate_batch_docx(safe_name, batch)
        elif safe_name.endswith('.pptx'):
            return _generate_batch_pptx(safe_name, batch)
        else:
            return error_response('유효하지 않은 파일 형식입니다.', status_code=400)

    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"[AutoDocument Error] Batch generation failed: {str(e)}")
        return error_response(f'일괄 문서 생성 실패: {str(e)}', status_code=500)


def _generate_batch_pptx(safe_name, batch):
    """PPT 일괄 생성 (항목별 슬라이드 묶음)"""
    template_path = os.path.join(TEMPLATE_DIR, safe_name)
    if not os.path.exists(template_path):
        return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)

    # 첫 번째 항목으로 기본 프레젠테이션 생성
    first_mappings = batch[0].get('mappings', {})
    first_placeholders = {f'{{{{{k}}}}}': str(v) if v is not None else '' for k, v in first_mappings.items()}

    merged_prs = Presentation(template_path)
    for slide in merged_prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    _replace_paragraph_placeholders(paragraph, first_placeholders)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            _replace_paragraph_placeholders(paragraph, first_placeholders)
    for slide in merged_prs.slides:
        for shape in slide.shapes:
            _auto_fit_text(slide, shape)

    # 나머지 항목들: 각각 템플릿을 로드하고 치환한 뒤 슬라이드를 병합
    for entry in batch[1:]:
        item_mappings = entry.get('mappings', {})
        item_placeholders = {f'{{{{{k}}}}}': str(v) if v is not None else '' for k, v in item_mappings.items()}

        src_prs = Presentation(template_path)
        for slide in src_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        _replace_paragraph_placeholders(paragraph, item_placeholders)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                _replace_paragraph_placeholders(paragraph, item_placeholders)
        for slide in src_prs.slides:
            for shape in slide.shapes:
                _auto_fit_text(slide, shape)

        # 슬라이드 복사 (XML 레벨)
        for slide in src_prs.slides:
            from copy import deepcopy

            slide_layout = merged_prs.slide_layouts[0]
            new_slide = merged_prs.slides.add_slide(slide_layout)

            for ph in list(new_slide.placeholders):
                sp = ph._element
                sp.getparent().remove(sp)

            for shape in slide.shapes:
                el = deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            src_bg = slide._element.find('{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}bg')
            if src_bg is None:
                src_bg = slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
            if src_bg is not None:
                new_bg = deepcopy(src_bg)
                existing_bg = new_slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                if existing_bg is not None:
                    new_slide._element.replace(existing_bg, new_bg)
                else:
                    new_slide._element.insert(0, new_bg)

    doc_name = os.path.splitext(safe_name)[0]
    safe_doc_name = "".join(c for c in doc_name if c not in r'\/:*?"<>|')
    filename = f"{safe_doc_name[:50]}_{len(batch)}건_생성본.pptx"

    buffer = io.BytesIO()
    merged_prs.save(buffer)
    buffer.seek(0)

    return send_file(
        buffer,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


def _generate_batch_docx(safe_name, batch):
    """Word 일괄 생성 (각 항목을 페이지 구분으로 병합)"""
    template_path = os.path.join(TEMPLATE_DIR_WORD, safe_name)
    if not os.path.exists(template_path):
        return error_response(f'템플릿을 찾을 수 없습니다: {safe_name}', status_code=404)

    from docx.oxml.ns import qn
    from copy import deepcopy

    # 첫 번째 항목으로 기본 문서 생성
    first_mappings = batch[0].get('mappings', {})
    first_placeholders = {f'{{{{{k}}}}}': str(v) if v is not None else '' for k, v in first_mappings.items()}

    merged_doc = DocxDocument(template_path)
    _replace_all_placeholders_docx(merged_doc, first_placeholders)

    # 나머지 항목들: 각각 템플릿을 로드, 치환 후 페이지 나누기로 병합
    for entry in batch[1:]:
        item_mappings = entry.get('mappings', {})
        item_placeholders = {f'{{{{{k}}}}}': str(v) if v is not None else '' for k, v in item_mappings.items()}

        src_doc = DocxDocument(template_path)
        _replace_all_placeholders_docx(src_doc, item_placeholders)

        # 페이지 나누기 추가
        merged_doc.add_page_break()

        # 소스 문서의 본문 요소를 병합 문서에 복사
        for element in src_doc.element.body:
            # sectPr(섹션 속성)은 복사하지 않음
            if element.tag == qn('w:sectPr'):
                continue
            merged_doc.element.body.append(deepcopy(element))

    doc_name = os.path.splitext(safe_name)[0]
    safe_doc_name = "".join(c for c in doc_name if c not in r'\/:*?"<>|')
    filename = f"{safe_doc_name[:50]}_{len(batch)}건_생성본.docx"

    buffer = io.BytesIO()
    merged_doc.save(buffer)
    buffer.seek(0)

    return send_file(
        buffer,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


def _extract_division(performance_name):
    """성과항목명에서 사업부 추출: '[MX] 항목명' → 'MX'"""
    if not performance_name:
        return '미지정'
    m = re.match(r'^\[([^\]]+)\]', performance_name)
    return m.group(1) if m else '미지정'


def _safe_float(val):
    """안전하게 숫자 변환"""
    if val is None or val == '' or val == '-':
        return None
    try:
        return float(val)
    except (ValueError, TypeError):
        return None


def _calc_stats(values):
    """숫자 리스트로 합계, 평균, 최소, 최대, 건수 계산"""
    nums = [v for v in values if v is not None]
    if not nums:
        return {'합계': 0, '평균': 0, '최소': 0, '최대': 0, '건수': 0}
    return {
        '합계': round(sum(nums), 2),
        '평균': round(sum(nums) / len(nums), 2),
        '최소': round(min(nums), 2),
        '최대': round(max(nums), 2),
        '건수': len(nums)
    }


@bp.route('/presets', methods=['GET'])
@jwt_required()
def get_presets():
    """프리셋 목록 조회"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        presets = _load_presets()
        mode = request.args.get('mode')
        if mode:
            presets = [p for p in presets if p.get('mode') == mode]
        return success_response(presets)
    except Exception as e:
        return error_response(f'프리셋 조회 실패: {str(e)}', status_code=500)


@bp.route('/presets', methods=['POST'])
@jwt_required()
def save_preset():
    """프리셋 저장"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        data = request.get_json()
        if not data or not data.get('name'):
            return error_response('프리셋 이름이 필요합니다.', status_code=400)

        presets = _load_presets()
        preset = {
            'id': str(uuid.uuid4()),
            'name': data['name'],
            'mode': data.get('mode', 'single'),
            'template': data.get('template'),
            'mappings': data.get('mappings', {}),
            'sourceType': data.get('sourceType', ''),
            'year': data.get('year'),
            # 시간대를 **명시**해 저장한다(예: 2026-08-02T05:17:47+09:00).
            #   - 예전 `datetime.now().isoformat()` 은 오프셋이 없어 "이게 KST 인가
            #     UTC 인가" 를 값만 봐서는 알 수 없었다. 화면은 로컬로 파싱하고
            #     (`new Date(iso)`) DB 의 다른 시각은 UTC 라 섞이면 9시간 어긋난다.
            #   - `utcnow()` 로 바꾸면 안 된다 — 오프셋 없는 UTC 문자열을 화면이
            #     로컬로 읽어 **표시가 하루 밀린다.** 오프셋을 붙이는 게 정답이다.
            'createdAt': datetime.now().astimezone().isoformat(),
        }
        presets.append(preset)
        _save_presets(presets)

        return success_response(preset, '프리셋이 저장되었습니다.')
    except Exception as e:
        return error_response(f'프리셋 저장 실패: {str(e)}', status_code=500)


@bp.route('/presets/<preset_id>', methods=['DELETE'])
@jwt_required()
def delete_preset(preset_id):
    """프리셋 삭제"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        presets = _load_presets()
        original_len = len(presets)
        presets = [p for p in presets if p.get('id') != preset_id]

        if len(presets) == original_len:
            return error_response('프리셋을 찾을 수 없습니다.', status_code=404)

        _save_presets(presets)
        return success_response(None, '프리셋이 삭제되었습니다.')
    except Exception as e:
        return error_response(f'프리셋 삭제 실패: {str(e)}', status_code=500)


@bp.route('/stats', methods=['GET'])
@jwt_required()
def get_statistics():
    """통계 데이터 계산 및 반환"""
    try:
        user = User.query.get(get_jwt_identity())
        if not _check_admin(user):
            return error_response('관리자 권한이 필요합니다.', status_code=403)

        from app.modules.digital_twin_dashboard.models import DashboardData

        server_data = DashboardData.query.first()
        projects = []
        performances = []

        if server_data:
            projects = [p for p in (server_data.projects or []) if not p.get('_permanentlyDeleted')]
            performances = [p for p in (server_data.performances or []) if not p.get('_deleted')]

        current_year = int(request.args.get('year', 0)) or None

        # 연도 필터 적용
        if current_year:
            year_projects = [p for p in projects if p.get('과제년도') == current_year]
        else:
            year_projects = projects

        # ─── 1. 사업부별 과제 현황 ───
        division_stats = {}
        for p in year_projects:
            div = p.get('사업부', '미지정')
            if div not in division_stats:
                division_stats[div] = {'과제수': 0, '진행률_목록': [], '진행상태': {}}
            division_stats[div]['과제수'] += 1

            prog = _safe_float(p.get('진행률'))
            if prog is not None:
                division_stats[div]['진행률_목록'].append(prog)

            status = p.get('진행상태', '미착수')
            division_stats[div]['진행상태'][status] = division_stats[div]['진행상태'].get(status, 0) + 1

        division_summary = {}
        for div, info in division_stats.items():
            division_summary[div] = {
                '과제수': info['과제수'],
                '평균진행률': round(sum(info['진행률_목록']) / len(info['진행률_목록']), 1) if info['진행률_목록'] else 0,
                '진행상태': info['진행상태']
            }

        # ─── 2. 전체 과제 통계 ───
        total_project_stats = {
            '전체과제수': len(year_projects),
            '사업부별': division_summary,
            '전체평균진행률': round(
                sum(p for p in ((_safe_float(pr.get('진행률')) or 0) for pr in year_projects)) / len(year_projects), 1
            ) if year_projects else 0,
        }

        # 진행상태별 집계
        status_counts = {}
        for p in year_projects:
            s = p.get('진행상태', '미착수')
            status_counts[s] = status_counts.get(s, 0) + 1
        total_project_stats['진행상태별'] = status_counts

        # ─── 3. 성과 통계 (사업부 × 대분류 × 소분류) ───
        perf_grouped = {}  # {사업부: {대분류: {소분류: [perf...]}}}

        for perf in performances:
            div = _extract_division(perf.get('성과항목', ''))
            cat = perf.get('대분류', '미지정')
            sub = perf.get('소분류', '미지정')

            perf_grouped.setdefault(div, {}).setdefault(cat, {}).setdefault(sub, []).append(perf)

        performance_stats = {}
        for div, cats in perf_grouped.items():
            performance_stats[div] = {}
            for cat, subs in cats.items():
                performance_stats[div][cat] = {}
                for sub, perfs in subs.items():
                    currents = [_safe_float(p.get('현재수준')) for p in perfs]
                    targets = [_safe_float(p.get('목표수준')) for p in perfs]
                    actuals = [_safe_float(p.get('실적수준')) for p in perfs]

                    performance_stats[div][cat][sub] = {
                        '건수': len(perfs),
                        '현재수준': _calc_stats(currents),
                        '목표수준': _calc_stats(targets),
                        '실적수준': _calc_stats(actuals),
                        '단위': perfs[0].get('단위', '') if perfs else '',
                        '항목목록': [p.get('성과항목', '') for p in perfs]
                    }

        # ─── 4. 플랫 통계 필드 (플레이스홀더 매핑용) ───
        flat_fields = {}

        # 전체
        flat_fields['전체과제수'] = total_project_stats['전체과제수']
        flat_fields['전체평균진행률'] = total_project_stats['전체평균진행률']
        flat_fields['전체성과수'] = len(performances)

        for status, cnt in status_counts.items():
            flat_fields[f'진행상태_{status}_과제수'] = cnt

        # 사업부별
        for div, info in division_summary.items():
            flat_fields[f'{div}_과제수'] = info['과제수']
            flat_fields[f'{div}_평균진행률'] = info['평균진행률']
            for status, cnt in info['진행상태'].items():
                flat_fields[f'{div}_{status}_과제수'] = cnt

        # 성과 사업부별/대분류별/소분류별
        for div, cats in performance_stats.items():
            div_total_count = 0
            div_all_currents = []
            div_all_targets = []
            div_all_actuals = []

            for cat, subs in cats.items():
                cat_count = 0
                cat_currents = []
                cat_targets = []
                cat_actuals = []

                for sub, info in subs.items():
                    cnt = info['건수']
                    cat_count += cnt
                    cat_currents.extend([_safe_float(v) for v in [info['현재수준']['합계']]])
                    cat_targets.extend([_safe_float(v) for v in [info['목표수준']['합계']]])
                    cat_actuals.extend([_safe_float(v) for v in [info['실적수준']['합계']]])

                    # 소분류 레벨
                    flat_fields[f'{div}_{cat}_{sub}_건수'] = cnt
                    flat_fields[f'{div}_{cat}_{sub}_현재수준_합계'] = info['현재수준']['합계']
                    flat_fields[f'{div}_{cat}_{sub}_현재수준_평균'] = info['현재수준']['평균']
                    flat_fields[f'{div}_{cat}_{sub}_목표수준_합계'] = info['목표수준']['합계']
                    flat_fields[f'{div}_{cat}_{sub}_목표수준_평균'] = info['목표수준']['평균']
                    flat_fields[f'{div}_{cat}_{sub}_실적수준_합계'] = info['실적수준']['합계']
                    flat_fields[f'{div}_{cat}_{sub}_실적수준_평균'] = info['실적수준']['평균']

                div_total_count += cat_count
                div_all_currents.extend(cat_currents)
                div_all_targets.extend(cat_targets)
                div_all_actuals.extend(cat_actuals)

                # 대분류 레벨
                flat_fields[f'{div}_{cat}_건수'] = cat_count

            flat_fields[f'{div}_성과수'] = div_total_count

        # ─── 4. KPI 대시보드 데이터 ───
        from app.modules.digital_twin_dashboard.models import KPIDashboardCard, KPICategory

        # ⚠️ 여기는 `utcnow()` 로 바꾸면 **안 된다.** 이건 시각이 아니라 "지금이 몇 년도
        #    업무인가" 다. KST 1월 1일 00~09시에는 UTC 가 아직 작년이라, utcnow 를 쓰면
        #    새해 첫 아침에 작년 KPI 를 보여준다. 업무 기준 연도는 로컬(KST)이 맞다.
        query_year = current_year or datetime.now().year
        kpi_cards = KPIDashboardCard.query.filter_by(year=query_year).order_by(KPIDashboardCard.order).all()
        kpi_categories = {str(c.id): c.name for c in KPICategory.query.filter_by(year=query_year).all()}

        # 성과를 key로 매핑 (uuid → 성과, id → 성과)
        perf_by_key = {}
        for p in performances:
            for key_field in ['uuid', 'id', '성과항목']:
                if p.get(key_field):
                    perf_by_key[p[key_field]] = p

        for card_data in kpi_cards:
            card = card_data.to_dict()
            card_name = card['name']
            selected_keys = card.get('selectedPerfKeys', [])
            logic = card.get('logic', '합계')

            # 선택된 성과 필터
            filtered = [perf_by_key[k] for k in selected_keys if k in perf_by_key]

            def _aggregate(vals, agg_logic):
                if not vals:
                    return None
                s = sum(vals)
                return round(s, 2) if agg_logic == '합계' else round(s / len(vals), 2)

            # 현재/목표: 전체 성과
            all_current = [v for v in (_safe_float(p.get('현재수준')) for p in filtered) if v is not None]
            all_target = [v for v in (_safe_float(p.get('목표수준')) for p in filtered) if v is not None]
            current_val = _aggregate(all_current, logic)
            target_val = _aggregate(all_target, logic)

            # 단일 실적
            single_perfs = [p for p in filtered if not (p.get('월별실적여부') and isinstance(p.get('월별실적'), list) and len(p.get('월별실적', [])) > 0)]
            single_vals = [v for v in (_safe_float(p.get('실적수준')) for p in single_perfs) if v is not None]
            actual_val = _aggregate(single_vals, logic)

            # 월별 실적
            monthly_perfs = [p for p in filtered if p.get('월별실적여부') and isinstance(p.get('월별실적'), list) and len(p.get('월별실적', [])) > 0]
            monthly_vals = {}
            if monthly_perfs:
                for m in range(12):
                    m_vals = [v for v in (_safe_float(p['월별실적'][m]) if m < len(p.get('월별실적', [])) else None for p in monthly_perfs) if v is not None]
                    agg = _aggregate(m_vals, logic)
                    if agg is not None:
                        monthly_vals[m + 1] = agg

            flat_fields[f'KPI_{card_name}_현재'] = current_val if current_val is not None else ''
            flat_fields[f'KPI_{card_name}_목표'] = target_val if target_val is not None else ''
            flat_fields[f'KPI_{card_name}_실적'] = actual_val if actual_val is not None else ''
            flat_fields[f'KPI_{card_name}_성과수'] = len(filtered)

            for month_num, month_val in monthly_vals.items():
                flat_fields[f'KPI_{card_name}_{month_num}월'] = month_val

        # ─── 5. 부서별 현황 데이터 ───
        dept_stats = {}
        dept_personnel = {}  # {부서: set(인력)}

        for p in year_projects:
            dept_list = p.get('담당부서목록', [])
            if not isinstance(dept_list, list):
                continue
            for dept in dept_list:
                if not dept:
                    continue
                if dept not in dept_stats:
                    dept_stats[dept] = {'과제수': 0, '성과수': 0, '진행률_목록': [], '진행상태': {}}
                dept_stats[dept]['과제수'] += 1
                dept_stats[dept]['성과수'] += len(p.get('성과목록', []))

                prog = _safe_float(p.get('진행률'))
                if prog is not None:
                    dept_stats[dept]['진행률_목록'].append(prog)

                status = p.get('진행상태', '미착수')
                dept_stats[dept]['진행상태'][status] = dept_stats[dept]['진행상태'].get(status, 0) + 1

            # 인력 집계
            personnel_list = p.get('과제참여인력목록', [])
            if isinstance(personnel_list, list):
                for person in personnel_list:
                    dept_name = person.get('부서', '')
                    name = person.get('이름', '')
                    if dept_name and name:
                        if dept_name not in dept_personnel:
                            dept_personnel[dept_name] = set()
                        dept_personnel[dept_name].add(name)

        for dept, info in dept_stats.items():
            avg_prog = round(sum(info['진행률_목록']) / len(info['진행률_목록']), 1) if info['진행률_목록'] else 0
            flat_fields[f'부서_{dept}_과제수'] = info['과제수']
            flat_fields[f'부서_{dept}_성과수'] = info['성과수']
            flat_fields[f'부서_{dept}_평균진행률'] = avg_prog
            flat_fields[f'부서_{dept}_인원수'] = len(dept_personnel.get(dept, set()))
            for status, cnt in info['진행상태'].items():
                flat_fields[f'부서_{dept}_{status}_과제수'] = cnt

        flat_fields['전체부서수'] = len(dept_stats)
        flat_fields['전체인원수'] = len(set().union(*dept_personnel.values())) if dept_personnel else 0

        result = {
            'projects': total_project_stats,
            'performances': performance_stats,
            'flat_fields': flat_fields
        }

        return success_response(result)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return error_response(f'통계 데이터 조회 실패: {str(e)}', status_code=500)
